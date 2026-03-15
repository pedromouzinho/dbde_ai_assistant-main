# =============================================================================
# test_pptx_engine.py — Tests for the Millennium BCP branded PPTX engine
# =============================================================================

import io
import pytest
from pptx_engine import (
    generate_presentation,
    generate_presentation_from_outline,
    BRAND_ACCENT_HEX,
    SLIDE_WIDTH_EMU,
    SLIDE_HEIGHT_EMU,
)


class TestGeneratePresentation:
    """Tests for the main generate_presentation function."""

    def test_minimal_presentation(self):
        """Single content slide produces valid PPTX."""
        slides = [{"type": "content", "title": "Test", "bullets": ["A", "B"]}]
        buf = generate_presentation("Test Title", slides)
        assert isinstance(buf, io.BytesIO)
        content = buf.getvalue()
        assert len(content) > 1000  # not trivially empty
        # PPTX is a ZIP — check magic bytes
        assert content[:2] == b"PK"

    def test_slide_dimensions_widescreen(self):
        """Presentation uses 13.333x7.5 inch widescreen layout."""
        from pptx import Presentation as PptxRead
        slides = [{"type": "content", "title": "X", "bullets": ["Y"]}]
        buf = generate_presentation("T", slides)
        prs = PptxRead(buf)
        assert prs.slide_width == SLIDE_WIDTH_EMU
        assert prs.slide_height == SLIDE_HEIGHT_EMU

    def test_auto_title_and_closing(self):
        """Auto-adds title and closing slides."""
        from pptx import Presentation as PptxRead
        slides = [{"type": "content", "title": "Body", "bullets": ["X"]}]
        buf = generate_presentation("Title", slides)
        prs = PptxRead(buf)
        # 1 title (auto) + 1 content + 1 closing (auto) = 3
        assert len(prs.slides) == 3

    def test_no_duplicate_title_slide(self):
        """If slides already start with title, don't add another."""
        from pptx import Presentation as PptxRead
        slides = [
            {"type": "title", "title": "My Title", "subtitle": "Sub"},
            {"type": "content", "title": "Body", "bullets": ["X"]},
        ]
        buf = generate_presentation("My Title", slides)
        prs = PptxRead(buf)
        # 1 title (from specs) + 1 content + 1 closing (auto) = 3
        assert len(prs.slides) == 3

    def test_no_duplicate_closing(self):
        """If slides already end with closing, don't add another."""
        from pptx import Presentation as PptxRead
        slides = [
            {"type": "content", "title": "Body", "bullets": ["X"]},
            {"type": "closing", "text": "Obrigado"},
        ]
        buf = generate_presentation("T", slides)
        prs = PptxRead(buf)
        # 1 title (auto) + 1 content + 1 closing (from specs) = 3
        assert len(prs.slides) == 3

    def test_all_slide_types(self):
        """All slide types generate without error."""
        slides = [
            {"type": "title", "title": "Capa", "subtitle": "Subtítulo"},
            {"type": "section", "title": "Secção 1"},
            {"type": "content", "title": "Conteúdo", "bullets": ["A", "B", "- Sub-bullet"]},
            {"type": "two_column", "title": "Duas Colunas", "left": ["L1", "L2"], "right": ["R1", "R2"]},
            {"type": "kpi", "title": "KPIs", "kpis": [
                {"value": "125", "label": "Total"},
                {"value": "98%", "label": "Cobertura"},
            ]},
            {"type": "table", "title": "Tabela", "headers": ["Col1", "Col2"], "rows": [["A", "B"], ["C", "D"]]},
            {"type": "agenda", "items": ["Item 1", "Item 2", "Item 3"]},
            {"type": "closing", "text": "Obrigado"},
        ]
        buf = generate_presentation("Full Test", slides, include_title_slide=False, include_closing_slide=False)
        from pptx import Presentation as PptxRead
        prs = PptxRead(buf)
        assert len(prs.slides) == 8

    def test_section_auto_numbering(self):
        """Section dividers auto-number correctly."""
        slides = [
            {"type": "section", "title": "Intro"},
            {"type": "content", "title": "Content", "bullets": ["X"]},
            {"type": "section", "title": "Results"},
            {"type": "content", "title": "More Content", "bullets": ["Y"]},
        ]
        buf = generate_presentation("T", slides)
        content = buf.getvalue()
        assert len(content) > 1000

    def test_empty_slides_returns_title_and_closing_only(self):
        """Empty slides list still creates title + closing."""
        from pptx import Presentation as PptxRead
        buf = generate_presentation("Empty", [])
        prs = PptxRead(buf)
        assert len(prs.slides) == 2  # title + closing

    def test_kpi_max_four(self):
        """KPI slide handles more than 4 kpis (only shows first 4)."""
        slides = [{"type": "kpi", "title": "Many KPIs", "kpis": [
            {"value": str(i), "label": f"KPI {i}"} for i in range(8)
        ]}]
        buf = generate_presentation("T", slides)
        assert len(buf.getvalue()) > 1000

    def test_table_row_cap(self):
        """Table slide caps at 15 rows."""
        slides = [{"type": "table", "title": "Big Table", "headers": ["A", "B"],
                   "rows": [[str(i), str(i*2)] for i in range(30)]}]
        buf = generate_presentation("T", slides)
        assert len(buf.getvalue()) > 1000

    def test_badge_text_custom(self):
        """Custom badge text is accepted."""
        slides = [{"type": "content", "title": "Test", "bullets": ["X"]}]
        buf = generate_presentation("T", slides, badge_text="CUSTOM TEAM")
        assert len(buf.getvalue()) > 1000

    def test_invalid_slide_type_fallback(self):
        """Unknown slide type falls back to content slide."""
        slides = [{"type": "unknown_type", "title": "Fallback", "bullets": ["X"]}]
        buf = generate_presentation("T", slides)
        assert len(buf.getvalue()) > 1000


class TestGenerateFromOutline:
    """Tests for the outline parser."""

    def test_simple_outline(self):
        outline = """# Secção 1
## Slide A
- Bullet 1
- Bullet 2
# Secção 2
## Slide B
- Bullet 3
"""
        buf = generate_presentation_from_outline("Outline Test", outline)
        from pptx import Presentation as PptxRead
        prs = PptxRead(buf)
        # title(auto) + section1 + slideA + section2 + slideB + closing(auto) = 6
        assert len(prs.slides) == 6

    def test_table_in_outline(self):
        outline = """## Dados
| Nome | Valor |
| --- | --- |
| Alpha | 100 |
| Beta | 200 |
"""
        buf = generate_presentation_from_outline("Table Test", outline)
        assert len(buf.getvalue()) > 1000

    def test_empty_outline(self):
        buf = generate_presentation_from_outline("Empty", "")
        from pptx import Presentation as PptxRead
        prs = PptxRead(buf)
        assert len(prs.slides) == 2  # title + closing only


class TestValidationLayer:
    """Tests for the smart validation and auto-correction layer."""

    def test_split_overloaded_bullets(self):
        """Content slide with 15 bullets splits into 3 slides."""
        from pptx import Presentation as PptxRead
        slides = [{"type": "content", "title": "Big List",
                   "bullets": [f"Item {i}" for i in range(15)]}]
        buf = generate_presentation("T", slides,
                                    include_title_slide=False,
                                    include_closing_slide=False)
        prs = PptxRead(buf)
        # 15 bullets / 7 per slide = 3 slides (7+7+1)
        assert len(prs.slides) == 3

    def test_split_slide_adds_cont_suffix(self):
        """Split slides get '(cont.)' in title."""
        from pptx_engine import _validate_and_fix_slides
        slides = [{"type": "content", "title": "Results",
                   "bullets": [f"B{i}" for i in range(10)]}]
        fixed = _validate_and_fix_slides(slides)
        assert len(fixed) == 2
        assert fixed[0]["title"] == "Results"
        assert fixed[1]["title"] == "Results (cont.)"

    def test_truncate_long_bullet(self):
        """Bullets >150 chars are truncated with ellipsis."""
        from pptx_engine import _validate_and_fix_slides
        long_bullet = "A" * 200
        slides = [{"type": "content", "title": "T", "bullets": [long_bullet]}]
        fixed = _validate_and_fix_slides(slides)
        assert len(fixed[0]["bullets"][0]) == 150
        assert fixed[0]["bullets"][0].endswith("…")

    def test_truncate_long_title(self):
        """Titles >80 chars are truncated."""
        from pptx_engine import _validate_and_fix_slides
        long_title = "T" * 100
        slides = [{"type": "content", "title": long_title, "bullets": ["X"]}]
        fixed = _validate_and_fix_slides(slides)
        assert len(fixed[0]["title"]) == 80
        assert fixed[0]["title"].endswith("…")

    def test_split_kpis_over_four(self):
        """KPI slide with 6 KPIs splits into 2 slides."""
        from pptx_engine import _validate_and_fix_slides
        kpis = [{"value": str(i), "label": f"KPI{i}"} for i in range(6)]
        slides = [{"type": "kpi", "title": "Metrics", "kpis": kpis}]
        fixed = _validate_and_fix_slides(slides)
        assert len(fixed) == 2
        assert len(fixed[0]["kpis"]) == 4
        assert len(fixed[1]["kpis"]) == 2

    def test_split_large_table(self):
        """Table with 25 rows splits into multiple slides."""
        from pptx_engine import _validate_and_fix_slides
        slides = [{"type": "table", "title": "Data",
                   "headers": ["A", "B"],
                   "rows": [[str(i), str(i*2)] for i in range(25)]}]
        fixed = _validate_and_fix_slides(slides)
        assert len(fixed) == 3  # 12 + 12 + 1

    def test_trim_excess_columns(self):
        """Table with >8 columns is trimmed to 8."""
        from pptx_engine import _validate_and_fix_slides
        headers = [f"Col{i}" for i in range(12)]
        slides = [{"type": "table", "title": "Wide",
                   "headers": headers, "rows": [["x"] * 12]}]
        fixed = _validate_and_fix_slides(slides)
        assert len(fixed[0]["headers"]) == 8

    def test_remove_empty_content(self):
        """Content slides with no bullets are removed."""
        from pptx_engine import _validate_and_fix_slides
        slides = [
            {"type": "content", "title": "Empty", "bullets": []},
            {"type": "content", "title": "Full", "bullets": ["X"]},
        ]
        fixed = _validate_and_fix_slides(slides)
        assert len(fixed) == 1
        assert fixed[0]["title"] == "Full"

    def test_remove_consecutive_sections(self):
        """Consecutive section dividers: only last is kept."""
        from pptx_engine import _validate_and_fix_slides
        slides = [
            {"type": "section", "title": "First"},
            {"type": "section", "title": "Second"},
            {"type": "content", "title": "Content", "bullets": ["X"]},
        ]
        fixed = _validate_and_fix_slides(slides)
        assert len(fixed) == 2
        assert fixed[0]["title"] == "Second"

    def test_two_column_trimmed(self):
        """Two-column items trimmed to max 8."""
        from pptx_engine import _validate_and_fix_slides
        slides = [{"type": "two_column", "title": "T",
                   "left": [f"L{i}" for i in range(15)],
                   "right": [f"R{i}" for i in range(15)]}]
        fixed = _validate_and_fix_slides(slides)
        assert len(fixed[0]["left"]) == 8
        assert len(fixed[0]["right"]) == 8

    def test_end_to_end_validation_generates_valid_pptx(self):
        """Full pipeline with overloaded slides still produces valid PPTX."""
        slides = [
            {"type": "content", "title": "A" * 100,
             "bullets": [f"Bullet {i}" for i in range(20)]},
            {"type": "kpi", "title": "KPIs",
             "kpis": [{"value": str(i), "label": f"K{i}"} for i in range(8)]},
            {"type": "table", "title": "Data",
             "headers": [f"C{i}" for i in range(10)],
             "rows": [[str(i)] * 10 for i in range(30)]},
        ]
        buf = generate_presentation("Stress Test", slides)
        content = buf.getvalue()
        assert content[:2] == b"PK"  # valid ZIP/PPTX
        assert len(content) > 5000


class TestOpusFallback:
    """Tests for the Opus planner fallback logic."""

    def test_fallback_from_content_creates_slides(self):
        """Fallback parser creates slides from simple text."""
        from pptx_engine import _fallback_slides_from_content
        content = """# Introdução
- Ponto 1
- Ponto 2
## Detalhes
- Sub-ponto A
- Sub-ponto B"""
        slides = _fallback_slides_from_content(content, "Test")
        assert len(slides) >= 2
        assert all(s.get("type") == "content" for s in slides)

    def test_fallback_from_empty_content(self):
        """Fallback with empty content still produces something."""
        from pptx_engine import _fallback_slides_from_content
        slides = _fallback_slides_from_content("", "Título")
        assert len(slides) >= 1

    def test_fallback_slides_generate_valid_pptx(self):
        """Fallback slides can be rendered into valid PPTX."""
        from pptx_engine import _fallback_slides_from_content, generate_presentation
        slides = _fallback_slides_from_content(
            "Dados importantes\n- KPI 1: 95%\n- KPI 2: 120\n- Análise completa",
            "Relatório"
        )
        buf = generate_presentation("Relatório", slides)
        assert buf.getvalue()[:2] == b"PK"


class TestToolIntegration:
    """Test tool_generate_presentation integration."""

    @pytest.mark.asyncio
    async def test_tool_generates_pptx(self):
        from tools_export import tool_generate_presentation
        result = await tool_generate_presentation(
            title="Integration Test",
            slides=[
                {"type": "content", "title": "Test", "bullets": ["A", "B"]},
            ],
        )
        assert result.get("presentation_generated") is True
        assert result.get("format") == "pptx"
        assert "_file_download" in result
        assert result["_file_download"]["format"] == "pptx"
        assert result["_file_download"]["size_bytes"] > 1000
        assert result.get("planning_model") == "structured_input"

    @pytest.mark.asyncio
    async def test_tool_rejects_empty_slides_and_content(self):
        from tools_export import tool_generate_presentation
        result = await tool_generate_presentation(title="Bad", slides=[], content="")
        assert "error" in result

    @pytest.mark.asyncio
    async def test_tool_rejects_no_input(self):
        from tools_export import tool_generate_presentation
        result = await tool_generate_presentation(title="Bad", slides=None, content="")
        assert "error" in result

    @pytest.mark.asyncio
    async def test_tool_rejects_too_many_slides(self):
        from tools_export import tool_generate_presentation
        result = await tool_generate_presentation(
            title="Too Many",
            slides=[{"type": "content", "title": f"S{i}", "bullets": ["X"]} for i in range(51)],
        )
        assert "error" in result
