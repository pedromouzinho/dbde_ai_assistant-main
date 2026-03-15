# =============================================================================
# app.py — FastAPI routes + wiring v7.2
# =============================================================================
# Thin routing layer: liga todos os módulos, expõe endpoints.
# Nenhuma lógica de negócio aqui — apenas routing e error handling.
# =============================================================================

import io
import json
import html
import base64
import zipfile
import os
import uuid
import asyncio
import logging
import traceback
import re
import hashlib
import time
import contextlib
from urllib.parse import urlsplit, unquote
from pathlib import Path
from collections import deque, OrderedDict
from dataclasses import dataclass
from datetime import datetime, timezone, timedelta
from typing import Optional, Dict, Any, List, Callable, Tuple

import httpx
from fastapi import FastAPI, Depends, HTTPException, UploadFile, File, Form, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, HTMLResponse, Response, JSONResponse
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from fastapi.staticfiles import StaticFiles


class JSONFormatter(logging.Formatter):
    def format(self, record):
        log_entry = {
            "timestamp": self.formatTime(record, self.datefmt),
            "level": record.levelname,
            "logger": record.name,
            "message": record.getMessage(),
        }
        if record.exc_info and record.exc_info[0]:
            log_entry["exception"] = self.formatException(record.exc_info)
        return json.dumps(log_entry, ensure_ascii=False)


def configure_logging(log_format: str = "text") -> None:
    handler = logging.StreamHandler()
    if str(log_format or "text").strip().lower() == "json":
        handler.setFormatter(JSONFormatter())
    else:
        handler.setFormatter(logging.Formatter("%(asctime)s %(levelname)s %(name)s %(message)s"))
    logging.basicConfig(level=logging.INFO, handlers=[handler], force=True)


# Bootstrap logging before optional imports that may emit warnings.
configure_logging(
    (os.getenv("LOG_FORMAT") or os.getenv("APPSETTING_LOG_FORMAT") or "text").strip().lower()
)

_pptx_import_traceback = ""

try:
    from pptx import Presentation
except Exception:
    Presentation = None
    _pptx_import_traceback = traceback.format_exc()
    logging.getLogger(__name__).warning(
        "[App] python-pptx import failed; uploads .pptx indisponíveis.\n%s",
        _pptx_import_traceback,
    )

from config import (
    APP_VERSION, APP_TITLE, APP_DESCRIPTION,
    AZURE_OPENAI_KEY,
    AZURE_SPEECH_ENABLED, AZURE_SPEECH_KEY, AZURE_SPEECH_REGION, AZURE_SPEECH_LANGUAGE,
    DEVOPS_INDEX, OMNI_INDEX, EXAMPLES_INDEX,
    DEVOPS_ORG, DEVOPS_PROJECT,
    SEARCH_SERVICE, SEARCH_KEY, API_VERSION_SEARCH,
    LLM_TIER_FAST, LLM_TIER_STANDARD, LLM_TIER_PRO, LLM_TIER_VISION, VISION_ENABLED,
    SPEECH_PROMPT_PRIMARY_SPEC, SPEECH_PROMPT_FALLBACK_SPEC,
    PROVIDER_GOVERNANCE_MODE, PROVIDER_EXTERNAL_MODEL_FAMILIES,
    PROVIDER_GOVERNANCE_EXPERIMENTAL_ALLOW_EXTERNAL,
    MODEL_ROUTER_ENABLED, MODEL_ROUTER_SPEC, MODEL_ROUTER_TARGET_TIERS, MODEL_ROUTER_NON_PROD_ONLY,
    IS_PRODUCTION,
    RERANK_ENABLED, RERANK_MODEL, RERANK_ENDPOINT, RERANK_TOP_N, RERANK_AUTH_MODE,
    ALLOWED_ORIGINS, LOG_FORMAT,
    AUTH_COOKIE_NAME, AUTH_COOKIE_SECURE, AUTH_COOKIE_MAX_AGE_SECONDS,
    UPLOAD_MAX_FILES_PER_CONVERSATION, UPLOAD_MAX_FILE_BYTES,
    UPLOAD_MAX_CONCURRENT_JOBS, UPLOAD_MAX_PENDING_JOBS_PER_USER,
    UPLOAD_MAX_IMAGES_PER_MESSAGE, UPLOAD_EMBEDDING_CONCURRENCY,
    UPLOAD_MAX_CHUNKS_PER_FILE, UPLOAD_JOB_STALE_SECONDS,
    UPLOAD_TABULAR_DEEP_INGEST_MAX_BYTES, UPLOAD_TABULAR_DEEP_INGEST_MAX_ROWS,
    UPLOAD_TABULAR_DEEP_INGEST_RECORD_LIMIT,
    UPLOAD_TABULAR_ARTIFACT_ENABLED,
    UPLOAD_MAX_BATCH_TOTAL_BYTES,
    UPLOAD_BLOB_CONTAINER_RAW, UPLOAD_BLOB_CONTAINER_TEXT, UPLOAD_BLOB_CONTAINER_CHUNKS, UPLOAD_BLOB_CONTAINER_ARTIFACTS,
    UPLOAD_INDEX_TOP, UPLOAD_INLINE_WORKER_ENABLED, UPLOAD_INLINE_WORKER_RUNTIME_ENABLED,
    UPLOAD_DEDICATED_WORKER_ENABLED,
    UPLOAD_WORKER_POLL_SECONDS,
    UPLOAD_WORKER_BATCH_SIZE, UPLOAD_ARTIFACT_RETENTION_HOURS, UPLOAD_TABULAR_RAW_RETENTION_HOURS,
    UPLOAD_TABULAR_READY_RAW_RETENTION_HOURS,
    UPLOAD_TABULAR_CHUNK_BACKFILL_BATCH_SIZE,
    UPLOAD_RETENTION_SWEEP_INTERVAL_SECONDS,
    UPLOAD_FRONTEND_ASYNC_THRESHOLD_BYTES,
    DOC_INTEL_ENABLED, DOC_INTEL_MODEL,
    CHAT_TOOLRESULT_BLOB_CONTAINER,
    EXPORT_AUTO_ASYNC_ENABLED, EXPORT_ASYNC_THRESHOLD_ROWS, EXPORT_MAX_CONCURRENT_JOBS,
    EXPORT_JOB_STALE_SECONDS, EXPORT_INLINE_WORKER_ENABLED, EXPORT_DEDICATED_WORKER_ENABLED,
    EXPORT_WORKER_POLL_SECONDS,
    EXPORT_WORKER_BATCH_SIZE,
    EXPORT_BRAND_COLOR, EXPORT_BRAND_NAME, EXPORT_AGENT_NAME,
    STARTUP_FAIL_FAST, TOKEN_QUOTA_CONFIG, CHAT_BUDGET_PER_MINUTE,
    STORY_LANE_ENABLED,
    WORKER_RUN_DIR, UPLOAD_WORKER_PID_FILE, EXPORT_WORKER_PID_FILE,
)
from models import (
    AgentChatRequest, AgentChatResponse,
    LoginRequest, CreateUserRequest, ChangePasswordRequest,
    FeedbackRequest, ExportRequest, SaveChatRequest, UpdateChatTitleRequest, PrivacyDeleteRequest,
    SpeechPromptNormalizeRequest, SpeechPromptNormalizeResponse, SpeechPromptTokenResponse,
    ModeSwitchRequest, ModeSwitchResponse,
    ClientErrorReport,
    UserStoryWorkspaceRequest,
    UserStoryValidateRequest,
    UserStoryPublishRequest,
    UserStoryFeedbackEventRequest,
    UserStoryCurationReviewRequest,
    UserStoryPromoteRequest,
    UserStorySearchSyncRequest,
    UserStoryDevOpsSyncRequest,
    UserStoryKnowledgeSyncRequest,
    UserStoryKnowledgeAssetUploadRequest,
    UserStoryKnowledgeAssetTextRequest,
    UserStoryKnowledgeAssetBundleRequest,
    UserStoryKnowledgeAssetReviewRequest,
)
from auth import (
    get_current_user, get_current_principal, jwt_encode, jwt_decode, hash_password, verify_password,
    set_request_cookie_token, reset_request_cookie_token,
    set_request_auth_payload, reset_request_auth_payload,
    set_request_auth_error, reset_request_auth_error,
    record_login_failure, is_account_locked, clear_login_attempts,
    cleanup_blacklist,
    _LOCKOUT_DURATION_MINUTES,
    principal_is_admin,
)
from storage import (
    init_http_client, ensure_tables_exist,
    table_insert, table_query, table_merge, table_delete,
    ensure_blob_containers,
    blob_upload_bytes, blob_upload_json, blob_download_bytes, blob_download_json, blob_delete,
    blob_upload_stream,
    parse_blob_ref, build_blob_ref,
)
from auth_runtime import (
    validate_request_token,
    is_account_locked_persistent,
    record_login_failure_persistent,
    clear_login_failures_persistent,
    revoke_token_persistent,
    persist_user_invalidation,
)
from tools import (
    get_embedding, get_generated_file,
    _store_generated_file,
    _devops_url, _devops_headers,
)
from http_helpers import devops_request_with_retry as _devops_request_with_retry
from tools_knowledge import _close_http_client as _close_knowledge_client
from tools_figma import _close_http_client as _close_figma_client
from tools_miro import _close_http_client as _close_miro_client
from pii_shield import close_http_client as _close_pii_http_client
from tool_registry import get_registered_tool_names
from learning import invalidate_prompt_rules_cache
from agent import (
    agent_chat as _agent_chat, agent_chat_stream,
    conversations, conversation_meta, uploaded_files_store,
    switch_conversation_mode,
)
from export_engine import to_csv, to_xlsx, to_pdf, to_svg_bar_chart, to_html_report
from llm_provider import llm_simple, close_all_providers
from rate_limit_storage import TableStorageRateLimit
from tabular_loader import (
    TabularLoaderError,
    detect_polymorphic_schema,
    get_tabular_upload_limit_bytes,
    get_tabular_upload_limits,
    is_tabular_filename,
    load_tabular_dataset,
    load_tabular_preview,
)
from tabular_artifacts import (
    build_tabular_artifact,
    iter_tabular_artifact_batches,
    load_tabular_artifact_dataset,
    load_tabular_artifact_preview,
)
from job_store import PersistentJobStore
from utils import odata_escape, safe_blob_component, create_logged_task
from tool_metrics import tool_metrics
from document_intelligence import analyze_document, tables_to_markdown
import token_quota as _tq_module
from token_quota import TokenQuotaManager
from user_story_lane import (
    build_context_preview as build_user_story_context_preview,
    build_user_story_eval_summary,
    generate_user_story as generate_user_story_draft,
    promote_user_story_to_curated_corpus,
    review_user_story_curated_candidate,
    sync_user_story_examples_search_index,
    validate_user_story_request,
    publish_user_story,
    record_user_story_feedback,
)
from story_devops_index import sync_story_devops_index
from story_index_admin import get_story_lane_index_status
from story_knowledge_assets import (
    create_story_knowledge_assets_from_bundle,
    create_story_knowledge_asset_from_text,
    create_story_knowledge_asset_from_upload,
    list_story_knowledge_assets,
    review_story_knowledge_asset,
)
from story_knowledge_index import sync_story_knowledge_index
from speech_prompt import normalize_spoken_prompt
from privacy_service import build_user_privacy_export, delete_user_personal_data
from provider_governance import evaluate_provider_governance

# =============================================================================
# APP SETUP
# =============================================================================

security = HTTPBearer(auto_error=False)
_allowed_origins = [o.strip().rstrip("/") for o in ALLOWED_ORIGINS.split(",") if o.strip()]
_allowed_origins_set = set(_allowed_origins)
_AUTH_EXEMPT_PATHS = {"/health", "/api/info", "/api/client-error", "/docs", "/openapi.json", "/redoc"}
logger = logging.getLogger(__name__)
BASE_DIR = Path(__file__).resolve().parent
_CHAT_BUDGET_LIMIT = f"{max(1, int(CHAT_BUDGET_PER_MINUTE or 10))}/minute"
_inline_worker_task: Optional[asyncio.Task] = None
_inline_export_worker_task: Optional[asyncio.Task] = None
_upload_retention_task: Optional[asyncio.Task] = None


def _client_ip(request: Request) -> str:
    xff = (request.headers.get("x-forwarded-for") or "").strip()
    if xff:
        return xff.split(",")[0].strip() or "unknown"
    if request.client and request.client.host:
        return request.client.host
    return "unknown"


def _extract_request_token(request: Request) -> str:
    cookie_token = (request.cookies.get(AUTH_COOKIE_NAME) or "").strip()
    authz = (request.headers.get("authorization") or "").strip()
    bearer_token = ""
    if authz.lower().startswith("bearer "):
        bearer_token = authz.split(" ", 1)[1].strip()
    return cookie_token or bearer_token


def _auth_payload_from_request(request: Request) -> dict:
    cached_payload = getattr(request.state, "auth_payload", None)
    if isinstance(cached_payload, dict) and cached_payload:
        return cached_payload
    if str(getattr(request.state, "auth_error", "") or "").strip():
        return {}

    token = _extract_request_token(request)
    if not token:
        return {}
    try:
        return jwt_decode(token)
    except Exception:
        return {}


def _is_admin_user(user: Optional[dict]) -> bool:
    return principal_is_admin(user)


async def _conversation_belongs_to_user(conversation_id: str, user_sub: str) -> bool:
    safe_conv = str(conversation_id or "").strip()
    safe_user = str(user_sub or "").strip()
    if not safe_conv or not safe_user:
        return False
    meta = conversation_meta.get(safe_conv, {})
    owner_sub = str((meta or {}).get("owner_sub", "") or (meta or {}).get("partition_key", "") or "").strip()
    if owner_sub:
        return owner_sub == safe_user
    try:
        rows = await table_query(
            "ChatHistory",
            f"PartitionKey eq '{odata_escape(safe_user)}' and RowKey eq '{odata_escape(safe_conv)}'",
            top=1,
        )
        return bool(rows)
    except Exception:
        return False


def _login_rate_key(request: Request) -> str:
    return f"ip:{_client_ip(request)}"


def _user_or_ip_rate_key(request: Request) -> str:
    payload = _auth_payload_from_request(request)
    sub = str(payload.get("sub", "")).strip()
    if sub:
        return f"user:{sub}"
    return f"ip:{_client_ip(request)}"


@dataclass(frozen=True)
class _RateLimitRule:
    max_requests: int
    window_seconds: int
    key_func: Callable[[Request], str]
    scope: Optional[str] = None


class _DecoratorRateLimiter:
    """Compat layer para manter @limiter.limit e @limiter.shared_limit."""

    def __init__(self, key_func: Callable[[Request], str]):
        self._default_key_func = key_func

    @staticmethod
    def _parse_limit(limit_spec: str) -> Tuple[int, int]:
        text = str(limit_spec or "").strip().lower()
        match = re.fullmatch(r"(\d+)\s*/\s*(second|seconds|minute|minutes|hour|hours|day|days)", text)
        if not match:
            raise ValueError(f"Rate limit inválido: {limit_spec}")
        amount = int(match.group(1))
        unit = match.group(2)
        factor = 1
        if unit.startswith("minute"):
            factor = 60
        elif unit.startswith("hour"):
            factor = 3600
        elif unit.startswith("day"):
            factor = 86400
        return amount, factor

    def limit(self, limit_spec: str, key_func: Optional[Callable[[Request], str]] = None):
        max_requests, window_seconds = self._parse_limit(limit_spec)
        resolved_key_func = key_func or self._default_key_func

        def decorator(fn):
            setattr(
                fn,
                "__dbde_rate_limit__",
                _RateLimitRule(
                    max_requests=max_requests,
                    window_seconds=window_seconds,
                    key_func=resolved_key_func,
                    scope=None,
                ),
            )
            return fn

        return decorator

    def shared_limit(
        self,
        limit_spec: str,
        scope: str,
        key_func: Optional[Callable[[Request], str]] = None,
    ):
        max_requests, window_seconds = self._parse_limit(limit_spec)
        resolved_key_func = key_func or self._default_key_func
        shared_scope = str(scope or "").strip()

        def decorator(fn):
            setattr(
                fn,
                "__dbde_rate_limit__",
                _RateLimitRule(
                    max_requests=max_requests,
                    window_seconds=window_seconds,
                    key_func=resolved_key_func,
                    scope=shared_scope if shared_scope else None,
                ),
            )
            return fn

        return decorator

    @staticmethod
    def resolve(request: Request) -> Optional[_RateLimitRule]:
        route = request.scope.get("route")
        endpoint = getattr(route, "endpoint", None)
        if endpoint is None:
            return None
        return getattr(endpoint, "__dbde_rate_limit__", None)


def _request_is_https(request: Request) -> bool:
    proto = (request.headers.get("x-forwarded-proto") or "").split(",")[0].strip().lower()
    return request.url.scheme == "https" or proto == "https"


def _normalize_origin_value(origin: str) -> str:
    raw = str(origin or "").strip().rstrip("/")
    if not raw:
        return ""
    parsed = urlsplit(raw)
    if not parsed.scheme or not parsed.netloc:
        return raw
    scheme = parsed.scheme.lower()
    host = (parsed.hostname or "").strip().lower()
    if not host:
        return raw
    port = parsed.port
    default_port = 443 if scheme == "https" else 80 if scheme == "http" else None
    suffix = "" if port in (None, default_port) else f":{port}"
    return f"{scheme}://{host}{suffix}"


def _request_origin(request: Request) -> str:
    host = (request.headers.get("host") or request.url.netloc or "").strip()
    if not host:
        return ""
    scheme = "https" if _request_is_https(request) else request.url.scheme
    return _normalize_origin_value(f"{scheme}://{host}")


def _origin_allowed_for_request(request: Request, origin: str) -> bool:
    normalized_origin = _normalize_origin_value(origin)
    if not normalized_origin:
        return True
    if normalized_origin in _allowed_origins_set:
        return True
    request_origin = _request_origin(request)
    return bool(request_origin and normalized_origin == request_origin)


_rate_limiter_backend = TableStorageRateLimit()
_last_rate_cache_cleanup = 0.0
_last_blacklist_cleanup = 0.0
limiter = _DecoratorRateLimiter(key_func=_user_or_ip_rate_key)

@contextlib.asynccontextmanager
async def lifespan(app_instance: FastAPI):
    await startup_event()
    try:
        yield
    finally:
        await shutdown_event()


app = FastAPI(
    title=APP_TITLE,
    version=APP_VERSION,
    description=APP_DESCRIPTION,
    lifespan=lifespan,
)
app.state.limiter = limiter
_static_dir = Path(__file__).resolve().parent / "static"
if _static_dir.exists():
    app.mount("/static", StaticFiles(directory=str(_static_dir)), name="static")
    _dist_dir = _static_dir / "dist"
    if _dist_dir.exists():
        app.mount("/dist", StaticFiles(directory=str(_dist_dir)), name="dist")
app.add_middleware(
    CORSMiddleware,
    allow_origins=_allowed_origins, allow_credentials=True,
    allow_methods=["GET", "POST", "PUT", "PATCH", "DELETE", "OPTIONS"],
    allow_headers=["Authorization", "Content-Type", "X-Requested-With"],
)

http_client: Optional[httpx.AsyncClient] = None
MAX_REQUEST_BODY_BYTES = 15 * 1024 * 1024
UPLOAD_MULTIPART_OVERHEAD_BYTES = 4 * 1024 * 1024
MAX_SINGLE_UPLOAD_REQUEST_BODY_BYTES = (
    max(
        int(UPLOAD_MAX_FILE_BYTES),
        int(get_tabular_upload_limit_bytes("sample.csv", UPLOAD_MAX_FILE_BYTES)),
        int(get_tabular_upload_limit_bytes("sample.tsv", UPLOAD_MAX_FILE_BYTES)),
        int(get_tabular_upload_limit_bytes("sample.xlsx", UPLOAD_MAX_FILE_BYTES)),
        int(get_tabular_upload_limit_bytes("sample.xlsb", UPLOAD_MAX_FILE_BYTES)),
        int(get_tabular_upload_limit_bytes("sample.xls", UPLOAD_MAX_FILE_BYTES)),
    )
    + UPLOAD_MULTIPART_OVERHEAD_BYTES
)
MAX_BATCH_UPLOAD_REQUEST_BODY_BYTES = int(UPLOAD_MAX_BATCH_TOTAL_BYTES) + UPLOAD_MULTIPART_OVERHEAD_BYTES


def _request_body_limit_bytes(request: Request) -> int:
    req_path = str(getattr(getattr(request, "url", None), "path", "") or "").strip()
    if req_path == "/upload/stream/async":
        encoded_filename = str(request.headers.get("x-upload-filename", "") or "").strip()
        filename = unquote(encoded_filename) if encoded_filename else ""
        if filename:
            return max(MAX_REQUEST_BODY_BYTES, _max_upload_bytes_for_file(filename))
        return MAX_REQUEST_BODY_BYTES
    if req_path in {"/upload", "/upload/async"}:
        return max(MAX_REQUEST_BODY_BYTES, MAX_SINGLE_UPLOAD_REQUEST_BODY_BYTES)
    if req_path == "/upload/batch/async":
        return max(MAX_REQUEST_BODY_BYTES, MAX_BATCH_UPLOAD_REQUEST_BODY_BYTES)
    return MAX_REQUEST_BODY_BYTES


@app.middleware("http")
async def enforce_allowed_origins(request: Request, call_next):
    global _last_rate_cache_cleanup, _last_blacklist_cleanup
    req_path = str(request.url.path or "").strip()
    is_exempt_path = req_path in _AUTH_EXEMPT_PATHS or req_path == "/health"

    token_ref = set_request_cookie_token((request.cookies.get(AUTH_COOKIE_NAME) or "").strip())
    auth_payload_ref = set_request_auth_payload(None)
    auth_error_ref = set_request_auth_error("")
    try:
        request.state.auth_payload = None
        request.state.auth_error = ""

        body_limit_bytes = _request_body_limit_bytes(request)
        content_length = (request.headers.get("content-length") or "").strip()
        if content_length:
            try:
                if int(content_length) > body_limit_bytes:
                    return JSONResponse(
                        status_code=413,
                        content={"detail": f"Payload demasiado grande (máximo {body_limit_bytes} bytes)"},
                    )
            except ValueError:
                pass
        received_body_size = 0
        original_receive = request._receive

        async def _size_limited_receive():
            nonlocal received_body_size
            message = await original_receive()
            if message.get("type") == "http.request":
                body_chunk = message.get("body", b"")
                received_body_size += len(body_chunk)
                if received_body_size > body_limit_bytes:
                    raise HTTPException(
                        status_code=413,
                        detail=f"Payload demasiado grande (máximo {body_limit_bytes} bytes)",
                    )
            return message

        request._receive = _size_limited_receive

        request_token = _extract_request_token(request)
        if request_token:
            resolved_payload, resolved_error = await validate_request_token(request_token)
            if resolved_payload:
                request.state.auth_payload = resolved_payload
                set_request_auth_payload(resolved_payload)
            elif resolved_error:
                request.state.auth_error = resolved_error
                set_request_auth_error(resolved_error)

        if not is_exempt_path:
            origin = request.headers.get("origin")
            if origin and _allowed_origins_set and not _origin_allowed_for_request(request, origin):
                return JSONResponse(status_code=403, content={"detail": "Origin não permitida"})

        now = time.time()
        if now - _last_rate_cache_cleanup > 60:
            _rate_limiter_backend.cleanup_local_cache()
            _last_rate_cache_cleanup = now
        if now - _last_blacklist_cleanup > 300:
            cleanup_blacklist()
            _last_blacklist_cleanup = now

        rule = limiter.resolve(request)
        if rule:
            route = request.scope.get("route")
            route_template = str(getattr(route, "path", request.url.path) or request.url.path)
            scope = rule.scope or f"route:{route_template}"
            rate_key = f"{scope}:{rule.key_func(request)}"
            limited = await _rate_limiter_backend.is_rate_limited(
                key=rate_key,
                limit=rule.max_requests,
                window_seconds=rule.window_seconds,
            )
            if limited:
                retry_after = max(1, int(rule.window_seconds))
                return JSONResponse(
                    status_code=429,
                    content={"detail": f"Limite de pedidos excedido. Tenta novamente em {retry_after} segundos."},
                    headers={"Retry-After": str(retry_after)},
                )
        try:
            response = await call_next(request)
        except HTTPException as exc:
            if exc.status_code == 413:
                return JSONResponse(status_code=413, content={"detail": str(exc.detail)})
            raise
        response.headers.setdefault("X-Content-Type-Options", "nosniff")
        response.headers.setdefault("X-Frame-Options", "DENY")
        response.headers.setdefault("Referrer-Policy", "strict-origin-when-cross-origin")
        response.headers.setdefault("Permissions-Policy", _build_permissions_policy())
        path = request.url.path or ""
        if not (path.startswith("/docs") or path.startswith("/redoc")):
            response.headers.setdefault("Content-Security-Policy", _build_content_security_policy())
        if _request_is_https(request):
            response.headers.setdefault("Strict-Transport-Security", "max-age=31536000; includeSubDomains")
        return response
    finally:
        reset_request_auth_error(auth_error_ref)
        reset_request_auth_payload(auth_payload_ref)
        reset_request_cookie_token(token_ref)


def _file_sha256(path: Path) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        while True:
            chunk = f.read(1024 * 1024)
            if not chunk:
                break
            h.update(chunk)
    return h.hexdigest()


def _load_pid_from_file(path: str) -> Optional[int]:
    try:
        txt = Path(path).read_text(encoding="utf-8").strip()
        if not txt:
            return None
        pid = int(txt)
        return pid if pid > 0 else None
    except Exception:
        return None


def _is_process_alive(pid: Optional[int]) -> bool:
    if not pid:
        return False
    try:
        os.kill(pid, 0)
        return True
    except Exception:
        return False


def _build_permissions_policy() -> str:
    # O microfone é necessário para a funcionalidade de voz do chat.
    return "camera=(), microphone=(self), geolocation=()"


def _build_content_security_policy() -> str:
    connect_src = ["'self'"]
    if AZURE_SPEECH_ENABLED and AZURE_SPEECH_REGION:
        region = AZURE_SPEECH_REGION.strip().lower()
        connect_src.extend(
            [
                f"https://{region}.stt.speech.microsoft.com",
                f"wss://{region}.stt.speech.microsoft.com",
                f"https://{region}.tts.speech.microsoft.com",
                "https://*.speech.microsoft.com",
                "wss://*.speech.microsoft.com",
            ]
        )
    return (
        "default-src 'self'; "
        "script-src 'self' https://cdnjs.cloudflare.com https://cdn.plot.ly; "
        "style-src 'self' 'unsafe-inline' https://fonts.googleapis.com; "
        "font-src 'self' https://fonts.gstatic.com; "
        "img-src 'self' data: https:; "
        f"connect-src {' '.join(connect_src)};"
    )


def _chat_partition_key_for_user(user: Optional[dict]) -> str:
    raw = (user or {}).get("sub") or (user or {}).get("username") or "anon"
    safe = "".join(c if c.isalnum() or c in "._-@" else "_" for c in str(raw))
    return (safe or "anon")[:100]


async def _load_conversation_messages_for_export(conv_id: str, user: Optional[dict]) -> List[dict]:
    safe_conv = odata_escape(conv_id)
    safe_pk = odata_escape(_chat_partition_key_for_user(user))
    rows = await table_query(
        "ChatHistory",
        f"PartitionKey eq '{safe_pk}' and RowKey eq '{safe_conv}'",
        top=1,
    )
    if not rows:
        return []
    raw_messages = rows[0].get("Messages", "[]")
    try:
        parsed = json.loads(raw_messages)
    except Exception:
        return []
    return parsed if isinstance(parsed, list) else []


async def _extract_export_data_from_tool_message(tool_msg: dict) -> Optional[dict]:
    blob_ref = str(tool_msg.get("result_blob_ref", "") or "").strip()
    if blob_ref:
        container, blob_name = parse_blob_ref(blob_ref)
        if container and blob_name:
            try:
                payload = await blob_download_json(container, blob_name)
                if isinstance(payload, dict) and payload:
                    return payload
            except Exception as e:
                logger.warning("[Export] blob load failed (%s): %s", blob_ref, e)

    raw_content = tool_msg.get("content", "")
    if isinstance(raw_content, str):
        try:
            parsed = json.loads(raw_content)
            if isinstance(parsed, dict) and parsed and not parsed.get("_persisted_summary"):
                return parsed
        except (json.JSONDecodeError, TypeError):
            pass
    return None


def _safe_export_title(raw_title: str) -> str:
    safe = "".join(c if c.isalnum() or c in " _-" else "_" for c in str(raw_title or "Export DBDE"))
    safe = safe.strip()
    return (safe or "Export_DBDE")[:60]


def _export_rows_count(data: dict) -> int:
    if not isinstance(data, dict):
        return 0
    if isinstance(data.get("items"), list):
        return len(data.get("items") or [])
    if isinstance(data.get("groups"), list):
        return len(data.get("groups") or [])
    if isinstance(data.get("timeline"), list):
        return len(data.get("timeline") or [])
    if isinstance(data.get("analysis_data"), list):
        return len(data.get("analysis_data") or [])
    return 0


async def _resolve_export_payload(export_request: ExportRequest, user: dict) -> dict:
    data = None

    blob_ref = str(export_request.result_blob_ref or "").strip()
    if blob_ref:
        container, blob_name = parse_blob_ref(blob_ref)
        if container and blob_name:
            try:
                payload = await blob_download_json(container, blob_name)
                if isinstance(payload, dict) and payload:
                    data = payload
            except Exception as e:
                logger.warning("[Export] explicit blob load failed (%s): %s", blob_ref, e)

    if data is None and isinstance(export_request.data, dict) and export_request.data:
        data = export_request.data

    if data is None:
        conv_id = str(export_request.conversation_id or "").strip()
        if not conv_id:
            raise HTTPException(400, "Conversa não indicada para export.")

        messages = conversations.get(conv_id, [])
        if not messages:
            messages = await _load_conversation_messages_for_export(conv_id, user)
        if not messages:
            raise HTTPException(
                400,
                "Conversa não encontrada nesta instância nem em persistência. "
                "Volta a enviar uma mensagem nesta conversa e tenta novamente.",
            )

        tool_msgs = [m for m in messages if m.get("role") == "tool"]
        idx = export_request.tool_call_index if export_request.tool_call_index is not None else -1
        if abs(idx) > len(tool_msgs):
            raise HTTPException(400, "Tool result não encontrado")
        selected = tool_msgs[idx]
        data = await _extract_export_data_from_tool_message(selected)
        if not data:
            raise HTTPException(
                400,
                "Dados de export não disponíveis para este resultado (histórico antigo/truncado). "
                "Reexecuta a query para gerar um resultado exportável persistido.",
            )

    if not data:
        raise HTTPException(400, "Sem dados para exportar")

    return data


def _build_export_zip_bytes(data: dict, title: str, summary: str = "") -> bytes:
    safe_title = _safe_export_title(title)
    csv_buf = to_csv(data)
    xlsx_buf = to_xlsx(data, safe_title)
    pdf_buf = to_pdf(data, safe_title, summary or "")

    out = io.BytesIO()
    with zipfile.ZipFile(out, mode="w", compression=zipfile.ZIP_DEFLATED) as zf:
        zf.writestr(f"{safe_title}.csv", csv_buf.getvalue())
        zf.writestr(f"{safe_title}.xlsx", xlsx_buf.getvalue())
        zf.writestr(f"{safe_title}.pdf", pdf_buf.getvalue())
        manifest = {
            "title": safe_title,
            "created_at": datetime.now(timezone.utc).isoformat(),
            "formats": ["csv", "xlsx", "pdf"],
            "rows": _export_rows_count(data),
        }
        zf.writestr("manifest.json", json.dumps(manifest, ensure_ascii=False, indent=2).encode("utf-8"))
    out.seek(0)
    return out.getvalue()


def _build_export_file(format_name: str, data: dict, title: str, summary: str = "") -> tuple[bytes, str, str]:
    fmt = str(format_name or "").strip().lower()
    safe_title = _safe_export_title(title)
    if fmt == "csv":
        return to_csv(data).getvalue(), "text/csv", f"{safe_title}.csv"
    if fmt == "xlsx":
        return (
            to_xlsx(data, safe_title).getvalue(),
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            f"{safe_title}.xlsx",
        )
    if fmt == "pdf":
        return to_pdf(data, safe_title, summary or "").getvalue(), "application/pdf", f"{safe_title}.pdf"
    if fmt == "svg":
        return to_svg_bar_chart(data, safe_title).encode("utf-8"), "image/svg+xml", f"{safe_title}.svg"
    if fmt == "html":
        return to_html_report(data, safe_title, summary or "").encode("utf-8"), "text/html; charset=utf-8", f"{safe_title}.html"
    if fmt == "zip":
        return _build_export_zip_bytes(data, safe_title, summary or ""), "application/zip", f"{safe_title}.zip"
    raise HTTPException(400, f"Formato não suportado: {fmt}")


def _render_chat_segment_to_html(raw_content: Any) -> str:
    if isinstance(raw_content, str):
        text = raw_content
    elif isinstance(raw_content, list):
        chunks = []
        for part in raw_content:
            if isinstance(part, dict) and part.get("type") == "text":
                chunks.append(str(part.get("text", "")))
        text = "\n".join(chunks)
    else:
        text = str(raw_content or "")

    safe = html.escape(text)
    parts = safe.split("```")
    if len(parts) == 1:
        return safe.replace("\n", "<br>")

    rendered = []
    for idx, part in enumerate(parts):
        if idx % 2 == 0:
            rendered.append(part.replace("\n", "<br>"))
        else:
            rendered.append(f"<pre><code>{part}</code></pre>")
    return "".join(rendered)


def _render_chat_html(messages: list, title: str) -> str:
    safe_title = html.escape(_safe_export_title(title or "Chat Export"))
    exported_at = datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M:%S UTC")
    rows = []
    for msg in (messages or []):
        if not isinstance(msg, dict):
            continue
        role = str(msg.get("role", "")).strip().lower()
        if role not in ("user", "assistant"):
            continue
        content_html = _render_chat_segment_to_html(msg.get("content", ""))
        role_label = "Utilizador" if role == "user" else "Assistente"
        bubble_class = "msg-user" if role == "user" else "msg-assistant"
        msg_ts = str(msg.get("timestamp", "") or msg.get("created_at", "") or "").strip()
        ts_html = f"<div class='msg-ts'>{html.escape(msg_ts)}</div>" if msg_ts else ""
        rows.append(
            f"<div class='msg-row {bubble_class}'>"
            f"<div class='msg-meta'>{role_label}</div>"
            f"<div class='msg-content'>{content_html}</div>"
            f"{ts_html}"
            "</div>"
        )

    body = "\n".join(rows) if rows else "<div class='msg-empty'>Sem mensagens para exportar.</div>"
    return f"""<!doctype html>
<html lang="pt">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>{safe_title}</title>
  <style>
    body {{ font-family: 'Montserrat', -apple-system, BlinkMacSystemFont, sans-serif; background:#f6f7fb; color:#1f2937; margin:0; }}
    .wrap {{ max-width: 980px; margin: 24px auto; padding: 0 16px; }}
    .header {{ background:#fff; border:1px solid #e5e7eb; border-radius:12px; padding:16px 18px; margin-bottom:16px; }}
    .title {{ margin:0; font-size:20px; color:{EXPORT_BRAND_COLOR}; }}
    .meta {{ margin-top:6px; font-size:12px; color:#6b7280; }}
    .chat {{ display:flex; flex-direction:column; gap:12px; }}
    .msg-row {{ max-width:80%; border-radius:14px; padding:10px 12px; border:1px solid #e5e7eb; background:#fff; }}
    .msg-user {{ margin-left:auto; background:#eef6ff; border-color:#bfdbfe; }}
    .msg-assistant {{ margin-right:auto; background:#f3f4f6; border-color:#e5e7eb; }}
    .msg-meta {{ font-size:11px; font-weight:700; letter-spacing:0.02em; color:#6b7280; margin-bottom:6px; text-transform:uppercase; }}
    .msg-content {{ font-size:14px; line-height:1.55; }}
    .msg-content pre {{ background:#111827; color:#f9fafb; padding:10px; border-radius:8px; overflow:auto; }}
    .msg-ts {{ margin-top:6px; font-size:11px; color:#9ca3af; }}
    .footer {{ margin-top:18px; font-size:12px; color:#6b7280; text-align:center; }}
  </style>
</head>
<body>
  <div class="wrap">
    <div class="header">
      <h1 class="title">{safe_title}</h1>
      <div class="meta">{html.escape(EXPORT_BRAND_NAME)} · exportado em {exported_at}</div>
    </div>
    <div class="chat">{body}</div>
    <div class="footer">Exportado por {html.escape(EXPORT_AGENT_NAME)} v{APP_VERSION}</div>
  </div>
</body>
</html>"""


def _export_job_from_storage_row(row: dict, fallback_job_id: str = "") -> dict:
    result_payload = None
    raw_result = row.get("ResultJson")
    if raw_result:
        try:
            result_payload = json.loads(raw_result)
        except Exception:
            result_payload = None
    return {
        "job_id": row.get("RowKey", fallback_job_id),
        "status": row.get("Status", "queued"),
        "user_sub": row.get("UserSub", ""),
        "conversation_id": row.get("ConversationId", ""),
        "format": row.get("Format", "xlsx"),
        "title": row.get("Title", "Export DBDE"),
        "summary": row.get("Summary", ""),
        "row_count": int(row.get("RowCount", 0) or 0),
        "payload_blob_ref": row.get("PayloadBlobRef", ""),
        "worker_id": row.get("WorkerId", ""),
        "claim_token": row.get("ClaimToken", ""),
        "created_at": row.get("CreatedAt", ""),
        "updated_at": row.get("UpdatedAt", ""),
        "started_at": row.get("StartedAt", ""),
        "finished_at": row.get("FinishedAt", ""),
        "error": row.get("Error", ""),
        "result": result_payload,
    }


def _export_job_public_view(job: dict) -> dict:
    return {
        "job_id": job.get("job_id"),
        "status": job.get("status"),
        "format": job.get("format"),
        "title": job.get("title"),
        "row_count": int(job.get("row_count", 0) or 0),
        "created_at": job.get("created_at"),
        "updated_at": job.get("updated_at"),
        "started_at": job.get("started_at"),
        "finished_at": job.get("finished_at"),
        "worker_id": job.get("worker_id", ""),
        "error": job.get("error", ""),
        "result": job.get("result"),
    }


def _cleanup_export_jobs() -> None:
    now = datetime.now(timezone.utc)
    expired_ids = []
    for job_id, meta in export_jobs_store.items():
        updated_at = meta.get("updated_at", meta.get("created_at"))
        ts = _parse_iso_dt(str(updated_at or ""))
        if not ts:
            continue
        if (now - ts).total_seconds() > EXPORT_JOB_TTL_SECONDS:
            expired_ids.append(job_id)
    for job_id in expired_ids:
        export_jobs_store.pop(job_id, None)
        create_logged_task(export_jobs_store.delete(job_id), name=f"export_job_delete_{job_id[:8]}")


async def _persist_export_job(job: dict) -> None:
    try:
        job_id = str(job.get("job_id", "") or "")
        if not job_id:
            return
        result_json = json.dumps(job.get("result") or {}, ensure_ascii=False, default=str)
        if len(result_json) > 20000:
            result_json = json.dumps({"_truncated": True, "status": str(job.get("status", ""))}, ensure_ascii=False)
        entity = {
            "PartitionKey": "export",
            "RowKey": job_id,
            "UserSub": str(job.get("user_sub", ""))[:120],
            "ConversationId": str(job.get("conversation_id", ""))[:120],
            "Status": str(job.get("status", "queued"))[:32],
            "Format": str(job.get("format", "xlsx"))[:20],
            "Title": str(job.get("title", "Export DBDE"))[:240],
            "Summary": str(job.get("summary", ""))[:1000],
            "RowCount": int(job.get("row_count", 0) or 0),
            "PayloadBlobRef": str(job.get("payload_blob_ref", ""))[:800],
            "WorkerId": str(job.get("worker_id", ""))[:120],
            "ClaimToken": str(job.get("claim_token", ""))[:120],
            "CreatedAt": str(job.get("created_at", datetime.now(timezone.utc).isoformat())),
            "UpdatedAt": str(job.get("updated_at", datetime.now(timezone.utc).isoformat())),
            "StartedAt": str(job.get("started_at", ""))[:64],
            "FinishedAt": str(job.get("finished_at", ""))[:64],
            "Error": str(job.get("error", ""))[:800],
            "ResultJson": result_json,
        }
        rows = await table_query(
            "ExportJobs",
            f"PartitionKey eq 'export' and RowKey eq '{odata_escape(job_id)}'",
            top=1,
        )
        if rows:
            await table_merge("ExportJobs", entity)
        else:
            await table_insert("ExportJobs", entity)
    except Exception as e:
        logger.warning("[Export] _persist_export_job failed: %s", e)


async def _load_export_job_from_storage(job_id: str) -> Optional[dict]:
    try:
        rows = await table_query(
            "ExportJobs",
            f"PartitionKey eq 'export' and RowKey eq '{odata_escape(job_id)}'",
            top=1,
        )
        if not rows:
            return None
        return _export_job_from_storage_row(rows[0], fallback_job_id=job_id)
    except Exception as e:
        logger.warning("[Export] _load_export_job_from_storage failed for %s: %s", job_id, e)
        return None


async def _load_fresh_export_job_state(job_id: str, in_memory_job: Optional[dict]) -> Optional[dict]:
    storage_job = await _load_export_job_from_storage(job_id)
    if not storage_job:
        return in_memory_job
    if not in_memory_job:
        await export_jobs_store.put(job_id, storage_job)
        return storage_job

    mem_status = str(in_memory_job.get("status", "")).lower()
    stg_status = str(storage_job.get("status", "")).lower()
    mem_ts = _parse_iso_dt(str(in_memory_job.get("updated_at", "") or "")) or _parse_iso_dt(
        str(in_memory_job.get("created_at", "") or "")
    )
    stg_ts = _parse_iso_dt(str(storage_job.get("updated_at", "") or "")) or _parse_iso_dt(
        str(storage_job.get("created_at", "") or "")
    )
    status_promoted = mem_status in ("queued", "processing") and stg_status in ("completed", "failed")
    storage_newer = stg_ts is not None and (mem_ts is None or stg_ts >= mem_ts)
    if status_promoted or storage_newer:
        await export_jobs_store.put(job_id, storage_job)
        return storage_job
    return in_memory_job


async def _queue_export_job(
    user_sub: str,
    conversation_id: str,
    format_name: str,
    title: str,
    summary: str,
    data: dict,
) -> dict:
    job_id = uuid.uuid4().hex
    now_iso = datetime.now(timezone.utc).isoformat()
    blob_name = (
        f"export-jobs/{safe_blob_component(user_sub or 'anon', max_len=90)}/"
        f"{safe_blob_component(job_id, max_len=90)}/payload.json"
    )
    payload = {
        "format": str(format_name or "xlsx").lower(),
        "title": title,
        "summary": summary or "",
        "data": data,
        "created_at": now_iso,
    }
    payload_blob = await blob_upload_json(CHAT_TOOLRESULT_BLOB_CONTAINER, blob_name, payload)
    row_count = _export_rows_count(data)
    job = {
        "job_id": job_id,
        "status": "queued",
        "user_sub": user_sub,
        "conversation_id": conversation_id,
        "format": str(format_name or "xlsx").lower(),
        "title": title,
        "summary": summary or "",
        "row_count": row_count,
        "payload_blob_ref": payload_blob.get("blob_ref", ""),
        "worker_id": "",
        "claim_token": "",
        "created_at": now_iso,
        "updated_at": now_iso,
        "started_at": "",
        "finished_at": "",
        "error": "",
        "result": None,
    }
    await export_jobs_store.put(job_id, job)
    await _persist_export_job(job)
    return job


async def _process_export_job(job: dict) -> None:
    payload_ref = str(job.get("payload_blob_ref", "") or "")
    container, blob_name = parse_blob_ref(payload_ref)
    if not container or not blob_name:
        raise RuntimeError("payload_blob_ref inválido no export job")
    payload = await blob_download_json(container, blob_name)
    if not isinstance(payload, dict):
        raise RuntimeError("payload de export não encontrado")

    fmt = str(payload.get("format", job.get("format", "xlsx"))).lower()
    title = _safe_export_title(str(payload.get("title", job.get("title", "Export DBDE"))))
    summary = str(payload.get("summary", job.get("summary", "")) or "")
    data = payload.get("data")
    if not isinstance(data, dict) or not data:
        raise RuntimeError("dados de export inválidos")

    content, mime_type, filename = _build_export_file(fmt, data, title, summary)
    if not content:
        raise RuntimeError("ficheiro de export vazio")

    download_id = await _store_generated_file(
        content,
        mime_type,
        filename,
        fmt,
        user_sub=str(job.get("user_sub", "") or ""),
        conversation_id=str(job.get("conversation_id", "") or ""),
        scope="export_job",
    )
    if not download_id:
        raise RuntimeError("falha ao persistir ficheiro gerado")

    job["status"] = "completed"
    job["error"] = ""
    job["result"] = {
        "file_generated": True,
        "download_id": download_id,
        "endpoint": f"/api/download/{download_id}",
        "filename": filename,
        "format": fmt,
        "mime_type": mime_type,
        "size_bytes": len(content),
    }


async def _run_export_job(job_id: str) -> None:
    _cleanup_export_jobs()
    job = await export_jobs_store.get_or_fetch(job_id)
    if not job:
        loaded = await _load_export_job_from_storage(job_id)
        if not loaded:
            return
        job = loaded
        await export_jobs_store.put(job_id, job)

    if str(job.get("status", "")).lower() in ("completed", "failed"):
        return

    now_iso = datetime.now(timezone.utc).isoformat()
    job["status"] = "processing"
    job["started_at"] = job.get("started_at") or now_iso
    job["updated_at"] = now_iso
    await _persist_export_job(job)

    async with _export_jobs_semaphore:
        try:
            await _process_export_job(job)
        except Exception as e:
            job["status"] = "failed"
            job["error"] = str(e)
            job["result"] = None
        finally:
            job["finished_at"] = datetime.now(timezone.utc).isoformat()
            job["updated_at"] = datetime.now(timezone.utc).isoformat()
            await export_jobs_store.put(str(job.get("job_id", "")), job)
            await _persist_export_job(job)


async def process_export_jobs_once(max_jobs: int = EXPORT_WORKER_BATCH_SIZE) -> dict:
    _cleanup_export_jobs()
    target = max(1, min(int(max_jobs or 1), 30))
    processed = 0
    claimed = 0
    skipped = 0
    try:
        rows = await table_query(
            "ExportJobs",
            "PartitionKey eq 'export'",
            top=max(target * 5, target),
        )
    except Exception as e:
        logger.warning("[Export] process_export_jobs_once query failed: %s", e)
        rows = []

    rows_sorted = sorted(rows, key=lambda r: str(r.get("CreatedAt", "")))
    for row in rows_sorted:
        if processed >= target:
            break
        job = _export_job_from_storage_row(row)
        job_id = str(job.get("job_id", "") or "")
        if not job_id:
            skipped += 1
            continue
        status = str(job.get("status", "")).lower()
        if status not in ("queued", "processing"):
            skipped += 1
            continue
        if status == "processing" and not _is_export_job_stale(job):
            skipped += 1
            continue

        claim_token = uuid.uuid4().hex
        job["worker_id"] = EXPORT_WORKER_INSTANCE_ID
        job["claim_token"] = claim_token
        job["status"] = "processing"
        job["started_at"] = datetime.now(timezone.utc).isoformat()
        job["updated_at"] = datetime.now(timezone.utc).isoformat()
        await export_jobs_store.put(job_id, job)
        await _persist_export_job(job)
        claimed += 1

        latest = await _load_export_job_from_storage(job_id)
        if latest:
            await export_jobs_store.put(job_id, latest)
            if str(latest.get("claim_token", "")) != claim_token:
                skipped += 1
                continue

        await _run_export_job(job_id)
        processed += 1
    return {"processed": processed, "claimed": claimed, "skipped": skipped}


async def _export_worker_loop() -> None:
    while True:
        try:
            await process_export_jobs_once(max_jobs=EXPORT_WORKER_BATCH_SIZE)
        except asyncio.CancelledError:
            raise
        except Exception as e:
            logger.warning("[Export] worker loop failed: %s", e)
        await asyncio.sleep(max(0.5, float(EXPORT_WORKER_POLL_SECONDS)))

async def startup_event():
    global http_client, _inline_worker_task, _inline_export_worker_task, _upload_retention_task
    configure_logging(LOG_FORMAT)
    http_client = httpx.AsyncClient(timeout=60)
    init_http_client(http_client)
    logger.info("HTTP client OK")

    critical_failures: list[str] = []

    if not AZURE_OPENAI_KEY:
        critical_failures.append("AZURE_OPENAI_KEY não configurada")

    try:
        await asyncio.wait_for(ensure_tables_exist(), timeout=15)
        logger.info("Table Storage OK")
    except asyncio.TimeoutError:
        msg = "Table Storage timeout (15s)"
        logger.warning(msg)
        critical_failures.append(msg)
    except Exception as e:
        msg = f"Table Storage error: {e}"
        logger.warning(msg)
        critical_failures.append(msg)

    try:
        await asyncio.wait_for(ensure_blob_containers(), timeout=15)
        logger.info("Blob Storage OK")
    except asyncio.TimeoutError:
        msg = "Blob Storage timeout (15s)"
        logger.warning(msg)
        critical_failures.append(msg)
    except Exception as e:
        msg = f"Blob Storage error: {e}"
        logger.warning(msg)
        critical_failures.append(msg)

    if critical_failures and STARTUP_FAIL_FAST:
        for f in critical_failures:
            logger.error("[STARTUP FAIL-FAST] %s", f)
        raise RuntimeError(
            f"Startup fail-fast: {len(critical_failures)} critical deps failed: "
            + "; ".join(critical_failures)
        )
    elif critical_failures:
        for f in critical_failures:
            logger.warning("[STARTUP] Non-blocking failure: %s", f)

    _tq_module.token_quota_manager = TokenQuotaManager(TOKEN_QUOTA_CONFIG)
    logger.info("[Startup] Token quota manager initialised: %s", list(TOKEN_QUOTA_CONFIG.keys()))

    _upload_retention_task = create_logged_task(_upload_retention_loop(), name="upload-retention-worker")
    logger.info(
        "Upload retention worker enabled (interval=%ss, retention=%sh, backfill batch=%s, startup=async)",
        UPLOAD_RETENTION_SWEEP_INTERVAL_SECONDS,
        UPLOAD_ARTIFACT_RETENTION_HOURS,
        UPLOAD_TABULAR_CHUNK_BACKFILL_BATCH_SIZE,
    )

    if INLINE_WORKER_ENABLED_EFFECTIVE:
        _inline_worker_task = create_logged_task(_upload_worker_loop(), name="upload-inline-worker")
        logger.info(
            "Inline upload worker enabled (poll=%.1fs, batch=%s)",
            UPLOAD_WORKER_POLL_SECONDS,
            UPLOAD_WORKER_BATCH_SIZE,
        )
    else:
        logger.info(
            "Inline upload worker disabled (configured=%s guard=%s)",
            UPLOAD_INLINE_WORKER_ENABLED,
            INLINE_WORKER_RUNTIME_GUARD,
        )
    if EXPORT_INLINE_WORKER_ENABLED_EFFECTIVE:
        _inline_export_worker_task = create_logged_task(_export_worker_loop(), name="export-inline-worker")
        logger.info(
            "Inline export worker enabled (poll=%.1fs, batch=%s)",
            EXPORT_WORKER_POLL_SECONDS,
            EXPORT_WORKER_BATCH_SIZE,
        )
    else:
        logger.info(
            "Inline export worker disabled (configured=%s guard=%s)",
            EXPORT_INLINE_WORKER_ENABLED,
            INLINE_WORKER_RUNTIME_GUARD,
        )
    logger.info("DBDE AI Agent v%s ready", APP_VERSION)

async def shutdown_event():
    global _inline_worker_task, _inline_export_worker_task, _upload_retention_task
    if _inline_worker_task:
        _inline_worker_task.cancel()
        try:
            await _inline_worker_task
        except (asyncio.CancelledError, Exception):
            pass
        _inline_worker_task = None
    if _inline_export_worker_task:
        _inline_export_worker_task.cancel()
        try:
            await _inline_export_worker_task
        except (asyncio.CancelledError, Exception):
            pass
        _inline_export_worker_task = None
    if _upload_retention_task:
        _upload_retention_task.cancel()
        try:
            await _upload_retention_task
        except (asyncio.CancelledError, Exception):
            pass
        _upload_retention_task = None
    if http_client:
        await http_client.aclose()
    await _close_pii_http_client()
    await _close_knowledge_client()
    await _close_figma_client()
    await _close_miro_client()
    await close_all_providers()

# =============================================================================
# LEARNING / FEW-SHOT HELPERS
# =============================================================================
feedback_memory = deque(maxlen=100)

async def _index_example(example_id, question, answer, rating, tools_used=None, feedback_note="", example_type="positive"):
    try:
        emb = await get_embedding(question)
        if not emb: return
        doc = {"id":example_id,"question":question[:2000],"answer":answer[:4000],"tools_used":",".join(tools_used) if tools_used else "","rating":rating,"feedback_note":feedback_note[:500],"example_type":example_type,"created_at":datetime.now(timezone.utc).isoformat(),"question_vector":emb}
        url = f"https://{SEARCH_SERVICE}.search.windows.net/indexes/{EXAMPLES_INDEX}/docs/index?api-version={API_VERSION_SEARCH}"
        payload = {"value": [{"@search.action": "mergeOrUpload", **doc}]}
        headers = {"api-key": SEARCH_KEY, "Content-Type": "application/json"}
        if http_client:
            await http_client.post(url, json=payload, headers=headers, timeout=30)
        else:
            async with httpx.AsyncClient(timeout=30) as c:
                await c.post(url, json=payload, headers=headers)
    except Exception as e:
        logger.error("[App] _index_example failed: %s", e)

def _audit_clip(value: Any, limit: int) -> str:
    return str(value or "")[: max(1, int(limit or 1))]


async def log_audit(user_id, action, question="", tools_used=None, tokens=None, duration_ms=0, metadata: Optional[dict] = None):
    try:
        ts = datetime.now(timezone.utc)
        safe_meta = metadata if isinstance(metadata, dict) else {}
        governance = evaluate_provider_governance(
            provider_used=safe_meta.get("provider_used", ""),
            model_used=safe_meta.get("model_used", ""),
            action=action,
            mode=safe_meta.get("mode", ""),
            tools_used=tools_used,
        )
        safe_meta = {
            **safe_meta,
            "provider_policy_mode": governance["policy_mode"],
            "provider_family": governance["provider_family"],
            "external_provider": governance["external_provider"],
            "data_sensitivity": governance["data_sensitivity"],
            "provider_policy_note": governance["policy_note"],
        }
        await table_insert(
            "AuditLog",
            {
                "PartitionKey": ts.strftime("%Y-%m"),
                "RowKey": f"{ts.strftime('%Y%m%d%H%M%S')}_{uuid.uuid4().hex[:8]}",
                "UserId": user_id or "anon",
                "Action": _audit_clip(action, 120),
                "Question": _audit_clip(question, 500),
                "ToolsUsed": ",".join(tools_used) if tools_used else "",
                "TotalTokens": tokens.get("total_tokens", 0) if tokens else 0,
                "DurationMs": int(duration_ms or 0),
                "Timestamp": ts.isoformat(),
                "Mode": _audit_clip(safe_meta.get("mode", ""), 32),
                "ModelUsed": _audit_clip(safe_meta.get("model_used", ""), 160),
                "ProviderUsed": _audit_clip(safe_meta.get("provider_used", ""), 80),
                "ProviderFamily": _audit_clip(safe_meta.get("provider_family", ""), 64),
                "ExternalProvider": _audit_clip(safe_meta.get("external_provider", ""), 8),
                "DataSensitivity": _audit_clip(safe_meta.get("data_sensitivity", ""), 32),
                "PolicyMode": _audit_clip(safe_meta.get("provider_policy_mode", ""), 32),
                "ConversationId": _audit_clip(safe_meta.get("conversation_id", ""), 128),
                "Confidence": _audit_clip(safe_meta.get("confidence", ""), 32),
                "MetadataJson": _audit_clip(json.dumps(safe_meta, ensure_ascii=False, default=str), 4000),
            },
        )
    except Exception as e:
        logger.error("[App] log_audit failed: %s", e)

# =============================================================================
# AGENT ENDPOINTS
# =============================================================================

@app.post("/chat/agent", response_model=AgentChatResponse)
@limiter.shared_limit(
    _CHAT_BUDGET_LIMIT,
    scope="chat_budget",
    key_func=_user_or_ip_rate_key,
)
async def agent_chat_endpoint(request: Request, chat_request: AgentChatRequest, credentials: HTTPAuthorizationCredentials = Depends(security)):
    user = get_current_user(credentials)

    result = await _agent_chat(chat_request, user)
    
    # Audit
    try:
        await log_audit(
            user.get("sub"),
            "agent_chat",
            chat_request.question,
            result.tools_used,
            result.tokens_used,
            result.total_time_ms,
            metadata={
                "mode": result.mode,
                "model_used": result.model_used,
                "conversation_id": result.conversation_id,
            },
        )
    except Exception as e:
        logger.error("[App] log_audit in chat failed: %s", e)
    
    return result

@app.post("/chat/agent/stream")
@limiter.shared_limit(
    _CHAT_BUDGET_LIMIT,
    scope="chat_budget",
    key_func=_user_or_ip_rate_key,
)
async def agent_chat_stream_endpoint(request: Request, chat_request: AgentChatRequest, credentials: HTTPAuthorizationCredentials = Depends(security)):
    """SSE streaming endpoint."""
    user = get_current_user(credentials)
    return StreamingResponse(agent_chat_stream(chat_request, user), media_type="text/event-stream")

@app.post("/chat/file", response_model=AgentChatResponse)
@limiter.shared_limit(
    _CHAT_BUDGET_LIMIT,
    scope="chat_budget",
    key_func=_user_or_ip_rate_key,
)
async def chat_with_file(request: Request, chat_request: AgentChatRequest, credentials: HTTPAuthorizationCredentials = Depends(security)):
    """Backward compat."""
    user = get_current_user(credentials)
    result = await _agent_chat(chat_request, user)
    try:
        await log_audit(
            user.get("sub"),
            "chat_file",
            chat_request.question,
            result.tools_used,
            result.tokens_used,
            result.total_time_ms,
            metadata={
                "mode": result.mode,
                "model_used": result.model_used,
                "conversation_id": result.conversation_id,
            },
        )
    except Exception as e:
        logger.error("[App] log_audit in chat_file failed: %s", e)
    return result

# =============================================================================
# MODE SWITCH
# =============================================================================

@app.post("/api/mode/switch", response_model=ModeSwitchResponse)
async def switch_mode(request: ModeSwitchRequest, credentials: HTTPAuthorizationCredentials = Depends(security)):
    get_current_user(credentials)
    ok = await switch_conversation_mode(request.conversation_id, request.mode)
    return ModeSwitchResponse(
        success=ok,
        message=f"Modo alterado para {request.mode}" if ok else "Conversa não encontrada",
        mode=request.mode,
        conversation_id=request.conversation_id,
    )


# =============================================================================
# USER STORY LANE
# =============================================================================

@app.post("/api/user-stories/context-preview")
@limiter.limit("30/minute", key_func=_user_or_ip_rate_key)
async def api_user_story_context_preview(
    request: Request,
    payload: UserStoryWorkspaceRequest,
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    user = get_current_user(credentials, request=request)
    try:
        result = await build_user_story_context_preview(payload.model_dump(), user_sub=str(user.get("sub", "") or ""))
    except RuntimeError as exc:
        raise HTTPException(503, str(exc))
    try:
        await log_audit(
            user.get("sub"),
            "user_story_context_preview",
            payload.objective,
            tools_used=["user_story_lane"],
            duration_ms=sum(int(stage.get("duration_ms", 0) or 0) for stage in result.get("stages", [])),
            metadata={
                "mode": "userstory",
                "conversation_id": payload.conversation_id or "",
                "team_scope": payload.team_scope,
            },
        )
    except Exception as exc:
        logger.warning("[App] user story context audit failed: %s", exc)
    return result


@app.post("/api/user-stories/generate")
@limiter.limit("20/minute", key_func=_user_or_ip_rate_key)
async def api_user_story_generate(
    request: Request,
    payload: UserStoryWorkspaceRequest,
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    user = get_current_user(credentials, request=request)
    try:
        result = await generate_user_story_draft(payload.model_dump(), user_sub=str(user.get("sub", "") or ""))
    except RuntimeError as exc:
        raise HTTPException(400, str(exc))
    try:
        await log_audit(
            user.get("sub"),
            "user_story_generate",
            payload.objective,
            tools_used=["user_story_lane"],
            duration_ms=sum(int(stage.get("duration_ms", 0) or 0) for stage in result.get("stages", [])),
            metadata={
                "mode": "userstory",
                "conversation_id": payload.conversation_id or "",
                "team_scope": payload.team_scope,
            },
        )
    except Exception as exc:
        logger.warning("[App] user story generate audit failed: %s", exc)
    return result


@app.post("/api/user-stories/validate")
@limiter.limit("30/minute", key_func=_user_or_ip_rate_key)
async def api_user_story_validate(
    request: Request,
    payload: UserStoryValidateRequest,
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    user = get_current_user(credentials, request=request)
    try:
        return await validate_user_story_request(payload.model_dump(), user_sub=str(user.get("sub", "") or ""))
    except RuntimeError as exc:
        raise HTTPException(400, str(exc))


@app.post("/api/user-stories/publish")
@limiter.limit("10/minute", key_func=_user_or_ip_rate_key)
async def api_user_story_publish(
    request: Request,
    payload: UserStoryPublishRequest,
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    user = get_current_user(credentials, request=request)
    try:
        result = await publish_user_story(payload.model_dump(), user_sub=str(user.get("sub", "") or ""))
    except RuntimeError as exc:
        raise HTTPException(400, str(exc))
    try:
        await log_audit(
            user.get("sub"),
            "user_story_publish",
            payload.draft_id,
            tools_used=["azure_devops_publish"],
            metadata={"mode": "userstory"},
        )
    except Exception as exc:
        logger.warning("[App] user story publish audit failed: %s", exc)
    return result


@app.post("/api/user-stories/feedback")
@limiter.limit("20/minute", key_func=_user_or_ip_rate_key)
async def api_user_story_feedback(
    request: Request,
    payload: UserStoryFeedbackEventRequest,
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    user = get_current_user(credentials, request=request)
    try:
        return await record_user_story_feedback(payload.model_dump(), user_sub=str(user.get("sub", "") or ""))
    except RuntimeError as exc:
        raise HTTPException(400, str(exc))

# =============================================================================
# FILE UPLOAD
# =============================================================================

MAX_FILES_PER_CONVERSATION = UPLOAD_MAX_FILES_PER_CONVERSATION
MAX_UPLOAD_FILE_BYTES = UPLOAD_MAX_FILE_BYTES
UPLOAD_JOB_TTL_SECONDS = 24 * 3600
MAX_CONCURRENT_UPLOAD_JOBS = UPLOAD_MAX_CONCURRENT_JOBS
WORKER_INSTANCE_ID = os.getenv("UPLOAD_WORKER_INSTANCE_ID") or f"web-{uuid.uuid4().hex[:8]}"
INLINE_WORKER_RUNTIME_GUARD = UPLOAD_INLINE_WORKER_RUNTIME_ENABLED
INLINE_WORKER_ENABLED_EFFECTIVE = bool(UPLOAD_INLINE_WORKER_ENABLED and INLINE_WORKER_RUNTIME_GUARD)
EXPORT_WORKER_INSTANCE_ID = os.getenv("EXPORT_WORKER_INSTANCE_ID") or f"export-web-{uuid.uuid4().hex[:8]}"
EXPORT_INLINE_WORKER_ENABLED_EFFECTIVE = bool(EXPORT_INLINE_WORKER_ENABLED and INLINE_WORKER_RUNTIME_GUARD)
EXPORT_JOB_TTL_SECONDS = 24 * 3600

upload_jobs_store = PersistentJobStore("UploadJobs", partition_key="upload-cache")
_upload_jobs_semaphore = asyncio.Semaphore(MAX_CONCURRENT_UPLOAD_JOBS)
_upload_worker_wake_event = asyncio.Event()
_upload_conv_locks: OrderedDict[str, asyncio.Lock] = OrderedDict()
_UPLOAD_CONV_LOCKS_MAX = 500
export_jobs_store = PersistentJobStore("ExportJobs", partition_key="export-cache")
_export_jobs_semaphore = asyncio.Semaphore(max(1, EXPORT_MAX_CONCURRENT_JOBS))


async def _read_upload_with_limit(upload: UploadFile, max_bytes: int) -> bytes:
    total = 0
    chunks: List[bytes] = []
    while True:
        chunk = await upload.read(1024 * 1024)
        if not chunk:
            break
        total += len(chunk)
        if total > max_bytes:
            raise HTTPException(413, f"Ficheiro excede limite máximo de {max_bytes} bytes")
        chunks.append(chunk)
    return b"".join(chunks)


def _max_upload_bytes_for_file(filename: str) -> int:
    return int(get_tabular_upload_limit_bytes(filename, MAX_UPLOAD_FILE_BYTES))


def _normalize_uploaded_conv_entry(conv_id: str) -> dict:
    current = uploaded_files_store.get(conv_id)
    if isinstance(current, dict) and isinstance(current.get("files"), list):
        return current

    files = []
    if isinstance(current, dict) and current:
        legacy = dict(current)
        legacy.pop("files", None)
        files = [legacy]

    normalized = {
        "files": files,
        "uploaded_at": datetime.now(timezone.utc).isoformat(),
    }
    uploaded_files_store[conv_id] = normalized
    return normalized


def _file_entry_summary(entry: dict) -> dict:
    polymorphic_schema = entry.get("polymorphic_schema")
    is_polymorphic = bool(isinstance(polymorphic_schema, dict) and polymorphic_schema.get("is_polymorphic"))
    return {
        "filename": entry.get("filename", ""),
        "rows": entry.get("row_count", 0),
        "columns": entry.get("col_names", []),
        "truncated": bool(entry.get("truncated")),
        "uploaded_at": entry.get("uploaded_at", ""),
        "has_chunks": isinstance(entry.get("chunks"), list) and len(entry.get("chunks", [])) > 0,
        "is_image": bool(entry.get("image_base64")),
        "polymorphic": is_polymorphic,
        "pivot_column": polymorphic_schema.get("pivot_column", "") if is_polymorphic else "",
        "pivot_values_count": int(polymorphic_schema.get("pivot_values_count", 0) or 0) if is_polymorphic else 0,
    }


def _append_uploaded_entry(conv_id: str, store_entry: dict) -> list:
    conv_entry = _normalize_uploaded_conv_entry(conv_id)
    files = conv_entry.get("files", [])
    if len(files) >= MAX_FILES_PER_CONVERSATION:
        raise HTTPException(
            400,
            f"Limite de {MAX_FILES_PER_CONVERSATION} ficheiros por conversa atingido. Remove alguns anexos e tenta novamente.",
        )
    files.append(store_entry)
    if len(files) > MAX_FILES_PER_CONVERSATION:
        files = files[-MAX_FILES_PER_CONVERSATION:]
    conv_entry["files"] = files
    conv_entry["updated_at"] = datetime.now(timezone.utc).isoformat()
    uploaded_files_store[conv_id] = conv_entry
    return [_file_entry_summary(f) for f in conv_entry.get("files", [])]


def _get_upload_conv_lock(conv_id: str) -> asyncio.Lock:
    lock = _upload_conv_locks.get(conv_id)
    if lock is None:
        while len(_upload_conv_locks) >= _UPLOAD_CONV_LOCKS_MAX:
            _upload_conv_locks.popitem(last=False)
        lock = asyncio.Lock()
        _upload_conv_locks[conv_id] = lock
    else:
        _upload_conv_locks.move_to_end(conv_id)
    return lock


async def _append_uploaded_entry_safe(conv_id: str, store_entry: dict) -> list:
    async with _get_upload_conv_lock(conv_id):
        return _append_uploaded_entry(conv_id, store_entry)


def _cleanup_upload_jobs() -> None:
    now = datetime.now(timezone.utc)
    expired = []
    for job_id, meta in upload_jobs_store.items():
        updated_at = meta.get("updated_at", meta.get("created_at"))
        if not updated_at:
            continue
        try:
            ts = datetime.fromisoformat(str(updated_at))
        except Exception:
            continue
        if (now - ts).total_seconds() > UPLOAD_JOB_TTL_SECONDS:
            expired.append(job_id)
    for job_id in expired:
        upload_jobs_store.pop(job_id, None)
        create_logged_task(upload_jobs_store.delete(job_id), name=f"upload_job_delete_{job_id[:8]}")


async def _count_pending_jobs_for_user(user_sub: str) -> int:
    if not user_sub:
        return 0
    # Source of truth é Table Storage para evitar drift entre instâncias.
    pending_job_ids = set()
    try:
        safe_user = odata_escape(user_sub)
        rows = await table_query(
            "UploadJobs",
            f"PartitionKey eq 'upload' and UserSub eq '{safe_user}'",
            top=500,
        )
        for row in rows:
            status = str(row.get("Status", "")).lower()
            if status in ("queued", "processing"):
                pending_job_ids.add(str(row.get("RowKey", "")))
    except Exception as e:
        logger.warning("[App] _count_pending_jobs_for_user table query failed: %s", e)
        # Fallback local-only quando storage estiver indisponível.
        for job in upload_jobs_store.values():
            if str(job.get("user_sub", "")) != user_sub:
                continue
            if str(job.get("status", "")).lower() in ("queued", "processing"):
                pending_job_ids.add(str(job.get("job_id", "")))
    return len([jid for jid in pending_job_ids if jid])


async def _count_pending_jobs_for_conversation(conv_id: str, user_sub: str = "", include_all_users: bool = False) -> dict:
    counts = {"queued": 0, "processing": 0}
    # Source of truth é Table Storage para evitar overcount por cache stale local.
    seen = set()
    try:
        safe_conv = odata_escape(conv_id)
        rows = await table_query("UploadJobs", f"PartitionKey eq 'upload' and ConversationId eq '{safe_conv}'", top=500)
        for row in rows:
            if not include_all_users and user_sub and str(row.get("UserSub", "")) != user_sub:
                continue
            status = str(row.get("Status", "")).lower()
            if status not in ("queued", "processing"):
                continue
            job_id = str(row.get("RowKey", ""))
            if job_id in seen:
                continue
            seen.add(job_id)
            counts[status] += 1
    except Exception as e:
        logger.warning("[App] _count_pending_jobs_for_conversation table query failed: %s", e)
        # Fallback local-only quando storage estiver indisponível.
        for job in upload_jobs_store.values():
            if str(job.get("conversation_id", "")) != conv_id:
                continue
            if not include_all_users and user_sub and str(job.get("user_sub", "")) != user_sub:
                continue
            status = str(job.get("status", "")).lower()
            if status not in ("queued", "processing"):
                continue
            job_id = str(job.get("job_id", ""))
            if job_id in seen:
                continue
            seen.add(job_id)
            counts[status] += 1
    return {"total": sum(counts.values()), "counts": counts}


def _parse_iso_dt(value: str) -> Optional[datetime]:
    try:
        txt = str(value or "").strip()
        if not txt:
            return None
        parsed = datetime.fromisoformat(txt)
        if parsed.tzinfo is None:
            return parsed.replace(tzinfo=timezone.utc)
        return parsed
    except Exception:
        return None


def _retention_until_iso(*, hours: int) -> str:
    safe_hours = max(1, int(hours or 1))
    return (datetime.now(timezone.utc) + timedelta(hours=safe_hours)).isoformat()


def _is_retention_expired(value: str, now: Optional[datetime] = None) -> bool:
    expires_at = _parse_iso_dt(value)
    if not expires_at:
        return False
    now_dt = now or datetime.now(timezone.utc)
    return now_dt >= expires_at


def _raw_blob_retention_until_iso(
    *,
    filename: str = "",
    artifact_blob_ref: str = "",
    has_chunks: bool = False,
    fallback_hours: int = UPLOAD_ARTIFACT_RETENTION_HOURS,
) -> str:
    safe_fallback_hours = max(1, int(fallback_hours or 1))
    if artifact_blob_ref and is_tabular_filename(str(filename or "")):
        raw_hours = UPLOAD_TABULAR_READY_RAW_RETENTION_HOURS if has_chunks else UPLOAD_TABULAR_RAW_RETENTION_HOURS
        return _retention_until_iso(hours=max(1, int(raw_hours or 1)))
    return _retention_until_iso(hours=safe_fallback_hours)


def _effective_raw_blob_retention_until(row: dict) -> str:
    return str(row.get("RawBlobRetentionUntil", "") or row.get("RetentionUntil", "") or "").strip()


def _is_job_stale(job: dict, now: Optional[datetime] = None) -> bool:
    status = str(job.get("status", "")).lower()
    if status not in ("queued", "processing"):
        return False
    now_dt = now or datetime.now(timezone.utc)
    ref_dt = _parse_iso_dt(str(job.get("updated_at", "") or "")) or _parse_iso_dt(str(job.get("created_at", "") or ""))
    if not ref_dt:
        return False
    return (now_dt - ref_dt).total_seconds() > UPLOAD_JOB_STALE_SECONDS


def _is_export_job_stale(job: dict, now: Optional[datetime] = None) -> bool:
    status = str(job.get("status", "")).lower()
    if status != "processing":
        return False
    now_dt = now or datetime.now(timezone.utc)
    ref_dt = _parse_iso_dt(str(job.get("updated_at", "") or "")) or _parse_iso_dt(str(job.get("created_at", "") or ""))
    if not ref_dt:
        return False
    return (now_dt - ref_dt).total_seconds() > EXPORT_JOB_STALE_SECONDS


def _job_public_view(job: dict) -> dict:
    return {
        "job_id": job.get("job_id"),
        "status": job.get("status"),
        "worker_id": job.get("worker_id", ""),
        "conversation_id": job.get("conversation_id"),
        "filename": job.get("filename"),
        "size_bytes": job.get("size_bytes", 0),
        "created_at": job.get("created_at"),
        "updated_at": job.get("updated_at"),
        "started_at": job.get("started_at"),
        "finished_at": job.get("finished_at"),
        "error": job.get("error"),
        "result": job.get("result"),
    }


async def _load_fresh_job_state(job_id: str, in_memory_job: Optional[dict]) -> Optional[dict]:
    """Prefer latest persisted state to avoid stale queued status under external workers."""
    storage_job = await _load_upload_job_from_storage(job_id)
    if not storage_job:
        return in_memory_job
    if not in_memory_job:
        await upload_jobs_store.put(job_id, storage_job)
        return storage_job

    mem_status = str(in_memory_job.get("status", "")).lower()
    stg_status = str(storage_job.get("status", "")).lower()
    mem_ts = _parse_iso_dt(str(in_memory_job.get("updated_at", "") or "")) or _parse_iso_dt(
        str(in_memory_job.get("created_at", "") or "")
    )
    stg_ts = _parse_iso_dt(str(storage_job.get("updated_at", "") or "")) or _parse_iso_dt(
        str(storage_job.get("created_at", "") or "")
    )

    status_promoted = mem_status in ("queued", "processing") and stg_status in ("completed", "failed")
    storage_newer = stg_ts is not None and (mem_ts is None or stg_ts >= mem_ts)
    if status_promoted or storage_newer:
        await upload_jobs_store.put(job_id, storage_job)
        return storage_job
    return in_memory_job


def _job_from_storage_row(row: dict, fallback_job_id: str = "") -> dict:
    result_payload = None
    raw_result = row.get("ResultJson")
    if raw_result:
        try:
            result_payload = json.loads(raw_result)
        except Exception:
            result_payload = None
    return {
        "job_id": row.get("RowKey", fallback_job_id),
        "status": row.get("Status", "queued"),
        "conversation_id": row.get("ConversationId", ""),
        "filename": row.get("Filename", ""),
        "content_type": row.get("ContentType", ""),
        "size_bytes": int(row.get("SizeBytes", 0) or 0),
        "created_at": row.get("CreatedAt", ""),
        "updated_at": row.get("UpdatedAt", ""),
        "started_at": row.get("StartedAt", ""),
        "finished_at": row.get("FinishedAt", ""),
        "error": row.get("Error", ""),
        "result": result_payload,
        "user_sub": row.get("UserSub", ""),
        "worker_id": row.get("WorkerId", ""),
        "claim_token": row.get("ClaimToken", ""),
        "raw_blob_ref": row.get("RawBlobRef", ""),
        "text_blob_name": row.get("TextBlobName", ""),
        "chunks_blob_name": row.get("ChunksBlobName", ""),
        "artifact_blob_name": row.get("ArtifactBlobName", ""),
        "artifact_blob_ref": row.get("ArtifactBlobRef", ""),
        "artifact_format": row.get("ArtifactFormat", ""),
        "raw_blob_retention_until": row.get("RawBlobRetentionUntil", ""),
        "raw_blob_purged_at": row.get("RawBlobPurgedAt", ""),
        "retention_until": row.get("RetentionUntil", ""),
    }


async def _upsert_upload_index(entity: dict) -> None:
    conv_id = str(entity.get("PartitionKey", "") or "")
    row_key = str(entity.get("RowKey", "") or "")
    if not conv_id or not row_key:
        return
    try:
        rows = await table_query(
            "UploadIndex",
            f"PartitionKey eq '{odata_escape(conv_id)}' and RowKey eq '{odata_escape(row_key)}'",
            top=1,
        )
        if rows:
            await table_merge("UploadIndex", entity)
        else:
            await table_insert("UploadIndex", entity)
    except Exception as e:
        logger.warning("[App] _upsert_upload_index failed for %s/%s: %s", conv_id, row_key, e)


async def _list_upload_index(conv_id: str, user_sub: str = "", top: int = UPLOAD_INDEX_TOP) -> list:
    safe_conv = odata_escape(conv_id)
    rows = await table_query("UploadIndex", f"PartitionKey eq '{safe_conv}'", top=max(1, min(top, 500)))
    if not user_sub:
        return rows
    user_sub_txt = str(user_sub or "")
    filtered = []
    for row in rows:
        owner_sub = str(row.get("UserSub", "") or "")
        if not owner_sub or owner_sub == user_sub_txt:
            filtered.append(row)
    return filtered


def _upload_index_row_to_memory_entry(row: dict) -> dict:
    col_names = []
    col_analysis = []
    polymorphic_schema = None
    try:
        col_names = json.loads(row.get("ColNamesJson", "[]") or "[]")
    except Exception:
        col_names = []
    try:
        col_analysis = json.loads(row.get("ColAnalysisJson", "[]") or "[]")
    except Exception:
        col_analysis = []
    try:
        parsed_schema = json.loads(row.get("PolymorphicSchemaJson", "") or "")
        if isinstance(parsed_schema, dict) and parsed_schema.get("is_polymorphic"):
            polymorphic_schema = parsed_schema
    except Exception:
        polymorphic_schema = None
    if polymorphic_schema is None:
        poly_summary = str(row.get("PolymorphicSummary", "") or "").strip()
        pivot_column = str(row.get("PivotColumn", "") or "").strip()
        if poly_summary:
            polymorphic_schema = {
                "is_polymorphic": True,
                "summary_text": poly_summary,
                "pivot_column": pivot_column,
                "pivot_values_count": int(row.get("PivotValuesCount", 0) or 0),
            }
    return {
        "filename": row.get("Filename", ""),
        "data_text": row.get("PreviewText", ""),
        "row_count": int(row.get("RowCount", 0) or 0),
        "col_names": col_names if isinstance(col_names, list) else [],
        "col_analysis": col_analysis if isinstance(col_analysis, list) else [],
        "truncated": bool(row.get("Truncated", False)),
        "uploaded_at": row.get("UploadedAt", ""),
        "has_chunks": str(row.get("HasChunks", "")).lower() in ("true", "1"),
        "chunks_blob_ref": row.get("ChunksBlobRef", ""),
        "extracted_blob_ref": row.get("ExtractedBlobRef", ""),
        "tabular_artifact_blob_ref": row.get("TabularArtifactBlobRef", ""),
        "tabular_artifact_format": row.get("TabularArtifactFormat", ""),
        "tabular_artifact_row_count": int(row.get("TabularArtifactRowCount", 0) or 0),
        "polymorphic_schema": polymorphic_schema,
    }


async def _refresh_conversation_files_from_index(conv_id: str, user_sub: str = "") -> list:
    rows = await _list_upload_index(conv_id, user_sub=user_sub, top=UPLOAD_INDEX_TOP)
    if not rows:
        return []
    rows_sorted = sorted(rows, key=lambda r: str(r.get("UploadedAt", "")))
    entries = [_upload_index_row_to_memory_entry(r) for r in rows_sorted[-MAX_FILES_PER_CONVERSATION:]]
    uploaded_files_store[conv_id] = {
        "files": entries,
        "uploaded_at": datetime.now(timezone.utc).isoformat(),
        "from_index": True,
    }
    return [_file_entry_summary(e) for e in entries]


async def _count_files_for_conversation(conv_id: str, user_sub: str = "") -> int:
    conv_entry = _normalize_uploaded_conv_entry(conv_id)
    memory_count = len(conv_entry.get("files", []))
    try:
        rows = await _list_upload_index(conv_id, user_sub=user_sub, top=UPLOAD_INDEX_TOP)
        indexed_count = len(rows)
    except Exception:
        indexed_count = 0
    return max(memory_count, indexed_count)


async def _count_reserved_slots_for_conversation(conv_id: str, user_sub: str = "") -> int:
    uploaded = await _count_files_for_conversation(conv_id, user_sub=user_sub)
    pending = await _count_pending_jobs_for_conversation(
        conv_id=conv_id,
        user_sub=user_sub,
        include_all_users=True,
    )
    return uploaded + int((pending or {}).get("total", 0) or 0)


async def _build_semantic_chunks(
    text: str,
    base_chunk_size: int = 4000,
    overlap: int = 200,
    max_chunks: int = UPLOAD_MAX_CHUNKS_PER_FILE,
) -> list:
    if not text or len(text) <= 50000:
        return []

    text_len = len(text)
    effective_chunk_size = max(base_chunk_size, overlap + 1)
    step = effective_chunk_size - overlap
    estimated_chunks = max(1, (text_len + step - 1) // step)
    if estimated_chunks > max_chunks:
        effective_chunk_size = max(overlap + 1, int((text_len / max_chunks) + overlap))
        step = effective_chunk_size - overlap

    raw_chunks = []
    start = 0
    chunk_index = 0
    while start < text_len:
        end = min(start + effective_chunk_size, text_len)
        chunk_text = text[start:end]
        if chunk_text.strip():
            raw_chunks.append((chunk_index, start, end, chunk_text))
        if end >= text_len:
            break
        start = end - overlap
        chunk_index += 1

    async def _embed_chunk(raw_chunk):
        idx, start_pos, end_pos, chunk_text = raw_chunk
        emb = await get_embedding(chunk_text)
        if not emb:
            return None
        return {
            "index": idx,
            "start": start_pos,
            "end": end_pos,
            "text": chunk_text,
            "embedding": emb,
        }

    built_chunks = []
    worker_concurrency = max(1, min(UPLOAD_EMBEDDING_CONCURRENCY, 8))
    batch_size = max(worker_concurrency * 3, 6)
    for i in range(0, len(raw_chunks), batch_size):
        batch = raw_chunks[i : i + batch_size]
        results = await asyncio.gather(*[_embed_chunk(rc) for rc in batch], return_exceptions=True)
        for item in results:
            if isinstance(item, Exception):
                logger.warning("[App] chunk embedding failed: %s", item)
                continue
            if item:
                built_chunks.append(item)
    return built_chunks


async def _build_tabular_semantic_chunks_from_artifact(
    artifact_bytes: bytes,
    *,
    columns: Optional[list[str]] = None,
    rows_per_chunk: int = 120,
    max_chunk_chars: int = 4000,
    max_chunks: int = UPLOAD_MAX_CHUNKS_PER_FILE,
) -> list:
    if not artifact_bytes:
        return []

    safe_rows_per_chunk = max(1, min(int(rows_per_chunk or 0), 500))
    safe_max_chunk_chars = max(50, int(max_chunk_chars or 0))
    safe_max_chunks = max(1, int(max_chunks or 1))
    selected_columns = [str(col or "").strip() for col in (columns or []) if str(col or "").strip()]
    header_line = "\t".join(selected_columns) if selected_columns else ""

    raw_chunks: list[tuple[int, int, int, str]] = []
    chunk_lines: list[str] = [header_line] if header_line else []
    chunk_chars = len(header_line)
    chunk_row_count = 0
    chunk_index = 0
    chunk_start_row = 0
    current_row_number = 0

    for batch in iter_tabular_artifact_batches(
        artifact_bytes,
        columns=selected_columns or None,
        batch_rows=min(max(safe_rows_per_chunk * 2, 200), 5000),
    ):
        for row in batch:
            values = [str((row or {}).get(column, "") or "") for column in (selected_columns or list(row.keys()))]
            line = "\t".join(values)
            projected_chars = chunk_chars + (1 if chunk_lines else 0) + len(line)
            if chunk_row_count > 0 and (
                chunk_row_count >= safe_rows_per_chunk or projected_chars > safe_max_chunk_chars
            ):
                chunk_text = "\n".join(chunk_lines).strip()
                if chunk_text:
                    raw_chunks.append((chunk_index, chunk_start_row, current_row_number, chunk_text))
                if len(raw_chunks) >= safe_max_chunks:
                    break
                chunk_index += 1
                chunk_start_row = current_row_number
                chunk_lines = [header_line] if header_line else []
                chunk_chars = len(header_line)
                chunk_row_count = 0

            if chunk_row_count == 0 and header_line and not chunk_lines:
                chunk_lines = [header_line]
                chunk_chars = len(header_line)

            chunk_lines.append(line)
            chunk_chars += (1 if chunk_chars else 0) + len(line)
            chunk_row_count += 1
            current_row_number += 1
        if len(raw_chunks) >= safe_max_chunks:
            break

    if chunk_lines and chunk_row_count > 0 and len(raw_chunks) < safe_max_chunks:
        chunk_text = "\n".join(chunk_lines).strip()
        if chunk_text:
            raw_chunks.append((chunk_index, chunk_start_row, current_row_number, chunk_text))

    if not raw_chunks:
        return []

    async def _embed_chunk(raw_chunk):
        idx, start_row, end_row, chunk_text = raw_chunk
        emb = await get_embedding(chunk_text)
        if not emb:
            return None
        return {
            "index": idx,
            "start": start_row,
            "end": end_row,
            "text": chunk_text,
            "embedding": emb,
        }

    built_chunks = []
    # Process in small batches to avoid overwhelming the embedding API.
    # The semaphore in get_embedding() handles concurrency limiting;
    # we just batch to avoid creating thousands of tasks at once.
    batch_size = max(1, min(UPLOAD_EMBEDDING_CONCURRENCY, 5))
    for i in range(0, len(raw_chunks), batch_size):
        batch = raw_chunks[i : i + batch_size]
        results = await asyncio.gather(*[_embed_chunk(item) for item in batch], return_exceptions=True)
        for item in results:
            if isinstance(item, Exception):
                logger.warning("[App] tabular artifact chunk embedding failed: %s", item)
                continue
            if item:
                built_chunks.append(item)
    return built_chunks


async def _persist_upload_job(job: dict) -> None:
    try:
        pk = "upload"
        rk = str(job.get("job_id", ""))
        if not rk:
            return
        result_json = json.dumps(job.get("result") or {}, ensure_ascii=False, default=str)
        if len(result_json) > 20000:
            result_json = json.dumps(
                {
                    "_truncated": True,
                    "status": str(job.get("status", "")),
                    "filename": str(job.get("filename", "")),
                    "conversation_id": str(job.get("conversation_id", "")),
                },
                ensure_ascii=False,
                default=str,
            )
        entity = {
            "PartitionKey": pk,
            "RowKey": rk,
            "ConversationId": str(job.get("conversation_id", ""))[:120],
            "UserSub": str(job.get("user_sub", ""))[:120],
            "Filename": str(job.get("filename", ""))[:240],
            "ContentType": str(job.get("content_type", ""))[:120],
            "Status": str(job.get("status", "queued"))[:32],
            "SizeBytes": int(job.get("size_bytes", 0) or 0),
            "CreatedAt": str(job.get("created_at", datetime.now(timezone.utc).isoformat())),
            "UpdatedAt": str(job.get("updated_at", datetime.now(timezone.utc).isoformat())),
            "StartedAt": str(job.get("started_at", ""))[:64],
            "FinishedAt": str(job.get("finished_at", ""))[:64],
            "Error": str(job.get("error", ""))[:800],
            "WorkerId": str(job.get("worker_id", ""))[:120],
            "ClaimToken": str(job.get("claim_token", ""))[:120],
            "RawBlobRef": str(job.get("raw_blob_ref", ""))[:800],
            "TextBlobName": str(job.get("text_blob_name", ""))[:500],
            "ChunksBlobName": str(job.get("chunks_blob_name", ""))[:500],
            "ArtifactBlobName": str(job.get("artifact_blob_name", ""))[:500],
            "ArtifactBlobRef": str(job.get("artifact_blob_ref", ""))[:800],
            "ArtifactFormat": str(job.get("artifact_format", ""))[:32],
            "RawBlobRetentionUntil": str(job.get("raw_blob_retention_until", ""))[:64],
            "RawBlobPurgedAt": str(job.get("raw_blob_purged_at", ""))[:64],
            "RetentionUntil": str(job.get("retention_until", ""))[:64],
            "ResultJson": result_json,
        }
        rows = await table_query(
            "UploadJobs",
            f"PartitionKey eq 'upload' and RowKey eq '{odata_escape(rk)}'",
            top=1,
        )
        if rows:
            await table_merge("UploadJobs", entity)
        else:
            await table_insert("UploadJobs", entity)
    except Exception as e:
        logger.warning("[App] _persist_upload_job failed: %s", e)


async def _load_upload_job_from_storage(job_id: str) -> Optional[dict]:
    try:
        rows = await table_query(
            "UploadJobs",
            f"PartitionKey eq 'upload' and RowKey eq '{odata_escape(job_id)}'",
            top=1,
        )
        if not rows:
            return None
        row = rows[0]
        return _job_from_storage_row(row, fallback_job_id=job_id)
    except Exception as e:
        logger.warning("[App] _load_upload_job_from_storage failed for %s: %s", job_id, e)
        return None


async def _mark_job_failed(job: dict, reason: str) -> dict:
    if not job:
        return job
    if str(job.get("status", "")).lower() in ("completed", "failed"):
        return job
    now_iso = datetime.now(timezone.utc).isoformat()
    job["status"] = "failed"
    job["error"] = reason
    job["finished_at"] = now_iso
    job["updated_at"] = now_iso
    await upload_jobs_store.put(str(job.get("job_id", "")), job)
    await _persist_upload_job(job)
    return job


async def _extract_upload_entry(
    filename: str,
    content: bytes,
    content_type: str = "",
    *,
    prebuilt_artifact: dict | None = None,
) -> tuple[dict, dict]:
    data_text, row_count, col_names, truncated = "", 0, [], False
    semantic_chunks = None
    col_analysis = []
    polymorphic_schema = None
    doc_intel_meta = {}
    image_base64 = None
    image_content_type = None
    detected_delimiter = ","
    full_text = ""
    tabular_ingest_mode = ""
    filename_lower = filename.lower()
    is_pptx_mime = content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    is_image_mime = (content_type or "").startswith("image/")
    ext = os.path.splitext(filename_lower)[1]

    allowed_extensions = {
        ".xlsx", ".xls", ".xlsb", ".csv", ".pdf", ".png", ".jpg", ".jpeg",
        ".gif", ".webp", ".bmp", ".pptx", ".svg", ".txt", ".md",
        ".json", ".xml", ".html", ".htm", ".log", ".tsv",
    }
    if ext and ext not in allowed_extensions and not is_image_mime and not is_pptx_mime:
        raise HTTPException(
            400,
            f"Extensão '{ext}' não permitida. Aceites: {', '.join(sorted(allowed_extensions))}",
        )

    magic_signatures = {
        ".xlsx": (b"PK\x03\x04",),
        ".xlsb": (b"PK\x03\x04",),
        ".xls": (b"\xd0\xcf\x11\xe0",),
        ".pdf": (b"%PDF",),
        ".png": (b"\x89PNG",),
        ".jpg": (b"\xff\xd8\xff",),
        ".jpeg": (b"\xff\xd8\xff",),
        ".gif": (b"GIF87a", b"GIF89a"),
        ".webp": (b"RIFF",),
        ".bmp": (b"BM",),
        ".pptx": (b"PK\x03\x04",),
    }
    if ext in magic_signatures:
        prefix = content[:8]
        if not any(prefix.startswith(sig) for sig in magic_signatures[ext]):
            raise HTTPException(
                400,
                f"Conteúdo do ficheiro não corresponde à extensão '{ext}'. Verifica o ficheiro e tenta novamente.",
            )

    if is_tabular_filename(filename_lower):
        # ── Fast path: if we already built a Parquet artifact, derive
        # preview + dataset from it (DuckDB reads Parquet in ms) instead
        # of re-parsing the raw XLSX/XLS with openpyxl again.
        if prebuilt_artifact and prebuilt_artifact.get("artifact_bytes"):
            artifact_bytes = prebuilt_artifact["artifact_bytes"]
            try:
                preview = await asyncio.to_thread(
                    load_tabular_artifact_preview,
                    artifact_bytes,
                )
            except Exception:
                preview = {}
            col_names = list(prebuilt_artifact.get("columns") or preview.get("columns") or [])
            row_count = int(prebuilt_artifact.get("row_count", 0) or preview.get("row_count", 0) or 0)
            data_text = str(preview.get("data_text", "") or "")
            detected_delimiter = str(preview.get("delimiter", "\t") or "\t")
            col_analysis = list(preview.get("col_analysis") or [])
            sample_records = list(preview.get("sample_records") or [])
            column_types = dict(preview.get("column_types") or {})
            polymorphic_schema = detect_polymorphic_schema(
                col_names,
                sample_records,
                column_types,
                row_count,
            )
            truncated = bool(preview.get("truncated", False))
            # Deep ingest from Parquet (instant via DuckDB)
            should_deep_ingest = row_count <= UPLOAD_TABULAR_DEEP_INGEST_MAX_ROWS
            if row_count > 0 and should_deep_ingest:
                try:
                    full_dataset = await asyncio.to_thread(
                        load_tabular_artifact_dataset,
                        artifact_bytes,
                        max_rows=UPLOAD_TABULAR_DEEP_INGEST_RECORD_LIMIT,
                    )
                    full_records = list(full_dataset.get("records") or [])
                    if full_records and col_names:
                        full_lines = [detected_delimiter.join(col_names)]
                        for rec in full_records:
                            full_lines.append(
                                detected_delimiter.join(str(rec.get(column, "")) for column in col_names)
                            )
                        full_text = "\n".join(full_lines)
                    else:
                        full_text = data_text
                    tabular_ingest_mode = "deep"
                except Exception:
                    full_text = data_text
                    tabular_ingest_mode = "preview_only"
            else:
                full_text = data_text
                tabular_ingest_mode = "preview_only"
                if row_count > len(sample_records):
                    truncated = True
            logger.info(
                "[Upload] artifact-first tabular extract filename=%s rows=%s mode=%s (zero extra openpyxl passes)",
                filename, row_count, tabular_ingest_mode,
            )
        else:
            # ── Legacy path: no pre-built artifact, parse raw bytes directly
            try:
                preview = await asyncio.to_thread(load_tabular_preview, content, filename_lower)
            except TabularLoaderError as exc:
                raise HTTPException(400, str(exc)) from exc
            col_names = list(preview.get("columns") or [])
            row_count = int(preview.get("row_count", 0) or 0)
            data_text = str(preview.get("data_text", "") or "")
            detected_delimiter = str(preview.get("delimiter", "\t") or "\t")
            col_analysis = list(preview.get("col_analysis") or [])
            polymorphic_schema = detect_polymorphic_schema(
                col_names,
                list(preview.get("sample_records") or []),
                dict(preview.get("column_types") or {}),
                row_count,
            )
            truncated = bool(preview.get("truncated", False))
            sample_records = list(preview.get("sample_records") or [])
            should_deep_ingest = (
                len(content) <= UPLOAD_TABULAR_DEEP_INGEST_MAX_BYTES
                and row_count <= UPLOAD_TABULAR_DEEP_INGEST_MAX_ROWS
            )
            if row_count > 0 and should_deep_ingest:
                try:
                    full_dataset = await asyncio.to_thread(
                        load_tabular_dataset,
                        content,
                        filename_lower,
                        UPLOAD_TABULAR_DEEP_INGEST_RECORD_LIMIT,
                    )
                    full_records = list(full_dataset.get("records") or [])
                    if full_records and col_names:
                        full_lines = [detected_delimiter.join(col_names)]
                        for rec in full_records:
                            full_lines.append(
                                detected_delimiter.join(str(rec.get(column, "")) for column in col_names)
                            )
                        full_text = "\n".join(full_lines)
                    else:
                        full_text = data_text
                    tabular_ingest_mode = "deep"
                except Exception:
                    full_text = data_text
                    tabular_ingest_mode = "preview_only"
            else:
                full_text = data_text
                tabular_ingest_mode = "preview_only"
                if row_count > len(sample_records):
                    truncated = True
            if tabular_ingest_mode == "preview_only":
                logger.info(
                    "[Upload] tabular preview-only ingest filename=%s size_bytes=%s rows=%s",
                    filename,
                    len(content),
                    row_count,
                )
    elif filename_lower.endswith(".pdf"):
        used_doc_intel = False
        if DOC_INTEL_ENABLED:
            di_result = await analyze_document(content, filename, model_id=DOC_INTEL_MODEL)
            extracted_text = (di_result.get("text") or "").strip()
            if extracted_text:
                data_text = extracted_text
                table_md = tables_to_markdown(di_result.get("tables") or [])
                if table_md:
                    data_text += "\n\n" + table_md
                key_values = di_result.get("key_values") or []
                if key_values:
                    kv_lines = ["Campos extraidos:"]
                    for kv in key_values[:100]:
                        key = str((kv or {}).get("key") or "").strip()
                        value = str((kv or {}).get("value") or "").strip()
                        if key:
                            kv_lines.append(f"- {key}: {value}")
                    if len(kv_lines) > 1:
                        data_text += "\n\n" + "\n".join(kv_lines)

                row_count = int(di_result.get("page_count") or 0)
                col_names = [f"páginas ({row_count})"] if row_count > 0 else ["páginas"]
                doc_intel_meta = {
                    "enabled": True,
                    "used": True,
                    "model": DOC_INTEL_MODEL,
                    "page_count": int(di_result.get("page_count") or 0),
                    "table_count": int(di_result.get("table_count") or 0),
                    "key_values": len(key_values),
                }
                used_doc_intel = True
            else:
                doc_intel_meta = {
                    "enabled": True,
                    "used": False,
                    "model": DOC_INTEL_MODEL,
                    "error": str(di_result.get("error") or "empty_text"),
                }

        if not used_doc_intel:
            def _parse_pdf_sync(pdf_bytes: bytes) -> tuple[str, int]:
                try:
                    from pypdf import PdfReader
                except ImportError:
                    from PyPDF2 import PdfReader
                reader = PdfReader(io.BytesIO(pdf_bytes))
                pages = [
                    f"[Pág {i+1}]\n{p.extract_text() or ''}"
                    for i, p in enumerate(reader.pages)
                    if (p.extract_text() or "").strip()
                ]
                return "\n\n".join(pages), len(reader.pages)
            data_text, row_count = await asyncio.to_thread(_parse_pdf_sync, content)
            col_names = [f"páginas ({row_count})"]
            if not data_text.strip():
                raise HTTPException(400, "PDF sem texto")
    elif filename_lower.endswith((".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp")) or is_image_mime:
        if is_image_mime:
            image_content_type = content_type
        elif filename_lower.endswith((".jpg", ".jpeg")):
            image_content_type = "image/jpeg"
        elif filename_lower.endswith(".png"):
            image_content_type = "image/png"
        elif filename_lower.endswith(".gif"):
            image_content_type = "image/gif"
        elif filename_lower.endswith(".webp"):
            image_content_type = "image/webp"
        elif filename_lower.endswith(".bmp"):
            image_content_type = "image/bmp"
        else:
            image_content_type = "image/png"
        image_base64 = base64.b64encode(content).decode("ascii")
        data_text = f"[Imagem carregada: {filename}]"
        row_count = 1
        col_names = ["imagem"]
    elif filename_lower.endswith(".pptx") or is_pptx_mime:
        if Presentation is None:
            logger.warning("[App] Upload .pptx recusado: python-pptx não disponível")
            raise HTTPException(503, "Suporte .pptx indisponível no servidor (python-pptx não instalado)")
        def _parse_pptx_sync(pptx_bytes: bytes) -> tuple[str, int]:
            prs = Presentation(io.BytesIO(pptx_bytes))
            slide_blocks = []
            for si, slide in enumerate(prs.slides, 1):
                shape_texts = []
                for shape in slide.shapes:
                    if not getattr(shape, "has_text_frame", False):
                        continue
                    text_frame = getattr(shape, "text_frame", None)
                    if not text_frame:
                        continue
                    paragraphs = []
                    for paragraph in text_frame.paragraphs:
                        txt = (paragraph.text or "").strip()
                        if txt:
                            paragraphs.append(txt)
                    if paragraphs:
                        shape_texts.append("\n".join(paragraphs))
                if shape_texts:
                    slide_blocks.append(f"[Slide {si}]\n" + "\n".join(shape_texts))
            return "\n\n".join(slide_blocks), len(prs.slides)
        data_text, row_count = await asyncio.to_thread(_parse_pptx_sync, content)
        col_names = [f"slides ({row_count})", "texto"]
        if not data_text.strip():
            raise HTTPException(400, "PPTX sem texto legível")
    elif filename_lower.endswith(".svg"):
        data_text = content.decode("utf-8", errors="replace")
        row_count = 0
        col_names = ["svg"]
    else:
        data_text = content.decode("utf-8", errors="replace")
        row_count = data_text.count("\n")
        col_names = ["texto"]

    if not is_tabular_filename(filename_lower):
        full_text = data_text
    # NOTE: semantic chunks are now deferred to background / first-query to avoid
    # blocking the upload response with expensive embedding API calls.
    # if not image_base64 and len(full_text) > 50000 and tabular_ingest_mode != "preview_only":
    #     semantic_chunks = await _build_semantic_chunks(full_text)

    if len(data_text) > 100000:
        data_text = data_text[:100000]
        truncated = True

    store_entry = {
        "filename": filename,
        "data_text": data_text,
        "row_count": row_count,
        "col_names": col_names,
        "truncated": truncated,
        "uploaded_at": datetime.now(timezone.utc).isoformat(),
    }
    if image_base64:
        store_entry["image_base64"] = image_base64
        store_entry["image_content_type"] = image_content_type or "image/png"
    if semantic_chunks is not None:
        store_entry["chunks"] = semantic_chunks
    if doc_intel_meta:
        store_entry["document_intelligence"] = doc_intel_meta
    if is_tabular_filename(filename_lower):
        store_entry["col_analysis"] = col_analysis
        store_entry["tabular_ingest_mode"] = tabular_ingest_mode or "preview_only"
        if polymorphic_schema:
            store_entry["polymorphic_schema"] = polymorphic_schema

    response_payload = {
        "filename": filename,
        "rows": row_count,
        "columns": col_names,
        "truncated": truncated,
        "preview": "\n".join(data_text.split("\n")[:6]),
        "col_analysis": col_analysis if col_analysis else None,
    }
    if doc_intel_meta:
        response_payload["document_intelligence"] = doc_intel_meta
    if polymorphic_schema:
        response_payload["polymorphic"] = True
        response_payload["pivot_column"] = str(polymorphic_schema.get("pivot_column", "") or "")
        response_payload["pivot_values_count"] = int(polymorphic_schema.get("pivot_values_count", 0) or 0)
    if is_tabular_filename(filename_lower):
        response_payload["tabular_ingest_mode"] = tabular_ingest_mode or "preview_only"
    return store_entry, response_payload


def _build_upload_blob_paths(conv_id: str, job_id: str, filename: str) -> dict:
    safe_name = safe_blob_component(filename or "upload.bin", max_len=180)
    prefix = f"{safe_blob_component(conv_id, max_len=90)}/{safe_blob_component(job_id, max_len=90)}"
    return {
        "raw_blob_name": f"{prefix}/raw/{safe_name}",
        "text_blob_name": f"{prefix}/extracted/text.txt",
        "chunks_blob_name": f"{prefix}/extracted/chunks.json",
        "artifact_blob_name": f"{prefix}/artifact/data.parquet",
    }


async def _process_upload_job(job: dict, *, inline_content: bytes | None = None) -> None:
    conv_id = str(job.get("conversation_id", "") or "")
    user_sub = str(job.get("user_sub", "") or "")
    filename = str(job.get("filename", "") or "unknown")
    content_type = str(job.get("content_type", "") or "")
    job_id = str(job.get("job_id", "") or "")
    raw_blob_ref = str(job.get("raw_blob_ref", "") or "")
    container, blob_name = parse_blob_ref(raw_blob_ref)
    if not container or not blob_name:
        # Backward compatibility: recover blob references for rows queued before
        # blob metadata fields were persisted in UploadJobs.
        fallback_paths = _build_upload_blob_paths(conv_id, job_id, filename)
        raw_blob_ref = build_blob_ref(UPLOAD_BLOB_CONTAINER_RAW, fallback_paths["raw_blob_name"])
        container, blob_name = parse_blob_ref(raw_blob_ref)
        if not container or not blob_name:
            raise RuntimeError("raw_blob_ref inválido no upload job")
        job["raw_blob_ref"] = raw_blob_ref
        if not str(job.get("text_blob_name", "") or ""):
            job["text_blob_name"] = fallback_paths["text_blob_name"]
        if not str(job.get("chunks_blob_name", "") or ""):
            job["chunks_blob_name"] = fallback_paths["chunks_blob_name"]

    existing_count = await _count_files_for_conversation(conv_id, user_sub=user_sub)
    if existing_count >= MAX_FILES_PER_CONVERSATION:
        raise RuntimeError(
            f"Limite de {MAX_FILES_PER_CONVERSATION} ficheiros por conversa atingido antes do processamento."
        )

    raw_bytes = inline_content if inline_content is not None else await blob_download_bytes(container, blob_name)
    if raw_bytes is None:
        raise RuntimeError("Blob original não encontrado para este upload job")

    # ── Build tabular artifact FIRST (single openpyxl/CSV pass) so that
    # _extract_upload_entry can derive preview + dataset from the fast
    # Parquet artifact instead of re-parsing the raw file again.
    artifact_blob_ref = ""
    artifact_format = ""
    artifact_row_count = 0
    artifact = {}
    if UPLOAD_TABULAR_ARTIFACT_ENABLED and is_tabular_filename(filename):
        try:
            artifact = await asyncio.to_thread(build_tabular_artifact, raw_bytes, filename)
            artifact_blob_name = str(job.get("artifact_blob_name", "") or "")
            if artifact_blob_name and artifact.get("artifact_bytes"):
                uploaded_artifact = await blob_upload_bytes(
                    UPLOAD_BLOB_CONTAINER_ARTIFACTS,
                    artifact_blob_name,
                    artifact.get("artifact_bytes") or b"",
                    content_type="application/octet-stream",
                )
                artifact_blob_ref = str(uploaded_artifact.get("blob_ref", "") or "")
                artifact_format = str(artifact.get("format", "parquet") or "parquet")
                artifact_row_count = int(artifact.get("row_count", 0) or 0)
                job["artifact_blob_ref"] = artifact_blob_ref
                job["artifact_format"] = artifact_format
        except Exception as exc:
            logger.warning("[Upload] tabular artifact build failed for %s: %s", filename, exc)

    store_entry, result_payload = await _extract_upload_entry(
        filename, raw_bytes, content_type,
        prebuilt_artifact=artifact if artifact.get("artifact_bytes") else None,
    )

    raw_blob_retention_until = str(
        job.get("raw_blob_retention_until", "") or _raw_blob_retention_until_iso(
            filename=filename,
            artifact_blob_ref=artifact_blob_ref,
            fallback_hours=UPLOAD_ARTIFACT_RETENTION_HOURS,
        )
    )[:64]
    job["raw_blob_retention_until"] = raw_blob_retention_until

    extracted_blob_ref = ""
    text_payload = str(store_entry.get("data_text", "") or "")
    if text_payload:
        text_blob_name = str(job.get("text_blob_name", "") or "")
        if text_blob_name:
            text_blob = await blob_upload_bytes(
                UPLOAD_BLOB_CONTAINER_TEXT,
                text_blob_name,
                text_payload.encode("utf-8", errors="replace"),
                content_type="text/plain; charset=utf-8",
            )
            extracted_blob_ref = text_blob.get("blob_ref", "")

    chunks_blob_ref = ""
    chunks = store_entry.pop("chunks", None)
    col_names = store_entry.get("col_names", [])
    col_analysis = store_entry.get("col_analysis", [])
    polymorphic_schema = store_entry.get("polymorphic_schema") if isinstance(store_entry.get("polymorphic_schema"), dict) else None
    # NOTE: Semantic chunk embedding is deferred to background / first-query.
    # The artifact is already stored and can be chunked lazily via the backfill
    # endpoint or on first semantic search.  This avoids blocking upload with
    # hundreds of embedding API calls that can take 30-120 seconds.
    if (
        not chunks
        and artifact_blob_ref
        and is_tabular_filename(filename)
        and artifact.get("artifact_bytes")
    ):
        store_entry["needs_artifact_semantic_chunks"] = True
        logger.info("[Upload] deferring semantic chunk build for %s (artifact ready)", filename)
    if isinstance(chunks, list) and chunks:
        chunks_blob_name = str(job.get("chunks_blob_name", "") or "")
        if chunks_blob_name:
            chunks_blob = await blob_upload_json(
                UPLOAD_BLOB_CONTAINER_CHUNKS,
                chunks_blob_name,
                {"chunks": chunks},
            )
            chunks_blob_ref = chunks_blob.get("blob_ref", "")

    if artifact_blob_ref and is_tabular_filename(filename):
        raw_blob_retention_until = str(
            _raw_blob_retention_until_iso(
                filename=filename,
                artifact_blob_ref=artifact_blob_ref,
                has_chunks=bool(chunks_blob_ref),
                fallback_hours=UPLOAD_ARTIFACT_RETENTION_HOURS,
            )
        )[:64]
        job["raw_blob_retention_until"] = raw_blob_retention_until

    store_entry["extracted_blob_ref"] = extracted_blob_ref
    store_entry["chunks_blob_ref"] = chunks_blob_ref
    store_entry["tabular_artifact_blob_ref"] = artifact_blob_ref
    store_entry["tabular_artifact_format"] = artifact_format
    store_entry["has_chunks"] = bool(chunks_blob_ref)
    store_entry["user_sub"] = user_sub

    all_summaries = await _append_uploaded_entry_safe(conv_id, store_entry)
    if conv_id in conversation_meta:
        conversation_meta[conv_id]["file_injected"] = False

    upload_index_entity = {
        "PartitionKey": conv_id,
        "RowKey": str(job.get("job_id", "")),
        "UserSub": user_sub[:120],
        "Filename": filename[:240],
        "ContentType": content_type[:120],
        "FileSizeBytes": int(job.get("size_bytes", 0) or 0),
        "UploadedAt": datetime.now(timezone.utc).isoformat(),
        "RowCount": int(store_entry.get("row_count", 0) or 0),
        "ColNamesJson": json.dumps(col_names if isinstance(col_names, list) else [], ensure_ascii=False)[:32000],
        "ColAnalysisJson": json.dumps(col_analysis if isinstance(col_analysis, list) else [], ensure_ascii=False)[:32000],
        "PreviewText": text_payload[:16000],
        "Truncated": bool(store_entry.get("truncated", False)),
        "HasChunks": bool(chunks_blob_ref),
        "ExtractedBlobRef": extracted_blob_ref,
        "ChunksBlobRef": chunks_blob_ref,
        "RawBlobRef": raw_blob_ref,
        "TabularArtifactBlobRef": artifact_blob_ref,
        "TabularArtifactFormat": artifact_format[:32],
        "TabularArtifactRowCount": int(artifact_row_count or 0),
        "PolymorphicSummary": str((polymorphic_schema or {}).get("summary_text", "") or "")[:4000],
        "PivotColumn": str((polymorphic_schema or {}).get("pivot_column", "") or "")[:240],
        "PivotValuesCount": int((polymorphic_schema or {}).get("pivot_values_count", 0) or 0),
        "PolymorphicSchemaJson": json.dumps(polymorphic_schema, ensure_ascii=False)[:32000] if polymorphic_schema else "",
        "RawBlobRetentionUntil": raw_blob_retention_until,
        "RawBlobPurgedAt": str(job.get("raw_blob_purged_at", "") or "")[:64],
        "RetentionUntil": str(job.get("retention_until", "") or _retention_until_iso(hours=UPLOAD_ARTIFACT_RETENTION_HOURS))[:64],
    }
    await _upsert_upload_index(upload_index_entity)

    result_payload["conversation_id"] = conv_id
    result_payload["status"] = "ok"
    result_payload["total_files"] = len(all_summaries)
    result_payload["all_files"] = all_summaries
    result_payload["has_chunks"] = bool(chunks_blob_ref)
    result_payload["tabular_artifact_ready"] = bool(artifact_blob_ref)
    if artifact_format:
        result_payload["tabular_artifact_format"] = artifact_format
    result_payload["index_row_key"] = upload_index_entity["RowKey"]

    job["status"] = "completed"
    job["result"] = result_payload
    job["error"] = ""


async def _run_upload_job(job_id: str, *, inline_content: bytes | None = None) -> None:
    _cleanup_upload_jobs()
    job = await upload_jobs_store.get_or_fetch(job_id)
    if not job:
        loaded = await _load_upload_job_from_storage(job_id)
        if not loaded:
            return
        job = loaded
        await upload_jobs_store.put(job_id, job)

    now_iso = datetime.now(timezone.utc).isoformat()
    if str(job.get("status", "")).lower() in ("completed", "failed"):
        return

    job["status"] = "processing"
    job["started_at"] = job.get("started_at") or now_iso
    job["updated_at"] = now_iso
    await _persist_upload_job(job)

    async with _upload_jobs_semaphore:
        try:
            await _process_upload_job(job, inline_content=inline_content)
        except Exception as e:
            job["status"] = "failed"
            job["error"] = str(e)
            job["result"] = None
        finally:
            job["finished_at"] = datetime.now(timezone.utc).isoformat()
            job["updated_at"] = datetime.now(timezone.utc).isoformat()
            await upload_jobs_store.put(str(job.get("job_id", "")), job)
            await _persist_upload_job(job)


async def _queue_upload_job(
    conv_id: str,
    user_sub: str,
    filename: str,
    content: bytes,
    content_type: str = "",
) -> dict:
    job_id = uuid.uuid4().hex
    blob_paths = _build_upload_blob_paths(conv_id, job_id, filename)
    raw_blob = await blob_upload_bytes(
        UPLOAD_BLOB_CONTAINER_RAW,
        blob_paths["raw_blob_name"],
        content,
        content_type=content_type or "application/octet-stream",
    )
    now_iso = datetime.now(timezone.utc).isoformat()
    job = {
        "job_id": job_id,
        "status": "queued",
        "conversation_id": conv_id,
        "filename": filename,
        "size_bytes": len(content),
        "created_at": now_iso,
        "updated_at": now_iso,
        "started_at": "",
        "finished_at": "",
        "error": "",
        "result": None,
        "user_sub": user_sub,
        "content_type": (content_type or "")[:120],
        "raw_blob_ref": raw_blob.get("blob_ref", build_blob_ref(UPLOAD_BLOB_CONTAINER_RAW, blob_paths["raw_blob_name"])),
        "text_blob_name": blob_paths["text_blob_name"],
        "chunks_blob_name": blob_paths["chunks_blob_name"],
        "artifact_blob_name": blob_paths["artifact_blob_name"],
        "artifact_blob_ref": "",
        "artifact_format": "",
        "raw_blob_retention_until": _retention_until_iso(hours=UPLOAD_ARTIFACT_RETENTION_HOURS),
        "raw_blob_purged_at": "",
        "retention_until": _retention_until_iso(hours=UPLOAD_ARTIFACT_RETENTION_HOURS),
    }
    await upload_jobs_store.put(job_id, job)
    await _persist_upload_job(job)
    # Processamento é feito por worker (inline loop ou worker externo).
    return job


async def _queue_upload_job_from_blob(
    conv_id: str,
    user_sub: str,
    filename: str,
    raw_blob_ref: str,
    size_bytes: int,
    content_type: str = "",
    *,
    job_id: Optional[str] = None,
) -> dict:
    resolved_job_id = str(job_id or uuid.uuid4().hex)
    blob_paths = _build_upload_blob_paths(conv_id, resolved_job_id, filename)
    now_iso = datetime.now(timezone.utc).isoformat()
    job = {
        "job_id": resolved_job_id,
        "status": "queued",
        "conversation_id": conv_id,
        "filename": filename,
        "size_bytes": int(size_bytes or 0),
        "created_at": now_iso,
        "updated_at": now_iso,
        "started_at": "",
        "finished_at": "",
        "error": "",
        "result": None,
        "user_sub": user_sub,
        "content_type": (content_type or "")[:120],
        "raw_blob_ref": raw_blob_ref,
        "text_blob_name": blob_paths["text_blob_name"],
        "chunks_blob_name": blob_paths["chunks_blob_name"],
        "artifact_blob_name": blob_paths["artifact_blob_name"],
        "artifact_blob_ref": "",
        "artifact_format": "",
        "raw_blob_retention_until": _retention_until_iso(hours=UPLOAD_ARTIFACT_RETENTION_HOURS),
        "raw_blob_purged_at": "",
        "retention_until": _retention_until_iso(hours=UPLOAD_ARTIFACT_RETENTION_HOURS),
    }
    await upload_jobs_store.put(resolved_job_id, job)
    await _persist_upload_job(job)
    return job


async def _nudge_upload_worker(max_jobs: int = 1, delay_seconds: float = 0.05) -> None:
    """Best-effort local pickup so async uploads don't depend only on sidecars."""
    # Wake the background worker loop immediately instead of waiting for the poll
    _upload_worker_wake_event.set()
    if delay_seconds > 0:
        await asyncio.sleep(delay_seconds)
    try:
        await process_upload_jobs_once(max_jobs=max_jobs)
    except Exception as e:
        logger.warning("[App] upload worker nudge failed: %s", e)


async def process_upload_jobs_once(max_jobs: int = UPLOAD_WORKER_BATCH_SIZE) -> dict:
    _cleanup_upload_jobs()
    target = max(1, min(int(max_jobs or 1), 50))
    processed = 0
    claimed = 0
    skipped = 0
    recovered = 0

    # Fetch both queued and processing jobs — processing jobs that are stale
    # (orphaned by a container restart) will be reclaimed and re-processed.
    try:
        rows_queued = await table_query(
            "UploadJobs",
            "PartitionKey eq 'upload' and Status eq 'queued'",
            top=max(target * 3, target),
        )
    except Exception as e:
        logger.warning("[App] process_upload_jobs_once query failed: %s", e)
        rows_queued = []

    # Recover orphaned "processing" jobs (stale from dead container instances)
    try:
        rows_processing = await table_query(
            "UploadJobs",
            "PartitionKey eq 'upload' and Status eq 'processing'",
            top=max(target * 2, target),
        )
    except Exception:
        rows_processing = []

    now = datetime.now(timezone.utc)
    orphan_threshold = 120  # seconds — if processing for >2 min, assume orphaned
    for row in rows_processing:
        job = _job_from_storage_row(row)
        updated_str = str(job.get("updated_at", "") or job.get("started_at", "") or "")
        try:
            updated_dt = datetime.fromisoformat(updated_str)
        except Exception:
            continue
        if (now - updated_dt).total_seconds() > orphan_threshold:
            job_id = str(job.get("job_id", "") or "")
            if job_id:
                job["status"] = "queued"
                job["updated_at"] = now.isoformat()
                job["worker_id"] = ""
                job["claim_token"] = ""
                await upload_jobs_store.put(job_id, job)
                await _persist_upload_job(job)
                rows_queued.append(row)
                recovered += 1
                logger.info("[Upload] recovered orphaned job %s (was processing for %ds)", job_id[:8], int((now - updated_dt).total_seconds()))

    rows_sorted = sorted(rows_queued, key=lambda r: str(r.get("CreatedAt", "")))
    for row in rows_sorted:
        if processed >= target:
            break
        job = _job_from_storage_row(row)
        job_id = str(job.get("job_id", "") or "")
        if not job_id:
            skipped += 1
            continue
        if str(job.get("status", "")).lower() not in ("queued",):
            skipped += 1
            continue

        claim_token = uuid.uuid4().hex
        job["worker_id"] = WORKER_INSTANCE_ID
        job["claim_token"] = claim_token
        job["status"] = "processing"
        job["started_at"] = datetime.now(timezone.utc).isoformat()
        job["updated_at"] = datetime.now(timezone.utc).isoformat()
        await upload_jobs_store.put(job_id, job)
        await _persist_upload_job(job)
        claimed += 1

        latest = await _load_upload_job_from_storage(job_id)
        if latest:
            await upload_jobs_store.put(job_id, latest)
            if str(latest.get("claim_token", "")) != claim_token:
                skipped += 1
                continue

        await _run_upload_job(job_id)
        processed += 1

    return {"processed": processed, "claimed": claimed, "skipped": skipped, "recovered": recovered}


async def _upload_worker_loop() -> None:
    poll_seconds = max(5.0, float(UPLOAD_WORKER_POLL_SECONDS))
    while True:
        try:
            await process_upload_jobs_once(max_jobs=UPLOAD_WORKER_BATCH_SIZE)
        except asyncio.CancelledError:
            raise
        except Exception as e:
            logger.warning("[App] upload worker loop failed: %s", e)
        # Sleep until next poll OR until woken up by a nudge event.
        _upload_worker_wake_event.clear()
        try:
            await asyncio.wait_for(_upload_worker_wake_event.wait(), timeout=poll_seconds)
        except asyncio.TimeoutError:
            pass


@app.post("/upload")
@limiter.limit("30/minute", key_func=_user_or_ip_rate_key)
async def upload_file(request: Request, file: UploadFile = File(...), conversation_id: Optional[str] = Form(None), credentials: HTTPAuthorizationCredentials = Depends(security)):
    user = get_current_user(credentials)
    conv_id = conversation_id or str(uuid.uuid4())
    filename = file.filename or "unknown"
    max_bytes = _max_upload_bytes_for_file(filename)
    content = await _read_upload_with_limit(file, max_bytes)
    user_sub = str(user.get("sub", "") or "")
    reserved_slots = await _count_reserved_slots_for_conversation(conv_id, user_sub=user_sub)
    if reserved_slots >= MAX_FILES_PER_CONVERSATION:
        raise HTTPException(
            400,
            (
                f"Limite de {MAX_FILES_PER_CONVERSATION} ficheiros por conversa atingido "
                "(incluindo uploads pendentes). Aguarda conclusão ou remove anexos."
            ),
        )
    _cleanup_upload_jobs()
    pending_jobs = await _count_pending_jobs_for_user(user_sub)
    if pending_jobs >= UPLOAD_MAX_PENDING_JOBS_PER_USER:
        raise HTTPException(
            429,
            (
                f"Limite de {UPLOAD_MAX_PENDING_JOBS_PER_USER} uploads pendentes atingido. "
                "Aguarda conclusão dos jobs em curso e tenta novamente."
            ),
        )
    try:
        job = await _queue_upload_job(conv_id, user_sub, filename, content, file.content_type or "")
        job_id = str(job.get("job_id", ""))
        # --- Inline processing for sync route: process immediately instead of
        # waiting for the background worker loop. This ensures small files
        # get an instant response without depending on the async worker. ---
        try:
            await _run_upload_job(job_id, inline_content=content)
            processed_job = await upload_jobs_store.get_or_fetch(job_id)
            if not processed_job:
                processed_job = await _load_upload_job_from_storage(job_id)
            if processed_job and str(processed_job.get("status", "")).lower() == "completed":
                return {
                    "status": "completed",
                    "job_id": job_id,
                    "conversation_id": conv_id,
                    "filename": filename,
                    "size_bytes": len(content),
                    "result": processed_job.get("result"),
                }
        except Exception as inline_err:
            logger.warning("[App] sync upload inline processing failed for %s, falling back to queued: %s", job_id, inline_err)
        # Fallback: return queued status (frontend will poll via waitUploadJob)
        create_logged_task(_nudge_upload_worker(), name=f"upload-nudge-{job_id[:8]}")
        return {
            "status": "queued",
            "job_id": job_id,
            "conversation_id": conv_id,
            "filename": filename,
            "size_bytes": len(content),
        }
    except Exception as e:
        raise HTTPException(400, str(e))


@app.post("/upload/async")
@limiter.limit("30/minute", key_func=_user_or_ip_rate_key)
async def upload_file_async(
    request: Request,
    file: UploadFile = File(...),
    conversation_id: Optional[str] = Form(None),
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    user = get_current_user(credentials)
    conv_id = conversation_id or str(uuid.uuid4())
    filename = file.filename or "unknown"
    max_bytes = _max_upload_bytes_for_file(filename)
    content = await _read_upload_with_limit(file, max_bytes)

    user_sub = str(user.get("sub", "") or "")
    reserved_slots = await _count_reserved_slots_for_conversation(conv_id, user_sub=user_sub)
    if reserved_slots >= MAX_FILES_PER_CONVERSATION:
        raise HTTPException(
            400,
            (
                f"Limite de {MAX_FILES_PER_CONVERSATION} ficheiros por conversa atingido "
                "(incluindo uploads pendentes). Aguarda conclusão ou remove anexos."
            ),
        )

    _cleanup_upload_jobs()
    pending_jobs = await _count_pending_jobs_for_user(user_sub)
    if pending_jobs >= UPLOAD_MAX_PENDING_JOBS_PER_USER:
        raise HTTPException(
            429,
            (
                f"Limite de {UPLOAD_MAX_PENDING_JOBS_PER_USER} uploads pendentes atingido. "
                "Aguarda conclusão dos jobs em curso e tenta novamente."
            ),
        )
    job = await _queue_upload_job(conv_id, user_sub, filename, content, file.content_type or "")
    create_logged_task(_nudge_upload_worker(), name=f"upload-nudge-{str(job.get('job_id', ''))[:8]}")
    return {
        "status": "queued",
        "job_id": job.get("job_id"),
        "conversation_id": conv_id,
        "filename": filename,
        "size_bytes": len(content),
    }


@app.post("/upload/stream/async")
@limiter.limit("20/minute", key_func=_user_or_ip_rate_key)
async def upload_file_stream_async(
    request: Request,
    conversation_id: Optional[str] = None,
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    user = get_current_user(credentials)
    conv_id = conversation_id or str(uuid.uuid4())
    encoded_filename = str(request.headers.get("x-upload-filename", "") or "").strip()
    filename = unquote(encoded_filename) if encoded_filename else "upload.bin"
    content_type = str(request.headers.get("content-type", "") or "").strip()

    if not filename or filename == "upload.bin":
        raise HTTPException(400, "Nome de ficheiro em falta")
    if not is_tabular_filename(filename):
        raise HTTPException(400, "O upload em streaming está disponível apenas para ficheiros tabulares")

    max_bytes = _max_upload_bytes_for_file(filename)
    declared_size = int(request.headers.get("content-length", "0") or 0)
    if declared_size and declared_size > max_bytes:
        raise HTTPException(413, f"Ficheiro excede limite máximo de {max_bytes} bytes")

    user_sub = str(user.get("sub", "") or "")
    reserved_slots = await _count_reserved_slots_for_conversation(conv_id, user_sub=user_sub)
    if reserved_slots >= MAX_FILES_PER_CONVERSATION:
        raise HTTPException(
            400,
            (
                f"Limite de {MAX_FILES_PER_CONVERSATION} ficheiros por conversa atingido "
                "(incluindo uploads pendentes). Aguarda conclusão ou remove anexos."
            ),
        )

    _cleanup_upload_jobs()
    pending_jobs = await _count_pending_jobs_for_user(user_sub)
    if pending_jobs >= UPLOAD_MAX_PENDING_JOBS_PER_USER:
        raise HTTPException(
            429,
            (
                f"Limite de {UPLOAD_MAX_PENDING_JOBS_PER_USER} uploads pendentes atingido. "
                "Aguarda conclusão dos jobs em curso e tenta novamente."
            ),
        )

    job_id = uuid.uuid4().hex
    blob_paths = _build_upload_blob_paths(conv_id, job_id, filename)
    try:
        raw_blob = await blob_upload_stream(
            UPLOAD_BLOB_CONTAINER_RAW,
            blob_paths["raw_blob_name"],
            request.stream(),
            content_type=content_type or "application/octet-stream",
            max_bytes=max_bytes,
        )
    except ValueError as exc:
        raise HTTPException(413, str(exc))
    except Exception as exc:
        raise HTTPException(400, str(exc))

    size_bytes = int(raw_blob.get("size_bytes", declared_size or 0) or 0)
    job = await _queue_upload_job_from_blob(
        conv_id,
        user_sub,
        filename,
        raw_blob.get("blob_ref", build_blob_ref(UPLOAD_BLOB_CONTAINER_RAW, blob_paths["raw_blob_name"])),
        size_bytes,
        content_type,
        job_id=job_id,
    )
    create_logged_task(_nudge_upload_worker(), name=f"upload-nudge-{str(job.get('job_id', ''))[:8]}")
    return {
        "status": "queued",
        "job_id": job.get("job_id"),
        "conversation_id": conv_id,
        "filename": filename,
        "size_bytes": size_bytes,
        "upload_mode": "stream",
    }


@app.post("/upload/batch/async")
@limiter.limit("20/minute", key_func=_user_or_ip_rate_key)
async def upload_files_async_batch(
    request: Request,
    files: List[UploadFile] = File(...),
    conversation_id: Optional[str] = Form(None),
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    user = get_current_user(credentials)
    if not files:
        raise HTTPException(400, "Nenhum ficheiro recebido")

    conv_id = conversation_id or str(uuid.uuid4())
    user_sub = str(user.get("sub", "") or "")
    reserved_slots = await _count_reserved_slots_for_conversation(conv_id, user_sub=user_sub)
    available_slots = max(0, MAX_FILES_PER_CONVERSATION - reserved_slots)
    if available_slots <= 0:
        raise HTTPException(
            400,
            (
                f"Limite de {MAX_FILES_PER_CONVERSATION} ficheiros por conversa atingido "
                "(incluindo uploads pendentes). Aguarda conclusão ou remove anexos."
            ),
        )

    files_to_queue = list(files[:available_slots])
    skipped = []
    for f in files[available_slots:]:
        skipped.append(
            {
                "filename": f.filename or "unknown",
                "reason": f"Limite de {MAX_FILES_PER_CONVERSATION} ficheiros por conversa",
            }
        )

    _cleanup_upload_jobs()
    pending_jobs = await _count_pending_jobs_for_user(user_sub)
    pending_capacity = max(0, UPLOAD_MAX_PENDING_JOBS_PER_USER - pending_jobs)
    if pending_capacity <= 0:
        raise HTTPException(
            429,
            (
                f"Limite de {UPLOAD_MAX_PENDING_JOBS_PER_USER} uploads pendentes atingido. "
                "Aguarda conclusão dos jobs em curso e tenta novamente."
            ),
        )

    if len(files_to_queue) > pending_capacity:
        overflow = files_to_queue[pending_capacity:]
        files_to_queue = files_to_queue[:pending_capacity]
        for f in overflow:
            skipped.append(
                {
                    "filename": f.filename or "unknown",
                    "reason": f"Limite de {UPLOAD_MAX_PENDING_JOBS_PER_USER} uploads pendentes por utilizador",
                }
            )

    queued_jobs = []
    batch_total_bytes = 0
    for uf in files_to_queue:
        filename = uf.filename or "unknown"
        max_bytes = _max_upload_bytes_for_file(filename)
        try:
            content = await _read_upload_with_limit(uf, max_bytes)
        except HTTPException as exc:
            skipped.append(
                {
                    "filename": filename,
                    "reason": exc.detail,
                }
            )
            continue
        if batch_total_bytes + len(content) > UPLOAD_MAX_BATCH_TOTAL_BYTES:
            skipped.append(
                {
                    "filename": filename,
                    "reason": "Lote excede tamanho total máximo permitido",
                }
            )
            continue
        batch_total_bytes += len(content)
        job = await _queue_upload_job(conv_id, user_sub, filename, content, uf.content_type or "")
        create_logged_task(_nudge_upload_worker(), name=f"upload-nudge-{str(job.get('job_id', ''))[:8]}")
        queued_jobs.append(
            {
                "job_id": job.get("job_id"),
                "filename": filename,
                "size_bytes": len(content),
            }
        )

    if not queued_jobs:
        reason = skipped[0]["reason"] if skipped else "Nenhum ficheiro elegível para processamento"
        raise HTTPException(400, reason)

    return {
        "status": "queued",
        "conversation_id": conv_id,
        "queued_count": len(queued_jobs),
        "total_requested": len(files),
        "queued_jobs": queued_jobs,
        "skipped": skipped,
    }


@app.get("/api/upload/status/{job_id}")
@limiter.limit("30/minute", key_func=_user_or_ip_rate_key)
async def upload_job_status(
    request: Request,
    job_id: str,
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    user = get_current_user(credentials)
    _cleanup_upload_jobs()

    job = await upload_jobs_store.get_or_fetch(job_id)
    job = await _load_fresh_job_state(job_id, job)

    if not job:
        raise HTTPException(404, "Upload job não encontrado")

    owner_sub = str(job.get("user_sub", "") or "")
    if user.get("role") != "admin" and owner_sub and owner_sub != user.get("sub"):
        raise HTTPException(403, "Sem permissão para este upload job")
    if _is_job_stale(job):
        stale_msg = "Upload interrompido por timeout/stale (possível restart do servidor). Reenvia o ficheiro."
        job = await _mark_job_failed(job, stale_msg)
    return _job_public_view(job)


@app.post("/api/upload/status/batch")
@limiter.limit("60/minute", key_func=_user_or_ip_rate_key)
async def upload_jobs_status_batch(
    request: Request,
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    user = get_current_user(credentials)
    try:
        body = await request.json()
    except Exception:
        raise HTTPException(400, "Payload JSON inválido")
    raw_ids = body.get("job_ids", []) if isinstance(body, dict) else []
    if not isinstance(raw_ids, list) or not raw_ids:
        raise HTTPException(400, "job_ids é obrigatório")

    job_ids = [str(x).strip() for x in raw_ids if str(x).strip()]
    if not job_ids:
        raise HTTPException(400, "job_ids é obrigatório")
    if len(job_ids) > 100:
        raise HTTPException(400, "Máximo de 100 job_ids por pedido")

    _cleanup_upload_jobs()
    stale_msg = "Upload interrompido por timeout/stale (possível restart do servidor). Reenvia o ficheiro."
    items = []
    for job_id in job_ids:
        job = await upload_jobs_store.get_or_fetch(job_id)
        job = await _load_fresh_job_state(job_id, job)
        if not job:
            items.append({"job_id": job_id, "status": "not_found", "error": "Upload job não encontrado"})
            continue

        owner_sub = str(job.get("user_sub", "") or "")
        if user.get("role") != "admin" and owner_sub and owner_sub != user.get("sub"):
            items.append({"job_id": job_id, "status": "forbidden", "error": "Sem permissão para este upload job"})
            continue

        if _is_job_stale(job):
            job = await _mark_job_failed(job, stale_msg)

        items.append(_job_public_view(job))

    return {
        "total": len(items),
        "items": items,
    }


@app.get("/api/upload/jobs")
@limiter.limit("30/minute", key_func=_user_or_ip_rate_key)
async def list_upload_jobs(
    request: Request,
    status: Optional[str] = None,
    limit: int = 50,
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    user = get_current_user(credentials)
    if user.get("role") != "admin":
        raise HTTPException(403, "Apenas admins")

    _cleanup_upload_jobs()
    allowed_status = {"queued", "processing", "completed", "failed"}
    status_filter = (status or "").strip().lower()
    if status_filter and status_filter not in allowed_status:
        raise HTTPException(400, "status inválido")
    top = max(1, min(limit, 200))

    try:
        rows = await table_query("UploadJobs", "PartitionKey eq 'upload'", top=top)
        jobs = [_job_from_storage_row(r) for r in rows]
    except Exception as e:
        logger.warning("[App] list_upload_jobs table query failed: %s", e)
        jobs = []

    if not jobs:
        jobs = list(upload_jobs_store.values())

    stale_msg = "Upload interrompido por timeout/stale (possível restart do servidor). Reenvia o ficheiro."
    for j in jobs:
        if _is_job_stale(j):
            await _mark_job_failed(j, stale_msg)

    jobs_sorted = sorted(
        jobs,
        key=lambda j: str(j.get("updated_at") or j.get("created_at") or ""),
        reverse=True,
    )
    if status_filter:
        jobs_sorted = [j for j in jobs_sorted if str(j.get("status", "")).lower() == status_filter]
    jobs_sorted = jobs_sorted[:top]

    counts = {"queued": 0, "processing": 0, "completed": 0, "failed": 0}
    for j in jobs:
        s = str(j.get("status", "")).lower()
        if s in counts:
            counts[s] += 1

    items = []
    for j in jobs_sorted:
        public = _job_public_view(j)
        public["user_sub"] = j.get("user_sub", "")
        items.append(public)

    return {
        "total": len(items),
        "status_filter": status_filter or None,
        "counts": counts,
        "items": items,
    }


@app.post("/api/upload/worker/run-once")
@limiter.limit("30/minute", key_func=_user_or_ip_rate_key)
async def upload_worker_run_once(
    request: Request,
    max_jobs: int = UPLOAD_WORKER_BATCH_SIZE,
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    user = get_current_user(credentials)
    if user.get("role") != "admin":
        raise HTTPException(403, "Apenas admins")
    result = await process_upload_jobs_once(max_jobs=max_jobs)
    return {
        "status": "ok",
        "worker_id": WORKER_INSTANCE_ID,
        **result,
    }


@app.get("/api/upload/pending/{conversation_id}")
@limiter.limit("60/minute", key_func=_user_or_ip_rate_key)
async def upload_pending_for_conversation(
    request: Request,
    conversation_id: str,
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    user = get_current_user(credentials)
    _cleanup_upload_jobs()
    stale_msg = "Upload interrompido por timeout/stale (possível restart do servidor). Reenvia o ficheiro."
    # Marca stale com base em estado persistido para consistência cross-instância.
    try:
        safe_conv = odata_escape(conversation_id)
        rows = await table_query("UploadJobs", f"PartitionKey eq 'upload' and ConversationId eq '{safe_conv}'", top=500)
        for row in rows:
            job = _job_from_storage_row(row)
            if _is_job_stale(job):
                await _mark_job_failed(job, stale_msg)
    except Exception as e:
        logger.warning("[App] upload_pending_for_conversation stale scan failed: %s", e)
        for job in list(upload_jobs_store.values()):
            if str(job.get("conversation_id", "")) == conversation_id and _is_job_stale(job):
                await _mark_job_failed(job, stale_msg)
    include_all_users = user.get("role") == "admin"
    pending = await _count_pending_jobs_for_conversation(
        conv_id=conversation_id,
        user_sub=str(user.get("sub", "") or ""),
        include_all_users=include_all_users,
    )
    return {
        "conversation_id": conversation_id,
        "pending_jobs": pending.get("total", 0),
        "queued": pending.get("counts", {}).get("queued", 0),
        "processing": pending.get("counts", {}).get("processing", 0),
    }


@app.get("/api/upload/index/{conversation_id}")
@limiter.limit("60/minute", key_func=_user_or_ip_rate_key)
async def upload_index_for_conversation(
    request: Request,
    conversation_id: str,
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    user = get_current_user(credentials)
    include_all = user.get("role") == "admin"
    rows = await _list_upload_index(
        conversation_id,
        user_sub="" if include_all else str(user.get("sub", "") or ""),
        top=UPLOAD_INDEX_TOP,
    )
    rows_sorted = sorted(rows, key=lambda r: str(r.get("UploadedAt", "")), reverse=True)
    items = []
    for row in rows_sorted:
        items.append(
            {
                "file_id": row.get("RowKey", ""),
                "filename": row.get("Filename", ""),
                "uploaded_at": row.get("UploadedAt", ""),
                "row_count": int(row.get("RowCount", 0) or 0),
                "has_chunks": bool(row.get("HasChunks", False)),
                "extracted_blob_ref": row.get("ExtractedBlobRef", ""),
                "chunks_blob_ref": row.get("ChunksBlobRef", ""),
            }
        )
    return {"conversation_id": conversation_id, "total": len(items), "items": items}

# =============================================================================
# DAILY DIGEST (Task 6.3)
# =============================================================================

_DIGEST_FIELDS = [
    "System.Id",
    "System.Title",
    "System.State",
    "System.WorkItemType",
    "System.AssignedTo",
    "System.CreatedDate",
]


def _digest_format_item(raw_item: dict) -> dict:
    fields = raw_item.get("fields", {}) if isinstance(raw_item, dict) else {}
    assigned = fields.get("System.AssignedTo", "")
    if isinstance(assigned, dict):
        assigned_to = assigned.get("displayName", "")
    else:
        assigned_to = str(assigned or "")
    wi_id = raw_item.get("id")
    return {
        "id": wi_id,
        "title": fields.get("System.Title", ""),
        "state": fields.get("System.State", ""),
        "type": fields.get("System.WorkItemType", ""),
        "assigned_to": assigned_to,
        "created_date": fields.get("System.CreatedDate", ""),
        "url": f"https://dev.azure.com/{DEVOPS_ORG}/{DEVOPS_PROJECT}/_workitems/edit/{wi_id}" if wi_id else "",
    }


async def _run_digest_section(section_name: str, wiql_query: str) -> dict:
    headers = _devops_headers()
    section = {"count": 0, "items": []}
    batch_errors = []

    try:
        wiql_resp = await _devops_request_with_retry(
            "POST",
            _devops_url("wit/wiql?api-version=7.1"),
            headers,
            {"query": wiql_query},
            max_retries=3,
            timeout=60,
        )
        if "error" in wiql_resp:
            section["error"] = wiql_resp["error"]
            return section

        ids = [wi.get("id") for wi in wiql_resp.get("workItems", []) if wi.get("id")]
        section["count"] = len(ids)
        if not ids:
            return section

        details = []
        for i in range(0, len(ids), 100):
            batch = ids[i:i + 100]
            batch_resp = await _devops_request_with_retry(
                "POST",
                _devops_url("wit/workitemsbatch?api-version=7.1"),
                headers,
                {"ids": batch, "fields": _DIGEST_FIELDS},
                max_retries=3,
                timeout=60,
            )
            if "error" in batch_resp:
                batch_errors.append(batch_resp["error"])
                continue
            details.extend(batch_resp.get("value", []))

        section["items"] = [_digest_format_item(item) for item in details]
        if batch_errors:
            section["error"] = "; ".join(batch_errors[:3])
        return section
    except Exception as e:
        section["error"] = f"{section_name} failed: {str(e)}"
        return section


@app.get("/api/digest")
@limiter.shared_limit(
    "10/minute",
    scope="chat_budget",
    key_func=_user_or_ip_rate_key,
)
async def api_digest(request: Request, credentials: HTTPAuthorizationCredentials = Depends(security)):
    get_current_user(credentials)

    sections_wiql = {
        "created_yesterday": (
            "SELECT [System.Id] FROM WorkItems "
            f"WHERE [System.TeamProject] = '{DEVOPS_PROJECT}' "
            "AND [System.WorkItemType] = 'User Story' "
            "AND [System.CreatedDate] >= @Today-1 "
            "AND [System.CreatedDate] < @Today "
            "ORDER BY [System.CreatedDate] DESC"
        ),
        "old_bugs": (
            "SELECT [System.Id] FROM WorkItems "
            f"WHERE [System.TeamProject] = '{DEVOPS_PROJECT}' "
            "AND [System.WorkItemType] = 'Bug' "
            "AND [System.State] = 'Active' "
            "AND [System.CreatedDate] < @Today-7 "
            "ORDER BY [System.CreatedDate] ASC"
        ),
        "unassigned": (
            "SELECT [System.Id] FROM WorkItems "
            f"WHERE [System.TeamProject] = '{DEVOPS_PROJECT}' "
            "AND [System.State] <> 'Closed' "
            "AND [System.State] <> 'Removed' "
            "AND [System.AssignedTo] = '' "
            "ORDER BY [System.ChangedDate] DESC"
        ),
        "closed_this_week": (
            "SELECT [System.Id] FROM WorkItems "
            f"WHERE [System.TeamProject] = '{DEVOPS_PROJECT}' "
            "AND [System.State] = 'Closed' "
            "AND [Microsoft.VSTS.Common.ClosedDate] >= @StartOfWeek "
            "ORDER BY [Microsoft.VSTS.Common.ClosedDate] DESC"
        ),
    }

    payload = {
        "date": datetime.now(timezone.utc).strftime("%Y-%m-%d"),
        "project": DEVOPS_PROJECT,
    }
    section_items = list(sections_wiql.items())
    section_results = await asyncio.gather(
        *[_run_digest_section(section_name, wiql_query) for section_name, wiql_query in section_items],
        return_exceptions=True,
    )
    for (section_name, _), result in zip(section_items, section_results):
        if isinstance(result, Exception):
            payload[section_name] = {
                "count": 0,
                "items": [],
                "error": f"{section_name} failed: {str(result)}",
            }
        else:
            payload[section_name] = result
    return payload

# =============================================================================
# EXPORT ENDPOINTS (Fase 3)
# =============================================================================

@app.post("/api/export")
@limiter.limit("10/minute", key_func=_user_or_ip_rate_key)
async def export_data(request: Request, export_request: ExportRequest, credentials: HTTPAuthorizationCredentials = Depends(security)):
    """Exporta dados de tool results (sync para leve, async para pesado)."""
    user = get_current_user(credentials)
    data = await _resolve_export_payload(export_request, user)
    title = _safe_export_title(export_request.title or "Export DBDE")
    summary = export_request.summary or ""
    fmt = str(export_request.format or "xlsx").strip().lower()
    row_count = _export_rows_count(data)

    heavy = bool(
        fmt == "zip"
        or (row_count >= EXPORT_ASYNC_THRESHOLD_ROWS and fmt in ("csv", "xlsx", "pdf", "html"))
    )
    if EXPORT_AUTO_ASYNC_ENABLED and export_request.prefer_async and heavy:
        job = await _queue_export_job(
            user_sub=str(user.get("sub", "") or ""),
            conversation_id=str(export_request.conversation_id or ""),
            format_name=fmt,
            title=title,
            summary=summary,
            data=data,
        )
        return JSONResponse(
            status_code=202,
            content={
                "status": "queued",
                "job_id": job.get("job_id"),
                "format": fmt,
                "row_count": row_count,
                "status_endpoint": f"/api/export/status/{job.get('job_id')}",
            },
        )

    content, mime_type, filename = _build_export_file(fmt, data, title, summary)
    if not content:
        raise HTTPException(500, "Falha ao gerar export")
    return StreamingResponse(
        io.BytesIO(content),
        media_type=mime_type,
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@app.post("/api/export/async")
@limiter.limit("10/minute", key_func=_user_or_ip_rate_key)
async def export_data_async(request: Request, export_request: ExportRequest, credentials: HTTPAuthorizationCredentials = Depends(security)):
    user = get_current_user(credentials)
    data = await _resolve_export_payload(export_request, user)
    title = _safe_export_title(export_request.title or "Export DBDE")
    summary = export_request.summary or ""
    fmt = str(export_request.format or "xlsx").strip().lower()
    job = await _queue_export_job(
        user_sub=str(user.get("sub", "") or ""),
        conversation_id=str(export_request.conversation_id or ""),
        format_name=fmt,
        title=title,
        summary=summary,
        data=data,
    )
    return {
        "status": "queued",
        "job_id": job.get("job_id"),
        "format": fmt,
        "row_count": _export_rows_count(data),
        "status_endpoint": f"/api/export/status/{job.get('job_id')}",
    }


@app.post("/api/export-chat")
@limiter.limit("10/minute", key_func=_user_or_ip_rate_key)
async def export_chat(request: Request, credentials: HTTPAuthorizationCredentials = Depends(security)):
    """Exporta histórico de chat em HTML/PDF."""
    user = get_current_user(credentials)
    body = await request.json()
    messages = body.get("messages", [])
    format_type = str(body.get("format", "html") or "html").strip().lower()
    title = str(body.get("title", "Chat Export") or "Chat Export")

    if not isinstance(messages, list) or not messages:
        return JSONResponse({"error": "Sem mensagens para exportar"}, status_code=400)

    html_content = _render_chat_html(messages, title)
    ts = datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")
    safe_title = _safe_export_title(title or "Chat Export")

    format_requested = format_type
    fallback_reason = None

    if format_type == "pdf":
        try:
            from export_engine import to_chat_pdf

            pdf_bytes = to_chat_pdf(messages, title)
            filename = f"{safe_title}_{ts}.pdf"
            download_id = await _store_generated_file(
                pdf_bytes,
                "application/pdf",
                filename,
                "pdf",
                user_sub=str(user.get("sub", "") or ""),
                scope="export_chat",
            )
            if not download_id:
                raise HTTPException(500, "Falha ao armazenar PDF exportado")
            return {
                "url": f"/api/download/{download_id}",
                "filename": filename,
                "format_requested": "pdf",
                "format_served": "pdf",
                "fallback_reason": None,
            }
        except Exception as e:
            logging.warning("[ExportChat] PDF generation failed, falling back to HTML: %s", e)
            format_type = "html"
            fallback_reason = f"PDF generation failed: {str(e)[:200]}"

    filename = f"{safe_title}_{ts}.html"
    download_id = await _store_generated_file(
        html_content.encode("utf-8"),
        "text/html; charset=utf-8",
        filename,
        "html",
        user_sub=str(user.get("sub", "") or ""),
        scope="export_chat",
    )
    if not download_id:
        raise HTTPException(500, "Falha ao armazenar HTML exportado")
    return {
        "url": f"/api/download/{download_id}",
        "filename": filename,
        "format_requested": format_requested,
        "format_served": "html",
        "fallback_reason": fallback_reason,
    }


@app.get("/api/export/status/{job_id}")
@limiter.limit("60/minute", key_func=_user_or_ip_rate_key)
async def export_job_status(request: Request, job_id: str, credentials: HTTPAuthorizationCredentials = Depends(security)):
    user = get_current_user(credentials)
    user_sub = str(user.get("sub", "") or "")
    is_admin = user.get("role") == "admin"

    _cleanup_export_jobs()
    job = await export_jobs_store.get_or_fetch(job_id)
    job = await _load_fresh_export_job_state(job_id, job)
    if not job:
        raise HTTPException(404, "Export job não encontrado")
    if not is_admin and str(job.get("user_sub", "")) != user_sub:
        raise HTTPException(403, "Sem permissão para este export job")
    return _export_job_public_view(job)


@app.post("/api/export/worker/run-once")
@limiter.limit("20/minute", key_func=_user_or_ip_rate_key)
async def export_worker_run_once(request: Request, credentials: HTTPAuthorizationCredentials = Depends(security)):
    user = get_current_user(credentials)
    if user.get("role") != "admin":
        raise HTTPException(403, "Apenas admins")
    result = await process_export_jobs_once(max_jobs=EXPORT_WORKER_BATCH_SIZE)
    return {"status": "ok", "worker_id": EXPORT_WORKER_INSTANCE_ID, **result}


@app.get("/api/download/{download_id}")
@limiter.limit("30/minute", key_func=_user_or_ip_rate_key)
async def download_generated_file(request: Request, download_id: str, credentials: HTTPAuthorizationCredentials = Depends(security)):
    """Download de ficheiro gerado por tool (armazenamento temporário em memória)."""
    user = get_current_user(credentials)
    entry = await get_generated_file(download_id)
    if not entry:
        raise HTTPException(404, "Ficheiro não encontrado ou expirado")
    owner_sub = str(entry.get("user_sub", "") or "").strip()
    conversation_id = str(entry.get("conversation_id", "") or "").strip()
    if not _is_admin_user(user):
        current_user_sub = str(user.get("sub", "") or "")
        if owner_sub and owner_sub != current_user_sub:
            raise HTTPException(403, "Sem permissão para este ficheiro")
        if not owner_sub and conversation_id:
            allowed = await _conversation_belongs_to_user(conversation_id, current_user_sub)
            if not allowed:
                raise HTTPException(403, "Sem permissão para este ficheiro")

    filename = entry.get("filename", "download.bin")
    safe_filename = "".join(c if c.isalnum() or c in " _-." else "_" for c in filename)[:80] or "download.bin"
    content = entry.get("content", b"")
    if not isinstance(content, (bytes, bytearray)) or len(content) == 0:
        raise HTTPException(410, "Ficheiro expirado ou inválido")

    return StreamingResponse(
        io.BytesIO(content),
        media_type=entry.get("mime_type", "application/octet-stream"),
        headers={"Content-Disposition": f'attachment; filename="{safe_filename}"'},
    )

# =============================================================================
# AUTH ENDPOINTS
# =============================================================================

@app.post("/api/auth/login")
@limiter.limit("5/minute", key_func=_login_rate_key)
async def login(request: Request, login_request: LoginRequest):
    if is_account_locked(login_request.username) or await is_account_locked_persistent(login_request.username):
        raise HTTPException(
            429,
            f"Conta temporariamente bloqueada. Tenta novamente em {_LOCKOUT_DURATION_MINUTES} minutos.",
        )

    safe_username = odata_escape(login_request.username)
    users = await table_query("Users", f"PartitionKey eq 'user' and RowKey eq '{safe_username}'", top=1)
    if not users:
        record_login_failure(login_request.username)
        await record_login_failure_persistent(login_request.username)
        raise HTTPException(401, "Credenciais inválidas")
    user = users[0]
    if not verify_password(login_request.password, user.get("PasswordHash","")):
        record_login_failure(login_request.username)
        await record_login_failure_persistent(login_request.username)
        raise HTTPException(401, "Credenciais inválidas")
    if user.get("IsActive") == False:
        raise HTTPException(403, "Conta desactivada")
    clear_login_attempts(login_request.username)
    await clear_login_failures_persistent(login_request.username)
    token = jwt_encode({"sub":login_request.username, "role":user.get("Role","user"), "name":user.get("DisplayName",login_request.username)})
    response = JSONResponse(
        content={
            "status": "ok",
            "username": login_request.username,
            "role": user.get("Role", "user"),
            "display_name": user.get("DisplayName", login_request.username),
        }
    )
    secure_cookie = AUTH_COOKIE_SECURE if _request_is_https(request) else False
    response.set_cookie(
        key=AUTH_COOKIE_NAME,
        value=token,
        httponly=True,
        secure=secure_cookie,
        samesite="lax",
        path="/",
        max_age=AUTH_COOKIE_MAX_AGE_SECONDS,
    )
    return response


@app.post("/api/auth/logout")
@limiter.limit("20/minute", key_func=_user_or_ip_rate_key)
async def logout(request: Request):
    token = _extract_request_token(request)
    if token:
        try:
            payload = jwt_decode(token)
            jti = str(payload.get("jti", "") or "")
            exp_raw = payload.get("exp")
            if jti and isinstance(exp_raw, str):
                exp = datetime.fromisoformat(exp_raw)
                if exp.tzinfo is None:
                    exp = exp.replace(tzinfo=timezone.utc)
                await revoke_token_persistent(jti, exp, username=str(payload.get("sub", "") or ""))
        except ValueError:
            pass

    response = JSONResponse(content={"status": "ok"})
    secure_cookie = AUTH_COOKIE_SECURE if _request_is_https(request) else False
    response.set_cookie(
        key=AUTH_COOKIE_NAME,
        value="",
        httponly=True,
        secure=secure_cookie,
        samesite="lax",
        path="/",
        max_age=0,
    )
    return response

@app.post("/api/auth/create-user")
@limiter.limit("10/minute", key_func=_user_or_ip_rate_key)
async def create_user(request: Request, payload: CreateUserRequest, credentials: HTTPAuthorizationCredentials = Depends(security)):
    user = get_current_user(credentials)
    if user.get("role") != "admin": raise HTTPException(403, "Apenas admins")
    safe_username = odata_escape(payload.username)
    existing = await table_query("Users", f"PartitionKey eq 'user' and RowKey eq '{safe_username}'", top=1)
    if existing: raise HTTPException(409, "Username já existe")
    entity = {"PartitionKey":"user","RowKey":payload.username,"PasswordHash":hash_password(payload.password),"DisplayName":payload.display_name or payload.username,"Role":payload.role or "user","IsActive":True,"CreatedAt":datetime.now(timezone.utc).isoformat(),"CreatedBy":user.get("sub")}
    created = await table_insert("Users", entity)
    if not created:
        raise HTTPException(500, "Falha ao criar utilizador")
    return {"status":"ok","username":payload.username}

@app.get("/api/auth/users")
@limiter.limit("30/minute", key_func=_user_or_ip_rate_key)
async def list_users(request: Request, credentials: HTTPAuthorizationCredentials = Depends(security)):
    user = get_current_user(credentials)
    if user.get("role") != "admin": raise HTTPException(403, "Apenas admins")
    users = await table_query("Users", "PartitionKey eq 'user'", top=100)
    return {"users":[{"username":u.get("RowKey"),"display_name":u.get("DisplayName"),"role":u.get("Role"),"is_active":u.get("IsActive",True)} for u in users]}

@app.delete("/api/auth/users/{username}")
@limiter.limit("20/minute", key_func=_user_or_ip_rate_key)
async def deactivate_user(request: Request, username: str, credentials: HTTPAuthorizationCredentials = Depends(security)):
    user = get_current_user(credentials)
    if user.get("role") != "admin": raise HTTPException(403, "Apenas admins")
    if username == user.get("sub"): raise HTTPException(400, "Não podes desactivar-te")
    await table_merge("Users", {"PartitionKey":"user","RowKey":username,"IsActive":False})
    await persist_user_invalidation(username)
    return {"status":"ok"}

@app.post("/api/auth/change-password")
@limiter.limit("20/minute", key_func=_user_or_ip_rate_key)
async def change_password(request: Request, payload: ChangePasswordRequest, credentials: HTTPAuthorizationCredentials = Depends(security)):
    user = get_current_user(credentials)
    username = user.get("sub")
    safe_username = odata_escape(username)
    users = await table_query("Users", f"PartitionKey eq 'user' and RowKey eq '{safe_username}'", top=1)
    if not users: raise HTTPException(404, "User não encontrado")
    if not verify_password(payload.current_password, users[0].get("PasswordHash","")): raise HTTPException(401, "Password actual incorrecta")
    await table_merge("Users", {"PartitionKey":"user","RowKey":username,"PasswordHash":hash_password(payload.new_password)})
    await persist_user_invalidation(username)
    return {"status":"ok"}

@app.post("/api/auth/reset-password/{username}")
@limiter.limit("15/minute", key_func=_user_or_ip_rate_key)
async def admin_reset_password(request: Request, username: str, payload: LoginRequest, credentials: HTTPAuthorizationCredentials = Depends(security)):
    user = get_current_user(credentials)
    if user.get("role") != "admin": raise HTTPException(403, "Apenas admins")
    await table_merge("Users", {"PartitionKey":"user","RowKey":username,"PasswordHash":hash_password(payload.password)})
    await persist_user_invalidation(username)
    return {"status":"ok"}

@app.post("/api/auth/force-logout/{username}")
@limiter.limit("10/minute", key_func=_user_or_ip_rate_key)
async def force_logout_user(request: Request, username: str, credentials: HTTPAuthorizationCredentials = Depends(security)):
    user = get_current_user(credentials)
    if user.get("role") != "admin":
        raise HTTPException(403, "Apenas admins")
    await persist_user_invalidation(username)
    return {"status": "ok", "message": f"Tokens de {username} invalidados"}

@app.get("/api/auth/me")
@limiter.limit("60/minute", key_func=_user_or_ip_rate_key)
async def get_me(request: Request, credentials: HTTPAuthorizationCredentials = Depends(security)):
    user = get_current_user(credentials)
    return {"username":user.get("sub"),"role":user.get("role"),"name":user.get("name")}


@app.post("/api/speech/prompt", response_model=SpeechPromptNormalizeResponse)
@limiter.limit("30/minute", key_func=_user_or_ip_rate_key)
async def normalize_speech_prompt(
    request: Request,
    payload: SpeechPromptNormalizeRequest,
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    user = get_current_user(credentials)
    transcript = str(payload.transcript or "").strip()
    if not transcript:
        raise HTTPException(400, "Transcrição vazia")

    started = time.perf_counter()
    result = await normalize_spoken_prompt(
        transcript=transcript,
        mode=str(payload.mode or "general"),
        language=str(payload.language or "pt-PT"),
    )
    logger.info(
        "[Speech] normalized prompt provider=%s confidence=%s auto_send=%s elapsed_ms=%d",
        result.get("provider_used", "unknown"),
        result.get("confidence", "unknown"),
        result.get("auto_send_allowed", False),
        int((time.perf_counter() - started) * 1000),
    )
    try:
        await log_audit(
            user.get("sub"),
            "speech_prompt",
            transcript,
            tools_used=["speech_prompt"],
            duration_ms=int((time.perf_counter() - started) * 1000),
            metadata={
                "mode": str(payload.mode or "general"),
                "provider_used": result.get("provider_used", "unknown"),
                "model_used": result.get("model_used", ""),
                "confidence": result.get("confidence", "unknown"),
                "auto_send_allowed": bool(result.get("auto_send_allowed", False)),
                "conversation_id": payload.conversation_id or "",
                "provider_family": result.get("provider_family", ""),
                "external_provider": bool(result.get("external_provider", False)),
                "data_sensitivity": result.get("data_sensitivity", ""),
                "provider_policy_mode": result.get("provider_policy_mode", ""),
                "provider_policy_note": result.get("provider_policy_note", ""),
            },
        )
    except Exception as exc:
        logger.warning("[App] speech prompt audit failed: %s", exc)
    return SpeechPromptNormalizeResponse(**result)


@app.post("/api/speech/token", response_model=SpeechPromptTokenResponse)
@limiter.limit("30/minute", key_func=_user_or_ip_rate_key)
async def issue_speech_token(
    request: Request,
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    get_current_user(credentials)
    if not (AZURE_SPEECH_ENABLED and AZURE_SPEECH_KEY and AZURE_SPEECH_REGION):
        raise HTTPException(503, "Azure Speech não está configurado.")

    token_url = f"https://{AZURE_SPEECH_REGION}.api.cognitive.microsoft.com/sts/v1.0/issuetoken"
    try:
        async with httpx.AsyncClient(timeout=10.0) as client:
            response = await client.post(
                token_url,
                headers={"Ocp-Apim-Subscription-Key": AZURE_SPEECH_KEY, "Content-Length": "0"},
            )
            response.raise_for_status()
    except Exception as exc:
        logger.error("[Speech] token issuance failed: %s", exc)
        raise HTTPException(502, "Falha ao obter token do Azure Speech")

    token = str(response.text or "").strip()
    if not token:
        raise HTTPException(502, "Azure Speech não devolveu token válido")

    return SpeechPromptTokenResponse(
        token=token,
        region=AZURE_SPEECH_REGION,
        language=AZURE_SPEECH_LANGUAGE,
        expires_in_seconds=600,
    )


@app.post("/api/speech/synthesize")
@limiter.limit("20/minute", key_func=_user_or_ip_rate_key)
async def synthesize_speech(
    request: Request,
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    """Azure TTS — converte texto curto em áudio MP3 (para clarificações em modo voz)."""
    get_current_user(credentials)
    if not (AZURE_SPEECH_ENABLED and AZURE_SPEECH_KEY and AZURE_SPEECH_REGION):
        raise HTTPException(503, "Azure Speech TTS não está configurado.")

    try:
        body = await request.json()
    except Exception:
        raise HTTPException(400, "Body JSON inválido.")

    text = str(body.get("text", "") or "").strip()
    if not text or len(text) > 1000:
        raise HTTPException(400, "Texto em falta ou demasiado longo (max 1000 chars).")

    voice = str(body.get("voice", "") or "").strip() or "pt-PT-RaquelNeural"
    language = str(body.get("language", "") or "").strip() or "pt-PT"

    # Validate voice/language against injection (only allow Azure Neural voice format)
    import re as _re
    if not _re.fullmatch(r'[a-zA-Z]{2,5}-[A-Z]{2,5}-[A-Za-z]+Neural', voice):
        voice = "pt-PT-RaquelNeural"
    if not _re.fullmatch(r'[a-zA-Z]{2,5}-[A-Z]{2,5}', language):
        language = "pt-PT"

    # Sanitize text for SSML
    safe_text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    ssml = (
        f'<speak version="1.0" xmlns="http://www.w3.org/2001/10/synthesis" xml:lang="{language}">'
        f'<voice name="{voice}">{safe_text}</voice>'
        f'</speak>'
    )
    tts_url = f"https://{AZURE_SPEECH_REGION}.tts.speech.microsoft.com/cognitiveservices/v1"
    try:
        async with httpx.AsyncClient(timeout=15.0) as client:
            resp = await client.post(
                tts_url,
                headers={
                    "Ocp-Apim-Subscription-Key": AZURE_SPEECH_KEY,
                    "Content-Type": "application/ssml+xml",
                    "X-Microsoft-OutputFormat": "audio-16khz-128kbitrate-mono-mp3",
                },
                content=ssml.encode("utf-8"),
            )
            resp.raise_for_status()
    except Exception as exc:
        logger.warning("[App] TTS synthesis failed: %s", exc)
        raise HTTPException(502, "Falha na síntese de voz.")

    return Response(content=resp.content, media_type="audio/mpeg")


# =============================================================================
# FEEDBACK
# =============================================================================

@app.post("/feedback")
@limiter.limit("30/minute", key_func=_user_or_ip_rate_key)
async def submit_feedback(request: Request, payload: FeedbackRequest, credentials: HTTPAuthorizationCredentials = Depends(security)):
    user = get_current_user(credentials)
    question, answer = "", ""
    cid = payload.conversation_id
    if not await _conversation_belongs_to_user(cid, str(user.get("sub", "") or "")):
        raise HTTPException(403, "Sem permissão para esta conversa")
    if cid in conversations:
        um = [m for m in conversations[cid] if m.get("role")=="user"]
        am = [m for m in conversations[cid] if m.get("role")=="assistant"]
        if um: question = um[-1].get("content","") if isinstance(um[-1].get("content"),str) else str(um[-1].get("content",""))
        if am: answer = am[-1].get("content","")
    
    ts = datetime.now(timezone.utc).isoformat()
    safe_conv = cid.replace("-","")[:32]
    entity = {"PartitionKey":safe_conv,"RowKey":f"{payload.message_index}_{ts.replace(':','').replace('-','').replace('.','')}", "Rating":payload.rating,"Note":payload.note or "","Question":question[:2000],"Answer":answer[:4000],"Timestamp_str":ts,"UserSub":str(user.get("sub", "") or "")}
    stored = await table_insert("feedback", entity)
    if not stored: feedback_memory.append(entity)

    if question and answer and (payload.rating >= 7 or payload.rating <= 3):
        etype = "positive" if payload.rating >= 7 else "negative"
        eid = f"{safe_conv}_{payload.message_index}"
        await table_insert("examples", {"PartitionKey":etype,"RowKey":eid,"Question":question[:2000],"Answer":answer[:4000],"Rating":payload.rating,"Note":payload.note or "","Timestamp_str":ts,"UserSub":str(user.get("sub", "") or "")})
        try: await _index_example(eid, question, answer, payload.rating, example_type=etype, feedback_note=payload.note or "")
        except Exception as e:
            logger.error("[App] _index_example in feedback failed: %s", e)
    
    return {"status":"ok","message":f"Feedback: {payload.rating}/10","persisted":"table_storage" if stored else "memory"}

@app.get("/feedback/stats")
@limiter.limit("30/minute", key_func=_user_or_ip_rate_key)
async def feedback_stats(request: Request, credentials: HTTPAuthorizationCredentials = Depends(security)):
    user = get_current_user(credentials)
    if not _is_admin_user(user):
        raise HTTPException(403, "Apenas admins")
    fbs = await table_query("feedback", top=1000)
    all_fb = fbs + list(feedback_memory)
    if not all_fb: return {"total":0,"average_rating":0}
    ratings = [f.get("Rating",0) for f in all_fb if f.get("Rating",0)>0]
    if not ratings: return {"total":0,"average_rating":0}
    return {"total":len(ratings),"average_rating":round(sum(ratings)/len(ratings),1),"distribution":{str(r):ratings.count(r) for r in range(1,11)}}

# =============================================================================
# CHAT PERSISTENCE
# =============================================================================

@app.post("/api/chats/save")
@limiter.limit("30/minute", key_func=_user_or_ip_rate_key)
async def save_chat(request: Request, payload: SaveChatRequest, credentials: HTTPAuthorizationCredentials = Depends(security)):
    user = get_current_user(credentials)
    uid = user.get("sub", payload.user_id)
    msgs = [{"role":m.get("role",""),"content":m.get("content","")} for m in payload.messages]
    msgs_json = json.dumps(msgs, ensure_ascii=False)
    while len(msgs_json)>60000 and len(msgs)>4: msgs.pop(1); msgs_json = json.dumps(msgs, ensure_ascii=False)
    entity = {"PartitionKey":uid,"RowKey":payload.conversation_id,"Title":(payload.title or "Nova conversa")[:100],"Messages":msgs_json,"MessageCount":len(payload.messages),"UpdatedAt":datetime.now(timezone.utc).isoformat()}
    if not await table_insert("ChatHistory", entity): await table_merge("ChatHistory", entity)
    return {"status":"ok","conversation_id":payload.conversation_id}

@app.get("/api/chats/{user_id}")
@limiter.limit("60/minute", key_func=_user_or_ip_rate_key)
async def list_chats(request: Request, user_id: str, credentials: HTTPAuthorizationCredentials = Depends(security)):
    user = get_current_user(credentials)
    uid = user.get("sub") if user.get("role")!="admin" else user_id
    safe_uid = odata_escape(uid)
    entities = await table_query("ChatHistory", f"PartitionKey eq '{safe_uid}'", top=100)
    chats = sorted([{"conversation_id":e.get("RowKey",""),"title":e.get("Title",""),"message_count":e.get("MessageCount",0),"updated_at":e.get("UpdatedAt","")} for e in entities], key=lambda c:c["updated_at"], reverse=True)
    return {"chats":chats}

@app.get("/api/chats/{user_id}/{conversation_id}")
@limiter.limit("60/minute", key_func=_user_or_ip_rate_key)
async def get_chat(request: Request, user_id: str, conversation_id: str, credentials: HTTPAuthorizationCredentials = Depends(security)):
    user = get_current_user(credentials)
    uid = user.get("sub") if user.get("role")!="admin" else user_id
    safe_uid = odata_escape(uid)
    safe_conv = odata_escape(conversation_id)
    es = await table_query("ChatHistory", f"PartitionKey eq '{safe_uid}' and RowKey eq '{safe_conv}'", top=1)
    if not es: raise HTTPException(404, "Não encontrada")
    return {"conversation_id":conversation_id,"title":es[0].get("Title",""),"messages":json.loads(es[0].get("Messages","[]")),"updated_at":es[0].get("UpdatedAt","")}


@app.post("/api/chats/{user_id}/{conversation_id}/title")
@limiter.limit("30/minute", key_func=_user_or_ip_rate_key)
async def update_chat_title(request: Request, user_id: str, conversation_id: str, payload: UpdateChatTitleRequest, credentials: HTTPAuthorizationCredentials = Depends(security)):
    user = get_current_user(credentials)
    uid = user.get("sub") if user.get("role") != "admin" else user_id
    safe_uid = odata_escape(uid)
    safe_conv = odata_escape(conversation_id)
    existing = await table_query("ChatHistory", f"PartitionKey eq '{safe_uid}' and RowKey eq '{safe_conv}'", top=1)
    if not existing:
        raise HTTPException(404, "Conversa não encontrada")
    title = str(payload.title or "").strip()
    if not title:
        raise HTTPException(400, "Título inválido")
    entity = {
        "PartitionKey": uid,
        "RowKey": conversation_id,
        "Title": title[:100],
        "UpdatedAt": datetime.now(timezone.utc).isoformat(),
    }
    await table_merge("ChatHistory", entity)
    return {"status": "ok", "conversation_id": conversation_id, "title": entity["Title"], "updated_at": entity["UpdatedAt"]}


async def _purge_upload_index_row(row: dict, *, delete_job: bool = True) -> dict:
    row_pk = str(row.get("PartitionKey", "") or "").strip()
    row_key = str(row.get("RowKey", "") or "").strip()
    blobs_deleted = 0
    blob_fields = ("RawBlobRef", "ExtractedBlobRef", "ChunksBlobRef", "TabularArtifactBlobRef")
    for field in blob_fields:
        ref = str(row.get(field, "") or "").strip()
        if not ref:
            continue
        container, blob_name = parse_blob_ref(ref)
        if not container or not blob_name:
            continue
        try:
            await blob_delete(container, blob_name)
            blobs_deleted += 1
        except Exception as e:
            logger.warning("[App] upload artifact blob delete failed for %s: %s", ref, e)
    row_deleted = 0
    if row_pk and row_key:
        try:
            await table_delete("UploadIndex", row_pk, row_key)
            row_deleted = 1
        except Exception as e:
            logger.warning("[App] upload index delete failed for %s/%s: %s", row_pk, row_key, e)
    if row_pk:
        uploaded_files_store.pop(row_pk, None)
        if row_pk in conversation_meta:
            conversation_meta[row_pk]["file_injected"] = False
    if delete_job and row_key:
        try:
            await upload_jobs_store.delete(row_key, partition_key="upload")
            job_deleted = 1
        except Exception as e:
            logger.warning("[App] upload job delete failed for %s: %s", row_key, e)
            job_deleted = 0
    else:
        job_deleted = 0
    return {"rows_deleted": row_deleted, "blobs_deleted": blobs_deleted, "jobs_deleted": job_deleted}


async def _clear_upload_index_raw_blob(row: dict, *, purged_at: str) -> int:
    row_pk = str(row.get("PartitionKey", "") or "").strip()
    row_key = str(row.get("RowKey", "") or "").strip()
    if not row_pk or not row_key:
        return 0
    try:
        await table_merge(
            "UploadIndex",
            {
                "PartitionKey": row_pk,
                "RowKey": row_key,
                "RawBlobRef": "",
                "RawBlobRetentionUntil": "",
                "RawBlobPurgedAt": purged_at[:64],
            },
        )
        return 1
    except Exception as e:
        logger.warning("[App] upload index raw blob clear failed for %s/%s: %s", row_pk, row_key, e)
        return 0


async def _clear_upload_job_raw_blob(job_id: str, *, purged_at: str) -> int:
    safe_job_id = str(job_id or "").strip()
    if not safe_job_id:
        return 0
    try:
        await table_merge(
            "UploadJobs",
            {
                "PartitionKey": "upload",
                "RowKey": safe_job_id,
                "RawBlobRef": "",
                "RawBlobRetentionUntil": "",
                "RawBlobPurgedAt": purged_at[:64],
            },
        )
        cached = upload_jobs_store.get(safe_job_id)
        if cached:
            cached["raw_blob_ref"] = ""
            cached["raw_blob_retention_until"] = ""
            cached["raw_blob_purged_at"] = purged_at[:64]
            await upload_jobs_store.put(safe_job_id, cached)
        return 1
    except Exception as e:
        logger.warning("[App] upload job raw blob clear failed for %s: %s", safe_job_id, e)
        return 0


async def _purge_expired_upload_artifacts(limit: int = 120) -> dict:
    now = datetime.now(timezone.utc)
    safe_limit = max(1, min(int(limit or 120), 500))
    rows_deleted = 0
    blobs_deleted = 0
    jobs_deleted = 0
    deleted_job_ids: set[str] = set()
    try:
        upload_index_rows = await table_query("UploadIndex", top=safe_limit)
    except Exception as e:
        logger.warning("[App] expired upload artifact sweep query failed: %s", e)
        upload_index_rows = []

    active_index_ids = set()
    for row in upload_index_rows:
        row_key = str(row.get("RowKey", "") or "").strip()
        if row_key:
            active_index_ids.add(row_key)
        retention_until = str(row.get("RetentionUntil", "") or "").strip()
        if not _is_retention_expired(retention_until, now=now):
            raw_blob_ref = str(row.get("RawBlobRef", "") or "").strip()
            raw_retention_until = _effective_raw_blob_retention_until(row)
            has_artifact = bool(str(row.get("TabularArtifactBlobRef", "") or "").strip())
            if raw_blob_ref and has_artifact and _is_retention_expired(raw_retention_until, now=now):
                purged_at = now.isoformat()
                container, blob_name = parse_blob_ref(raw_blob_ref)
                if container and blob_name:
                    try:
                        await blob_delete(container, blob_name)
                        blobs_deleted += 1
                    except Exception as e:
                        logger.warning("[App] upload raw blob early delete failed for %s: %s", raw_blob_ref, e)
                        continue
                await _clear_upload_index_raw_blob(row, purged_at=purged_at)
                if row_key:
                    await _clear_upload_job_raw_blob(row_key, purged_at=purged_at)
            continue
        result = await _purge_upload_index_row(row)
        rows_deleted += int(result.get("rows_deleted", 0) or 0)
        blobs_deleted += int(result.get("blobs_deleted", 0) or 0)
        jobs_deleted += int(result.get("jobs_deleted", 0) or 0)
        if row_key:
            active_index_ids.discard(row_key)
            if int(result.get("jobs_deleted", 0) or 0) > 0:
                deleted_job_ids.add(row_key)

    try:
        upload_job_rows = await table_query("UploadJobs", "PartitionKey eq 'upload'", top=safe_limit)
    except Exception as e:
        logger.warning("[App] expired upload jobs sweep query failed: %s", e)
        upload_job_rows = []

    for row in upload_job_rows:
        job_id = str(row.get("RowKey", "") or "").strip()
        if not job_id:
            continue
        if job_id in deleted_job_ids:
            continue
        status = str(row.get("Status", "") or "").strip().lower()
        if status not in ("completed", "failed"):
            continue
        if job_id in active_index_ids:
            continue
        retention_until = str(row.get("RetentionUntil", "") or "").strip()
        if not _is_retention_expired(retention_until, now=now):
            raw_blob_ref = str(row.get("RawBlobRef", "") or "").strip()
            raw_retention_until = _effective_raw_blob_retention_until(row)
            has_artifact = bool(str(row.get("ArtifactBlobRef", "") or row.get("TabularArtifactBlobRef", "") or "").strip())
            if raw_blob_ref and has_artifact and _is_retention_expired(raw_retention_until, now=now):
                purged_at = now.isoformat()
                container, blob_name = parse_blob_ref(raw_blob_ref)
                if container and blob_name:
                    try:
                        await blob_delete(container, blob_name)
                        blobs_deleted += 1
                    except Exception as e:
                        logger.warning("[App] upload job raw blob early delete failed for %s: %s", raw_blob_ref, e)
                        continue
                await _clear_upload_job_raw_blob(job_id, purged_at=purged_at)
            continue
        for field in ("RawBlobRef", "TabularArtifactBlobRef"):
            blob_ref = str(row.get(field, "") or "").strip()
            if not blob_ref:
                continue
            container, blob_name = parse_blob_ref(blob_ref)
            if container and blob_name:
                try:
                    await blob_delete(container, blob_name)
                    blobs_deleted += 1
                except Exception as e:
                    logger.warning("[App] upload job blob delete failed for %s: %s", blob_ref, e)
        try:
            await upload_jobs_store.delete(job_id, partition_key="upload")
            jobs_deleted += 1
        except Exception as e:
            logger.warning("[App] stale upload job delete failed for %s: %s", job_id, e)

    return {
        "rows_deleted": rows_deleted,
        "blobs_deleted": blobs_deleted,
        "jobs_deleted": jobs_deleted,
    }


async def _backfill_tabular_artifact_chunks(limit: int = UPLOAD_TABULAR_CHUNK_BACKFILL_BATCH_SIZE) -> dict:
    safe_limit = max(1, min(int(limit or 1), 25))
    queried = 0
    completed = 0
    skipped = 0

    try:
        rows = await table_query("UploadIndex", top=min(max(safe_limit * 10, 50), 250))
    except Exception as e:
        logger.warning("[App] tabular chunk backfill query failed: %s", e)
        return {"queried": 0, "completed": 0, "skipped": 0}

    candidates = []
    for row in rows:
        filename = str(row.get("Filename", "") or "").strip()
        artifact_blob_ref = str(row.get("TabularArtifactBlobRef", "") or "").strip()
        has_chunks = bool(row.get("HasChunks", False))
        chunks_blob_ref = str(row.get("ChunksBlobRef", "") or "").strip()
        if not filename or not is_tabular_filename(filename):
            continue
        if not artifact_blob_ref or has_chunks or chunks_blob_ref:
            continue
        candidates.append(row)
        if len(candidates) >= safe_limit:
            break

    for row in candidates:
        queried += 1
        conv_id = str(row.get("PartitionKey", "") or "").strip()
        job_id = str(row.get("RowKey", "") or "").strip()
        filename = str(row.get("Filename", "") or "").strip()
        artifact_blob_ref = str(row.get("TabularArtifactBlobRef", "") or "").strip()
        if not conv_id or not job_id or not filename or not artifact_blob_ref:
            skipped += 1
            continue

        container, blob_name = parse_blob_ref(artifact_blob_ref)
        if not container or not blob_name:
            skipped += 1
            continue

        try:
            artifact_bytes = await blob_download_bytes(container, blob_name)
            columns = []
            try:
                parsed_cols = json.loads(row.get("ColNamesJson", "[]") or "[]")
                if isinstance(parsed_cols, list):
                    columns = [str(col or "").strip() for col in parsed_cols if str(col or "").strip()]
            except Exception:
                columns = []
            chunks = await _build_tabular_semantic_chunks_from_artifact(
                artifact_bytes,
                columns=columns or None,
            )
            if not chunks:
                skipped += 1
                continue

            chunks_blob_ref = str(row.get("ChunksBlobRef", "") or "").strip()
            if chunks_blob_ref:
                chunks_container, chunks_blob_name = parse_blob_ref(chunks_blob_ref)
            else:
                fallback_paths = _build_upload_blob_paths(conv_id, job_id, filename)
                chunks_container = UPLOAD_BLOB_CONTAINER_CHUNKS
                chunks_blob_name = fallback_paths["chunks_blob_name"]
            if not chunks_container or not chunks_blob_name:
                skipped += 1
                continue

            chunks_blob = await blob_upload_json(
                chunks_container,
                chunks_blob_name,
                {"chunks": chunks},
            )
            new_chunks_blob_ref = str(chunks_blob.get("blob_ref", "") or "")
            if not new_chunks_blob_ref:
                skipped += 1
                continue

            raw_blob_retention_until = _raw_blob_retention_until_iso(
                filename=filename,
                artifact_blob_ref=artifact_blob_ref,
                has_chunks=True,
                fallback_hours=UPLOAD_ARTIFACT_RETENTION_HOURS,
            )
            await table_merge(
                "UploadIndex",
                {
                    "PartitionKey": conv_id,
                    "RowKey": job_id,
                    "ChunksBlobRef": new_chunks_blob_ref,
                    "HasChunks": True,
                    "RawBlobRetentionUntil": raw_blob_retention_until[:64],
                    "RawBlobPurgedAt": "",
                },
            )
            await table_merge(
                "UploadJobs",
                {
                    "PartitionKey": "upload",
                    "RowKey": job_id,
                    "ChunksBlobName": str(chunks_blob_name)[:500],
                    "RawBlobRetentionUntil": raw_blob_retention_until[:64],
                    "RawBlobPurgedAt": "",
                    "UpdatedAt": datetime.now(timezone.utc).isoformat(),
                },
            )
            cached_job = upload_jobs_store.get(job_id)
            if cached_job:
                cached_job = dict(cached_job)
                cached_job["chunks_blob_name"] = chunks_blob_name
                cached_job["raw_blob_retention_until"] = raw_blob_retention_until
                cached_job["raw_blob_purged_at"] = ""
                await upload_jobs_store.put(job_id, cached_job)
            completed += 1
        except Exception as e:
            logger.warning("[App] tabular chunk backfill failed for %s/%s: %s", conv_id, job_id, e)
            skipped += 1

    return {"queried": queried, "completed": completed, "skipped": skipped}


async def _upload_retention_loop() -> None:
    sleep_for = max(300, int(UPLOAD_RETENTION_SWEEP_INTERVAL_SECONDS or 1800))
    while True:
        try:
            await _purge_expired_upload_artifacts()
            await _backfill_tabular_artifact_chunks()
        except asyncio.CancelledError:
            raise
        except Exception as e:
            logger.warning("[App] upload retention loop failed: %s", e)
        await asyncio.sleep(float(sleep_for))


async def _purge_upload_artifacts_for_conversation(conversation_id: str, user_sub: str = "", include_all_users: bool = False) -> dict:
    safe_conv = str(conversation_id or "").strip()
    if not safe_conv:
        return {"rows_deleted": 0, "blobs_deleted": 0, "jobs_deleted": 0}
    rows_deleted = 0
    blobs_deleted = 0
    jobs_deleted = 0
    try:
        rows = await table_query("UploadIndex", f"PartitionKey eq '{odata_escape(safe_conv)}'", top=500)
    except Exception as e:
        logger.warning("[App] upload artifact purge query failed for %s: %s", safe_conv, e)
        return {"rows_deleted": 0, "blobs_deleted": 0, "jobs_deleted": 0}

    allowed_user = str(user_sub or "").strip()
    for row in rows:
        row_owner = str(row.get("UserSub", "") or "").strip()
        if not include_all_users and allowed_user and row_owner and row_owner != allowed_user:
            continue
        result = await _purge_upload_index_row(row)
        rows_deleted += int(result.get("rows_deleted", 0) or 0)
        blobs_deleted += int(result.get("blobs_deleted", 0) or 0)
        jobs_deleted += int(result.get("jobs_deleted", 0) or 0)
    uploaded_files_store.pop(safe_conv, None)
    return {"rows_deleted": rows_deleted, "blobs_deleted": blobs_deleted, "jobs_deleted": jobs_deleted}

@app.delete("/api/chats/{user_id}/{conversation_id}")
@limiter.limit("30/minute", key_func=_user_or_ip_rate_key)
async def delete_chat(request: Request, user_id: str, conversation_id: str, credentials: HTTPAuthorizationCredentials = Depends(security)):
    user = get_current_user(credentials)
    is_admin = _is_admin_user(user)
    uid = user.get("sub") if not is_admin else user_id
    await table_delete("ChatHistory", uid, conversation_id)
    await _purge_upload_artifacts_for_conversation(
        conversation_id,
        user_sub=str(uid or ""),
        include_all_users=is_admin,
    )
    conversation_meta.pop(conversation_id, None)
    uploaded_files_store.pop(conversation_id, None)
    if conversation_id in conversations:
        del conversations[conversation_id]
    return {"status":"ok"}


@app.get("/api/privacy/export")
@limiter.limit("5/hour", key_func=_user_or_ip_rate_key)
async def export_my_data(request: Request, credentials: HTTPAuthorizationCredentials = Depends(security)):
    principal = get_current_principal(credentials, request=request)
    export_payload = await build_user_privacy_export(principal.sub)
    filename = f"dbde-privacy-export-{safe_blob_component(principal.sub, 'user')}.json"
    content = json.dumps(export_payload, ensure_ascii=False, indent=2, default=str).encode("utf-8")
    download_id = await _store_generated_file(
        content,
        "application/json",
        filename,
        "json",
        user_sub=principal.sub,
        scope="privacy_export",
    )
    await log_audit(
        principal.sub,
        "privacy_export",
        metadata={
            "mode": "privacy",
            "provider_used": "internal",
            "conversation_id": "",
            "summary": export_payload.get("summary", {}),
        },
    )
    return {
        "status": "ok",
        "download_id": download_id,
        "url": f"/api/download/{download_id}" if download_id else "",
        "filename": filename,
        "summary": export_payload.get("summary", {}),
    }


@app.post("/api/privacy/delete")
@limiter.limit("2/hour", key_func=_user_or_ip_rate_key)
async def delete_my_data(
    request: Request,
    payload: PrivacyDeleteRequest,
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    principal = get_current_principal(credentials, request=request)
    if payload.confirmation != "DELETE_MY_DATA":
        raise HTTPException(400, "Confirmação inválida.")
    result = await delete_user_personal_data(principal.sub, delete_account=bool(payload.delete_account))
    for conversation_id in result.get("conversation_ids_deleted", []) or []:
        safe_conv = str(conversation_id or "").strip()
        if not safe_conv:
            continue
        conversations.pop(safe_conv, None)
        conversation_meta.pop(safe_conv, None)
        uploaded_files_store.pop(safe_conv, None)
    if payload.delete_account:
        await persist_user_invalidation(principal.sub)
    await log_audit(
        principal.sub,
        "privacy_delete",
        metadata={
            "mode": "privacy",
            "provider_used": "internal",
            "conversation_id": "",
            "delete_account": bool(payload.delete_account),
            "summary": result,
        },
    )
    return {"status": "ok", **result}

# =============================================================================
# LEARNING ENDPOINTS
# =============================================================================

@app.post("/api/learning/rules")
@limiter.limit("20/minute", key_func=_user_or_ip_rate_key)
async def add_rule(request: Request, rule_text: str, category: str = "general", credentials: HTTPAuthorizationCredentials = Depends(security)):
    user = get_current_user(credentials)
    if user.get("role") != "admin": raise HTTPException(403, "Admin only")
    rid = f"rule_{uuid.uuid4().hex[:8]}"
    await table_insert("PromptRules", {"PartitionKey":"active","RowKey":rid,"RuleText":rule_text,"Category":category,"CreatedBy":user.get("sub"),"CreatedAt":datetime.now(timezone.utc).isoformat()})
    invalidate_prompt_rules_cache()
    return {"status":"ok","rule_id":rid}

@app.get("/api/learning/rules")
@limiter.limit("60/minute", key_func=_user_or_ip_rate_key)
async def list_rules(request: Request, credentials: HTTPAuthorizationCredentials = Depends(security)):
    user = get_current_user(credentials)
    if not _is_admin_user(user):
        raise HTTPException(403, "Admin only")
    rules = await table_query("PromptRules", "PartitionKey eq 'active'", top=50)
    return {"rules":[{"id":r.get("RowKey"),"text":r.get("RuleText"),"category":r.get("Category"),"created_by":r.get("CreatedBy")} for r in rules]}

@app.delete("/api/learning/rules/{rule_id}")
@limiter.limit("20/minute", key_func=_user_or_ip_rate_key)
async def delete_rule(request: Request, rule_id: str, credentials: HTTPAuthorizationCredentials = Depends(security)):
    user = get_current_user(credentials)
    if user.get("role") != "admin": raise HTTPException(403, "Admin only")
    await table_delete("PromptRules", "active", rule_id)
    invalidate_prompt_rules_cache()
    return {"status":"ok"}

@app.post("/api/learning/analyze")
@limiter.limit("10/minute", key_func=_user_or_ip_rate_key)
async def analyze_feedback(request: Request, credentials: HTTPAuthorizationCredentials = Depends(security)):
    user = get_current_user(credentials)
    if user.get("role") != "admin": raise HTTPException(403, "Admin only")
    fbs = await table_query("feedback", top=500)
    if not fbs: return {"analysis":"Sem feedback suficiente.","suggestions":[]}
    
    neg = [f for f in fbs if f.get("Rating",10) <= 3]
    pos = [f for f in fbs if f.get("Rating",0) >= 8]
    summary = f"Total: {len(fbs)} feedbacks. Positivos(8+): {len(pos)}. Negativos(3-): {len(neg)}.\n\n"
    if neg:
        summary += "FEEDBACK NEGATIVO:\n"
        for f in neg[:10]: summary += f"- Q: {f.get('Question','')[:80]}... Rating: {f.get('Rating')}, Nota: {f.get('Note','')}\n"
    
    try:
        analysis = await llm_simple(f"Analisa feedback de agente AI e sugere melhorias:\n\n{summary}", tier="standard", max_tokens=1500)
    except Exception as e:
        logger.warning("[App] analyze_feedback LLM failed, using summary fallback: %s", e)
        analysis = summary
    
    return {"analysis":analysis, "total":len(fbs), "positive":len(pos), "negative":len(neg)}

# =============================================================================
# INFO / HEALTH / DEBUG
# =============================================================================

def _build_admin_info_payload() -> dict:
    return {
        "service": APP_TITLE,
        "version": APP_VERSION,
        "status": "running",
        "models": {"fast": LLM_TIER_FAST, "standard": LLM_TIER_STANDARD, "pro": LLM_TIER_PRO, "vision": LLM_TIER_VISION},
        "routing": {
            "model_router_enabled": MODEL_ROUTER_ENABLED,
            "model_router_effective": bool(MODEL_ROUTER_ENABLED and (not MODEL_ROUTER_NON_PROD_ONLY or not IS_PRODUCTION)),
            "model_router_spec": MODEL_ROUTER_SPEC,
            "model_router_target_tiers": list(MODEL_ROUTER_TARGET_TIERS),
            "model_router_non_prod_only": MODEL_ROUTER_NON_PROD_ONLY,
        },
        "rerank": {
            "enabled": RERANK_ENABLED,
            "model": RERANK_MODEL,
            "endpoint_configured": bool(str(RERANK_ENDPOINT or "").strip()),
            "top_n": RERANK_TOP_N,
            "auth_mode": RERANK_AUTH_MODE,
        },
        "indexes": {"devops": DEVOPS_INDEX, "omni": OMNI_INDEX, "examples": EXAMPLES_INDEX},
        "active_tools": get_registered_tool_names(),
        "capabilities": [
            "multi_model",
            "streaming_sse",
            "jwt_cookie_auth",
            "agent_routing",
            "parallel_tools",
            "export_csv_xlsx_pdf_svg_html_zip",
            "feedback",
            "file_upload",
            "chat_persistence",
            "adaptive_learning",
        ],
        "upload_limits": {
            "max_files_per_conversation": MAX_FILES_PER_CONVERSATION,
            "max_images_per_message": UPLOAD_MAX_IMAGES_PER_MESSAGE,
            "max_file_bytes": MAX_UPLOAD_FILE_BYTES,
            "max_file_bytes_by_extension": get_tabular_upload_limits(),
            "max_batch_total_bytes": UPLOAD_MAX_BATCH_TOTAL_BYTES,
            "frontend_async_threshold_bytes": UPLOAD_FRONTEND_ASYNC_THRESHOLD_BYTES,
            "max_concurrent_jobs": MAX_CONCURRENT_UPLOAD_JOBS,
            "max_pending_jobs_per_user": UPLOAD_MAX_PENDING_JOBS_PER_USER,
            "embedding_concurrency": UPLOAD_EMBEDDING_CONCURRENCY,
            "max_chunks_per_file": UPLOAD_MAX_CHUNKS_PER_FILE,
            "job_stale_seconds": UPLOAD_JOB_STALE_SECONDS,
            "index_top": UPLOAD_INDEX_TOP,
        },
        "upload_storage": {
            "raw_container": UPLOAD_BLOB_CONTAINER_RAW,
            "text_container": UPLOAD_BLOB_CONTAINER_TEXT,
            "chunks_container": UPLOAD_BLOB_CONTAINER_CHUNKS,
            "inline_worker_enabled": INLINE_WORKER_ENABLED_EFFECTIVE,
            "inline_worker_configured": UPLOAD_INLINE_WORKER_ENABLED,
            "inline_worker_runtime_guard": INLINE_WORKER_RUNTIME_GUARD,
            "dedicated_worker_sidecar": UPLOAD_DEDICATED_WORKER_ENABLED,
            "worker_poll_seconds": UPLOAD_WORKER_POLL_SECONDS,
            "worker_batch_size": UPLOAD_WORKER_BATCH_SIZE,
        },
        "export_storage": {
            "jobs_table": "ExportJobs",
            "payload_container": CHAT_TOOLRESULT_BLOB_CONTAINER,
            "auto_async_enabled": EXPORT_AUTO_ASYNC_ENABLED,
            "async_threshold_rows": EXPORT_ASYNC_THRESHOLD_ROWS,
            "max_concurrent_jobs": EXPORT_MAX_CONCURRENT_JOBS,
            "job_stale_seconds": EXPORT_JOB_STALE_SECONDS,
            "inline_worker_enabled": EXPORT_INLINE_WORKER_ENABLED_EFFECTIVE,
            "inline_worker_configured": EXPORT_INLINE_WORKER_ENABLED,
            "dedicated_worker_sidecar": EXPORT_DEDICATED_WORKER_ENABLED,
            "worker_poll_seconds": EXPORT_WORKER_POLL_SECONDS,
            "worker_batch_size": EXPORT_WORKER_BATCH_SIZE,
        },
        "pptx_status": "ok" if Presentation is not None else "unavailable",
    }


@app.get("/api/info")
@limiter.limit("120/minute", key_func=_login_rate_key)
async def api_info(request: Request):
    return {
        "service": APP_TITLE,
        "version": APP_VERSION,
        "status": "running",
        "mode": "public",
        "timestamp_utc": datetime.now(timezone.utc).isoformat(),
        "upload_limits": {
            "max_files_per_conversation": MAX_FILES_PER_CONVERSATION,
            "max_images_per_message": UPLOAD_MAX_IMAGES_PER_MESSAGE,
            "max_file_bytes": MAX_UPLOAD_FILE_BYTES,
            "max_file_bytes_by_extension": get_tabular_upload_limits(),
            "max_batch_total_bytes": UPLOAD_MAX_BATCH_TOTAL_BYTES,
            "frontend_async_threshold_bytes": UPLOAD_FRONTEND_ASYNC_THRESHOLD_BYTES,
        },
        "features": {
            "user_story_lane": STORY_LANE_ENABLED,
            "speech_prompt": True,
            "speech_provider": "azure_speech" if (AZURE_SPEECH_ENABLED and AZURE_SPEECH_KEY and AZURE_SPEECH_REGION) else "browser_fallback",
            "speech_prompt_primary": SPEECH_PROMPT_PRIMARY_SPEC,
            "speech_prompt_fallback": SPEECH_PROMPT_FALLBACK_SPEC,
            "speech_submit_modes": ["auto", "text"],
            "provider_governance_mode": PROVIDER_GOVERNANCE_MODE,
            "provider_external_model_families": list(PROVIDER_EXTERNAL_MODEL_FAMILIES),
            "provider_external_models_experimental_allowed": PROVIDER_GOVERNANCE_EXPERIMENTAL_ALLOW_EXTERNAL,
        },
    }


@app.get("/api/admin/info")
@limiter.limit("60/minute", key_func=_user_or_ip_rate_key)
async def api_admin_info(request: Request, credentials: HTTPAuthorizationCredentials = Depends(security)):
    user = get_current_user(credentials)
    if user.get("role") != "admin":
        raise HTTPException(403, "Apenas admins")
    payload = _build_admin_info_payload()
    payload["mode"] = "admin"
    return payload


@app.get("/api/admin/tool-metrics")
async def api_admin_tool_metrics(credentials: HTTPAuthorizationCredentials = Depends(security)):
    user = get_current_user(credentials)
    if user.get("role") != "admin":
        raise HTTPException(403, "Apenas admins")
    return tool_metrics.snapshot()


@app.get("/api/admin/token-quotas")
async def api_admin_token_quotas(credentials: HTTPAuthorizationCredentials = Depends(security)):
    user = get_current_user(credentials)
    if user.get("role") != "admin":
        raise HTTPException(403, "Apenas admins")
    mgr = _tq_module.token_quota_manager
    if not mgr:
        return {"error": "Token quota manager not initialised"}
    return await mgr.snapshot()


@app.get("/api/admin/user-stories/eval-summary")
async def api_admin_user_story_eval_summary(
    credentials: HTTPAuthorizationCredentials = Depends(security),
    top: int = 250,
    user_sub: Optional[str] = None,
):
    user = get_current_user(credentials)
    if user.get("role") != "admin":
        raise HTTPException(403, "Apenas admins")
    return await build_user_story_eval_summary(user_sub=str(user_sub or "").strip(), top=top)


@app.post("/api/admin/user-stories/promote-candidate")
async def api_admin_user_story_promote_candidate(
    payload: UserStoryPromoteRequest,
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    user = get_current_user(credentials)
    if user.get("role") != "admin":
        raise HTTPException(403, "Apenas admins")
    try:
        return await promote_user_story_to_curated_corpus(
            draft_id=payload.draft_id,
            source_user_sub=payload.user_sub,
            promoted_by=str(user.get("sub", "") or user.get("username", "") or "admin"),
            note=str(payload.note or "").strip(),
        )
    except RuntimeError as exc:
        raise HTTPException(400, str(exc))


@app.post("/api/admin/user-stories/review-candidate")
async def api_admin_user_story_review_candidate(
    payload: UserStoryCurationReviewRequest,
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    user = get_current_user(credentials)
    if user.get("role") != "admin":
        raise HTTPException(403, "Apenas admins")
    try:
        return await review_user_story_curated_candidate(
            draft_id=payload.draft_id,
            action=payload.action,
            reviewed_by=str(user.get("sub", "") or user.get("username", "") or "admin"),
            note=str(payload.note or "").strip(),
        )
    except RuntimeError as exc:
        raise HTTPException(400, str(exc))


@app.post("/api/admin/user-stories/sync-search-index")
async def api_admin_user_story_sync_search_index(
    payload: UserStorySearchSyncRequest,
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    user = get_current_user(credentials)
    if user.get("role") != "admin":
        raise HTTPException(403, "Apenas admins")
    try:
        return await sync_user_story_examples_search_index(
            draft_id=str(payload.draft_id or "").strip(),
            top=int(payload.top or 200),
        )
    except RuntimeError as exc:
        raise HTTPException(400, str(exc))


@app.post("/api/admin/user-stories/sync-devops-index")
async def api_admin_user_story_sync_devops_index(
    payload: UserStoryDevOpsSyncRequest,
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    user = get_current_user(credentials)
    if user.get("role") != "admin":
        raise HTTPException(403, "Apenas admins")
    try:
        return await sync_story_devops_index(
            since_iso=str(payload.since_iso or "").strip(),
            since_days=int(payload.since_days or 30),
            top=int(payload.top or 1200),
            update_cursor=bool(payload.update_cursor),
        )
    except RuntimeError as exc:
        raise HTTPException(400, str(exc))


@app.post("/api/admin/user-stories/sync-knowledge-index")
async def api_admin_user_story_sync_knowledge_index(
    payload: UserStoryKnowledgeSyncRequest,
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    user = get_current_user(credentials)
    if user.get("role") != "admin":
        raise HTTPException(403, "Apenas admins")
    try:
        return await sync_story_knowledge_index(
            max_docs=int(payload.max_docs or 1500),
            batch_size=int(payload.batch_size or 150),
            update_state=bool(payload.update_state),
        )
    except RuntimeError as exc:
        raise HTTPException(400, str(exc))


@app.get("/api/admin/user-stories/index-status")
async def api_admin_user_story_index_status(
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    user = get_current_user(credentials)
    if user.get("role") != "admin":
        raise HTTPException(403, "Apenas admins")
    return await get_story_lane_index_status()


@app.get("/api/admin/user-stories/knowledge-assets")
async def api_admin_user_story_list_knowledge_assets(
    top: int = 100,
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    user = get_current_user(credentials)
    if user.get("role") != "admin":
        raise HTTPException(403, "Apenas admins")
    return await list_story_knowledge_assets(top=max(1, min(int(top or 100), 500)))


@app.post("/api/admin/user-stories/knowledge-assets/import-upload")
async def api_admin_user_story_import_knowledge_asset_upload(
    payload: UserStoryKnowledgeAssetUploadRequest,
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    user = get_current_user(credentials)
    if user.get("role") != "admin":
        raise HTTPException(403, "Apenas admins")
    try:
        return await create_story_knowledge_asset_from_upload(
            conversation_id=payload.conversation_id,
            file_id=payload.file_id,
            imported_by=str(user.get("sub", "") or user.get("username", "") or "admin"),
            title=str(payload.title or "").strip(),
            domain=str(payload.domain or "").strip(),
            journey=str(payload.journey or "").strip(),
            flow=str(payload.flow or "").strip(),
            team_scope=str(payload.team_scope or "").strip(),
            note=str(payload.note or "").strip(),
        )
    except RuntimeError as exc:
        raise HTTPException(400, str(exc))


@app.post("/api/admin/user-stories/knowledge-assets/import-text")
async def api_admin_user_story_import_knowledge_asset_text(
    payload: UserStoryKnowledgeAssetTextRequest,
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    user = get_current_user(credentials)
    if user.get("role") != "admin":
        raise HTTPException(403, "Apenas admins")
    try:
        return await create_story_knowledge_asset_from_text(
            title=payload.title,
            content=payload.content,
            imported_by=str(user.get("sub", "") or user.get("username", "") or "admin"),
            asset_key=str(payload.asset_key or "").strip(),
            domain=str(payload.domain or "").strip(),
            journey=str(payload.journey or "").strip(),
            flow=str(payload.flow or "").strip(),
            team_scope=str(payload.team_scope or "").strip(),
            note=str(payload.note or "").strip(),
        )
    except RuntimeError as exc:
        raise HTTPException(400, str(exc))


@app.post("/api/admin/user-stories/knowledge-assets/import-bundle")
async def api_admin_user_story_import_knowledge_asset_bundle(
    payload: UserStoryKnowledgeAssetBundleRequest,
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    user = get_current_user(credentials)
    if user.get("role") != "admin":
        raise HTTPException(403, "Apenas admins")
    try:
        return await create_story_knowledge_assets_from_bundle(
            entries=[item.model_dump() for item in payload.items],
            imported_by=str(user.get("sub", "") or user.get("username", "") or "admin"),
        )
    except RuntimeError as exc:
        raise HTTPException(400, str(exc))


@app.post("/api/admin/user-stories/knowledge-assets/review")
async def api_admin_user_story_review_knowledge_asset(
    payload: UserStoryKnowledgeAssetReviewRequest,
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    user = get_current_user(credentials)
    if user.get("role") != "admin":
        raise HTTPException(403, "Apenas admins")
    try:
        return await review_story_knowledge_asset(
            asset_id=payload.asset_id,
            action=payload.action,
            reviewed_by=str(user.get("sub", "") or user.get("username", "") or "admin"),
            note=str(payload.note or "").strip(),
        )
    except RuntimeError as exc:
        raise HTTPException(400, str(exc))


@app.post("/api/client-error")
@limiter.limit("10/minute", key_func=_user_or_ip_rate_key)
async def report_client_error(request: Request, report: ClientErrorReport):
    """Receive and log frontend errors for observability."""
    user = None
    try:
        user = _auth_payload_from_request(request)
    except Exception:
        user = None
    user_id = str((user or {}).get("sub", "") or "anonymous")

    logger.warning(
        json.dumps(
            {
                "event": "client_error",
                "user_id": user_id,
                "error_type": report.error_type,
                "message": report.message[:500],
                "component": report.component or "",
                "url": report.url or "",
                "timestamp": report.timestamp or "",
                "has_stack": bool(report.stack),
            },
            ensure_ascii=False,
        )
    )
    return {"status": "logged"}


@app.get("/api/runtime/check")
async def runtime_check(credentials: HTTPAuthorizationCredentials = Depends(security)):
    user = get_current_user(credentials)
    if user.get("role") != "admin":
        raise HTTPException(403, "Apenas admins")

    targets = [
        BASE_DIR / "app.py",
        BASE_DIR / "agent.py",
        BASE_DIR / "tools.py",
        BASE_DIR / "config.py",
        BASE_DIR / "static" / "index.html",
    ]
    files = {}
    for p in targets:
        key = str(p.relative_to(BASE_DIR))
        try:
            full_sha = _file_sha256(p) if p.exists() else ""
            files[key] = {
                "exists": p.exists(),
                "sha256": full_sha[:24],
                "sha256_full": full_sha,
                "size_bytes": p.stat().st_size if p.exists() else 0,
            }
        except Exception as e:
            files[key] = {"exists": False, "error": str(e)}

    markers = {}
    try:
        index_txt = (BASE_DIR / "static" / "index.html").read_text(encoding="utf-8", errors="replace")
        markers = {
            "max_files_10": "MAX_FILES_PER_CONVERSATION = 10" in index_txt,
            "upload_async_frontend": "/upload/async" in index_txt,
            "upload_status_frontend": "/api/upload/status/" in index_txt,
        }
    except Exception as e:
        markers = {"error": str(e)}

    manifest_result = {}
    manifest_path = BASE_DIR / "deploy" / "runtime-manifest.json"
    if manifest_path.exists():
        try:
            expected = json.loads(manifest_path.read_text(encoding="utf-8"))
            expected_files = expected.get("files", {}) if isinstance(expected, dict) else {}
            drift = {}
            for rel_path, exp_hash in expected_files.items():
                actual = (files.get(rel_path) or {}).get("sha256_full", "")
                drift[rel_path] = {
                    "expected": str(exp_hash or ""),
                    "actual": str(actual or ""),
                    "match": bool(actual) and str(actual) == str(exp_hash),
                }
            manifest_result = {
                "path": str(manifest_path.relative_to(BASE_DIR)),
                "has_manifest": True,
                "all_match": all(v.get("match") for v in drift.values()) if drift else False,
                "files": drift,
            }
        except Exception as e:
            manifest_result = {"path": str(manifest_path.relative_to(BASE_DIR)), "has_manifest": True, "error": str(e)}
    else:
        manifest_result = {"path": str(manifest_path.relative_to(BASE_DIR)), "has_manifest": False}

    return {
        "service": APP_TITLE,
        "version": APP_VERSION,
        "runtime_utc": datetime.now(timezone.utc).isoformat(),
        "files": files,
        "markers": markers,
        "manifest_check": manifest_result,
    }


@app.get("/api/debug/upload-jobs")
async def debug_upload_jobs(
    request: Request,
    credentials: HTTPAuthorizationCredentials = Depends(security),
):
    """Temporary diagnostic endpoint — shows recent upload job states."""
    user = get_current_user(credentials, request=request)
    if user.get("role") != "admin":
        raise HTTPException(status_code=403, detail="Admin only")
    now = datetime.now(timezone.utc)
    jobs_local = []
    for job_id, job in upload_jobs_store.items():
        jobs_local.append({
            "job_id": str(job_id)[:12],
            "status": str(job.get("status", "")),
            "filename": str(job.get("filename", ""))[:40],
            "size_mb": round(int(job.get("size_bytes", 0) or 0) / 1024 / 1024, 1),
            "worker_id": str(job.get("worker_id", ""))[:16],
            "created_at": str(job.get("created_at", ""))[:25],
            "updated_at": str(job.get("updated_at", ""))[:25],
            "error": str(job.get("error", ""))[:100],
        })
    # Also check Table Storage
    storage_jobs = []
    try:
        rows = await table_query(
            "UploadJobs",
            "PartitionKey eq 'upload'",
            top=20,
        )
        for row in rows:
            status = str(row.get("Status", ""))
            if status.lower() in ("completed", "failed"):
                age = ""
                try:
                    updated = datetime.fromisoformat(str(row.get("UpdatedAt", "")))
                    age = f"{int((now - updated).total_seconds())}s ago"
                except (ValueError, TypeError):
                    pass
                if age and int((now - updated).total_seconds()) > 300:
                    continue  # skip old completed/failed
            storage_jobs.append({
                "job_id": str(row.get("RowKey", ""))[:12],
                "status": status,
                "filename": str(row.get("Filename", ""))[:40],
                "size_mb": round(int(row.get("SizeBytes", 0) or 0) / 1024 / 1024, 1),
                "worker_id": str(row.get("WorkerId", ""))[:16],
                "created_at": str(row.get("CreatedAt", ""))[:25],
                "updated_at": str(row.get("UpdatedAt", ""))[:25],
            })
    except Exception as e:
        storage_jobs = [{"error": str(e)[:200]}]
    return {
        "now_utc": now.isoformat(),
        "worker_instance": WORKER_INSTANCE_ID[:16],
        "inline_worker_enabled": INLINE_WORKER_ENABLED_EFFECTIVE,
        "worker_task_alive": _inline_worker_task is not None and not _inline_worker_task.done() if _inline_worker_task else False,
        "local_jobs": jobs_local,
        "storage_jobs": storage_jobs,
    }


@app.get("/health")
@limiter.limit("120/minute", key_func=_login_rate_key)
async def health(
    request: Request,
    deep: bool = False,
    credentials: Optional[HTTPAuthorizationCredentials] = Depends(security),
):
    result = {"status": "healthy", "mode": "basic", "checks": {"app": "ok"}}
    if not deep:
        return result

    try:
        user = get_current_user(credentials)
    except Exception:
        return JSONResponse(
            status_code=401,
            content={"detail": "Token inválido para deep health check"},
        )
    if not _is_admin_user(user):
        raise HTTPException(403, "Apenas admins")

    result["mode"] = "deep"
    checks = {}

    # 1) Table Storage
    try:
        await table_query("feedback", top=1)
        checks["table_storage"] = "ok"
    except Exception as e:
        checks["table_storage"] = f"error: {str(e)[:100]}"

    # 2) Blob Storage
    try:
        probe_blob = "__health_probe__.txt"
        _ = await blob_download_bytes(UPLOAD_BLOB_CONTAINER_RAW, probe_blob)
        checks["blob_storage"] = "ok"
    except Exception as e:
        checks["blob_storage"] = f"error: {str(e)[:100]}"

    # 3) Azure OpenAI (tier fast)
    try:
        _ = await llm_simple("ping", tier="fast", max_tokens=5)
        checks["llm_fast"] = "ok"
    except Exception as e:
        checks["llm_fast"] = f"error: {str(e)[:100]}"

    # 3b) Azure OpenAI (tier vision)
    try:
        if VISION_ENABLED:
            _ = await llm_simple("ping", tier="vision", max_tokens=5)
            checks["llm_vision"] = "ok"
        else:
            checks["llm_vision"] = "disabled"
    except Exception as e:
        checks["llm_vision"] = f"error: {str(e)[:100]}"

    # 4) Azure AI Search
    try:
        from http_helpers import search_request_with_retry
        url = f"https://{SEARCH_SERVICE}.search.windows.net/indexes/{DEVOPS_INDEX}/docs/search?api-version={API_VERSION_SEARCH}"
        headers = {"Content-Type": "application/json", "api-key": SEARCH_KEY}
        payload = {"search": "*", "top": 1}
        search_resp = await search_request_with_retry(url=url, headers=headers, json_body=payload, max_retries=2)
        checks["ai_search"] = "ok" if "error" not in search_resp else f"error: {str(search_resp.get('error', 'unknown'))[:100]}"
    except Exception as e:
        checks["ai_search"] = f"error: {str(e)[:100]}"

    # 5) Rerank
    try:
        if RERANK_ENABLED:
            endpoint_configured = bool(str(RERANK_ENDPOINT or "").strip())
            checks["rerank"] = "configured" if endpoint_configured else "error: missing endpoint"
        else:
            checks["rerank"] = "disabled"
    except Exception as e:
        checks["rerank"] = f"error: {str(e)[:100]}"

    # 6) Dedicated workers (when enabled)
    try:
        upload_worker_enabled = UPLOAD_DEDICATED_WORKER_ENABLED
        if upload_worker_enabled:
            pid = _load_pid_from_file(UPLOAD_WORKER_PID_FILE)
            checks["upload_worker"] = "ok" if _is_process_alive(pid) else "error: worker_not_running"
        else:
            checks["upload_worker"] = "disabled"
    except Exception as e:
        checks["upload_worker"] = f"error: {str(e)[:100]}"

    try:
        export_worker_enabled = EXPORT_DEDICATED_WORKER_ENABLED
        if export_worker_enabled:
            pid = _load_pid_from_file(EXPORT_WORKER_PID_FILE)
            checks["export_worker"] = "ok" if _is_process_alive(pid) else "error: worker_not_running"
        else:
            checks["export_worker"] = "disabled"
    except Exception as e:
        checks["export_worker"] = f"error: {str(e)[:100]}"

    result["checks"] = checks
    all_ok = all(v == "ok" or v in ("configured", "disabled") for v in checks.values())
    result["status"] = "healthy" if all_ok else "degraded"
    return JSONResponse(status_code=200 if all_ok else 503, content=result)

@app.get("/debug/conversations")
async def debug_conversations(credentials: HTTPAuthorizationCredentials = Depends(security)):
    user = get_current_user(credentials)
    if user.get("role") != "admin": raise HTTPException(403)
    return {cid: {"mode":conversation_meta.get(cid,{}).get("mode"), "msgs":len(msgs), "has_file":cid in uploaded_files_store} for cid,msgs in conversations.items()}

# =============================================================================
# FRONTEND
# =============================================================================

@app.get("/")
async def root():
    index_candidates = [
        BASE_DIR / "static" / "index.html",
        Path("/home/site/wwwroot/static/index.html"),
    ]
    for index_path in index_candidates:
        if index_path.exists():
            return HTMLResponse(content=index_path.read_text(encoding="utf-8", errors="replace"))
    return HTMLResponse(content=f"<h1>{APP_TITLE} v{APP_VERSION}</h1><p>Frontend not deployed. Use /docs for API.</p>")
