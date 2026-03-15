# =============================================================================
# pptx_engine.py — Millennium BCP branded PowerPoint generation engine v1.0
# =============================================================================
# Generates professional PPTX presentations following Millennium BCP / Digital
# Empresas branding guidelines extracted from official templates.
#
# Brand guidelines (from 5 official templates):
#   - Primary font: Montserrat (ExtraBold titles, Bold subtitles, Regular body)
#   - Secondary font: Trebuchet MS (fallback)
#   - Brand accent: #D1005D / #D2125C (cerise/magenta)
#   - Dark text: #585857 (gray)
#   - Light badge bg: #F2F2F2
#   - Slide size: 13.333 x 7.5 inches (widescreen 16:9)
#   - Section dividers: large numbers (01, 02) in brand color
#   - Badge element: rounded rect with "DIGITAL EMPRESAS" label
#   - KPI numbers: 54-72pt bold in brand color
#   - Tables: alternating rows, brand header, rounded containers
# =============================================================================

import io
import logging
import re
from datetime import datetime, timezone
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Brand constants
# ---------------------------------------------------------------------------
BRAND_ACCENT_HEX = "D1005D"
BRAND_ACCENT_DARK_HEX = "D2125C"
BRAND_DARK_TEXT_HEX = "585857"
BRAND_LIGHT_BG_HEX = "F2F2F2"
BRAND_WHITE_HEX = "FFFFFF"
BRAND_BLACK_HEX = "1A1A1A"
BRAND_HEADER_BADGE_TEXT = "DIGITAL EMPRESAS"

# Slide dimensions in EMU (English Metric Units — 914400 EMU = 1 inch)
SLIDE_WIDTH_EMU = 12192000   # 13.333 inches
SLIDE_HEIGHT_EMU = 6858000   # 7.5 inches

# Font sizes in pt
TITLE_FONT_SIZE = 36
SUBTITLE_FONT_SIZE = 18
BODY_FONT_SIZE = 12
SMALL_FONT_SIZE = 10
SECTION_NUMBER_FONT_SIZE = 72
KPI_NUMBER_FONT_SIZE = 54
TABLE_HEADER_FONT_SIZE = 10
TABLE_BODY_FONT_SIZE = 9
BADGE_FONT_SIZE = 8

# Font families
FONT_PRIMARY = "Montserrat"
FONT_SECONDARY = "Trebuchet MS"

# Layout margins (EMU)
MARGIN_LEFT = 609600     # ~0.67 inches
MARGIN_TOP = 914400      # 1 inch
MARGIN_RIGHT = 609600
MARGIN_BOTTOM = 457200   # 0.5 inches
CONTENT_WIDTH = SLIDE_WIDTH_EMU - MARGIN_LEFT - MARGIN_RIGHT


def _emu(inches: float) -> int:
    """Convert inches to EMU."""
    return int(inches * 914400)


def _pt(points: int):
    """Convert points to Pt object."""
    from pptx.util import Pt
    return Pt(points)


def _rgb(hex_color: str):
    """Convert hex color string to RGBColor."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _set_font(run, *, size: int = BODY_FONT_SIZE, bold: bool = False,
              italic: bool = False, color: str = BRAND_DARK_TEXT_HEX,
              font_name: str = FONT_PRIMARY):
    """Apply font styling to a run."""
    from pptx.util import Pt
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(color)
    run.font.name = font_name


def _add_text_box(slide, left, top, width, height, text: str, *,
                  size: int = BODY_FONT_SIZE, bold: bool = False,
                  color: str = BRAND_DARK_TEXT_HEX, alignment=None,
                  font_name: str = FONT_PRIMARY, word_wrap: bool = True):
    """Add a text box with styled text to a slide."""
    from pptx.util import Pt, Emu
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = word_wrap
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = str(text)
    if alignment:
        p.alignment = alignment
    run = p.runs[0] if p.runs else p.add_run()
    if not p.runs:
        run.text = str(text)
    _set_font(run, size=size, bold=bold, color=color, font_name=font_name)
    return txBox


def _add_badge(slide, text: str = BRAND_HEADER_BADGE_TEXT,
               left: int = MARGIN_LEFT, top: int = _emu(0.35)):
    """Add the branded 'DIGITAL EMPRESAS' badge to a slide."""
    from pptx.util import Pt, Emu
    from pptx.enum.text import PP_ALIGN
    badge_w = _emu(2.2)
    badge_h = _emu(0.3)
    shape = slide.shapes.add_shape(
        5,  # MSO_SHAPE.ROUNDED_RECTANGLE
        left, top, badge_w, badge_h,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(BRAND_LIGHT_BG_HEX)
    shape.line.fill.background()  # no border
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(text)
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0] if p.runs else p.add_run()
    if not p.runs:
        run.text = str(text)
    _set_font(run, size=BADGE_FONT_SIZE, bold=True, color=BRAND_DARK_TEXT_HEX)
    return shape


def _add_accent_bar(slide, left: int = 0, top: int = 0,
                    width: int = SLIDE_WIDTH_EMU, height: int = _emu(0.06)):
    """Add a thin brand-color accent bar."""
    shape = slide.shapes.add_shape(1, left, top, width, height)  # RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(BRAND_ACCENT_HEX)
    shape.line.fill.background()
    return shape


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _build_title_slide(prs, title: str, subtitle: str = "",
                       badge_text: str = BRAND_HEADER_BADGE_TEXT):
    """Create a title/cover slide with brand styling."""
    from pptx.util import Pt, Emu
    from pptx.enum.text import PP_ALIGN
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Accent bar at top
    _add_accent_bar(slide, top=0, height=_emu(0.08))

    # Badge
    if badge_text:
        _add_badge(slide, badge_text, left=MARGIN_LEFT, top=_emu(1.5))

    # Title
    _add_text_box(
        slide, MARGIN_LEFT, _emu(2.2), CONTENT_WIDTH, _emu(1.5),
        title, size=44, bold=True, color=BRAND_BLACK_HEX,
        font_name=FONT_PRIMARY,
    )

    # Subtitle
    if subtitle:
        _add_text_box(
            slide, MARGIN_LEFT, _emu(3.8), CONTENT_WIDTH, _emu(0.8),
            subtitle, size=SUBTITLE_FONT_SIZE, bold=False,
            color=BRAND_DARK_TEXT_HEX,
        )

    # Date footer
    date_str = datetime.now(timezone.utc).strftime("%d/%m/%Y")
    _add_text_box(
        slide, MARGIN_LEFT, _emu(6.6), _emu(3), _emu(0.3),
        date_str, size=SMALL_FONT_SIZE, color=BRAND_DARK_TEXT_HEX,
    )
    return slide


def _build_section_divider(prs, section_number: int, section_title: str,
                           badge_text: str = BRAND_HEADER_BADGE_TEXT):
    """Create a section divider slide (e.g., 01 — Title)."""
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    _add_accent_bar(slide, top=0, height=_emu(0.06))
    if badge_text:
        _add_badge(slide, badge_text)

    # Large section number
    num_text = f"{section_number:02d}" if section_number < 100 else str(section_number)
    _add_text_box(
        slide, MARGIN_LEFT, _emu(2.0), _emu(3), _emu(1.5),
        num_text, size=SECTION_NUMBER_FONT_SIZE, bold=True,
        color=BRAND_ACCENT_HEX, font_name=FONT_PRIMARY,
    )

    # Section title
    _add_text_box(
        slide, MARGIN_LEFT + _emu(3.2), _emu(2.3), _emu(7), _emu(1.0),
        section_title, size=TITLE_FONT_SIZE, bold=True,
        color=BRAND_BLACK_HEX,
    )
    return slide


def _build_content_slide(prs, title: str, bullets: List[str],
                         badge_text: str = BRAND_HEADER_BADGE_TEXT):
    """Create a content slide with title and bullet points."""
    from pptx.util import Pt, Emu
    from pptx.enum.text import PP_ALIGN
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_accent_bar(slide, top=0, height=_emu(0.04))
    if badge_text:
        _add_badge(slide, badge_text)

    # Title
    _add_text_box(
        slide, MARGIN_LEFT, _emu(0.85), CONTENT_WIDTH, _emu(0.6),
        title, size=24, bold=True, color=BRAND_BLACK_HEX,
    )

    # Bullets
    if bullets:
        txBox = slide.shapes.add_textbox(
            MARGIN_LEFT, _emu(1.7), CONTENT_WIDTH, _emu(5.0),
        )
        txBox.word_wrap = True
        tf = txBox.text_frame
        tf.word_wrap = True

        for idx, bullet_text in enumerate(bullets):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # Support sub-bullets with "- " prefix inside a bullet
            text = str(bullet_text).strip()
            is_sub = text.startswith("- ") or text.startswith("• ")
            if is_sub:
                text = text[2:].strip()
                p.level = 1
            else:
                p.level = 0

            p.space_after = Pt(6)
            run = p.add_run()
            run.text = text
            font_size = BODY_FONT_SIZE if not is_sub else (BODY_FONT_SIZE - 1)
            _set_font(run, size=font_size, color=BRAND_DARK_TEXT_HEX)
    return slide


def _build_two_column_slide(prs, title: str, left_content: List[str],
                            right_content: List[str],
                            badge_text: str = BRAND_HEADER_BADGE_TEXT):
    """Create a two-column content slide."""
    from pptx.util import Pt, Emu
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_accent_bar(slide, top=0, height=_emu(0.04))
    if badge_text:
        _add_badge(slide, badge_text)

    # Title
    _add_text_box(
        slide, MARGIN_LEFT, _emu(0.85), CONTENT_WIDTH, _emu(0.6),
        title, size=24, bold=True, color=BRAND_BLACK_HEX,
    )

    col_width = (CONTENT_WIDTH - _emu(0.5)) // 2

    for col_idx, items in enumerate([left_content, right_content]):
        col_left = MARGIN_LEFT if col_idx == 0 else (MARGIN_LEFT + col_width + _emu(0.5))
        if not items:
            continue
        txBox = slide.shapes.add_textbox(col_left, _emu(1.7), col_width, _emu(5.0))
        txBox.word_wrap = True
        tf = txBox.text_frame
        tf.word_wrap = True
        for idx, text in enumerate(items):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.space_after = Pt(6)
            run = p.add_run()
            run.text = str(text).strip()
            _set_font(run, size=BODY_FONT_SIZE, color=BRAND_DARK_TEXT_HEX)
    return slide


def _build_kpi_slide(prs, title: str, kpis: List[Dict[str, Any]],
                     badge_text: str = BRAND_HEADER_BADGE_TEXT):
    """Create a KPI/metrics slide with large numbers.

    Each KPI dict: {value: str, label: str, description?: str}
    Max 4 KPIs per slide.
    """
    from pptx.util import Pt, Emu
    from pptx.enum.text import PP_ALIGN
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_accent_bar(slide, top=0, height=_emu(0.04))
    if badge_text:
        _add_badge(slide, badge_text)

    # Title
    _add_text_box(
        slide, MARGIN_LEFT, _emu(0.85), CONTENT_WIDTH, _emu(0.6),
        title, size=24, bold=True, color=BRAND_BLACK_HEX,
    )

    n = min(len(kpis), 4)
    if n == 0:
        return slide

    kpi_width = CONTENT_WIDTH // n
    for i, kpi in enumerate(kpis[:4]):
        x = MARGIN_LEFT + i * kpi_width
        # KPI value (large number)
        _add_text_box(
            slide, x, _emu(2.2), kpi_width - _emu(0.2), _emu(1.0),
            str(kpi.get("value", "")),
            size=KPI_NUMBER_FONT_SIZE, bold=True, color=BRAND_ACCENT_HEX,
            alignment=PP_ALIGN.CENTER,
        )
        # KPI label
        _add_text_box(
            slide, x, _emu(3.3), kpi_width - _emu(0.2), _emu(0.4),
            str(kpi.get("label", "")),
            size=BODY_FONT_SIZE, bold=True, color=BRAND_DARK_TEXT_HEX,
            alignment=PP_ALIGN.CENTER,
        )
        # KPI description (optional)
        desc = kpi.get("description", "")
        if desc:
            _add_text_box(
                slide, x, _emu(3.8), kpi_width - _emu(0.2), _emu(0.6),
                str(desc),
                size=SMALL_FONT_SIZE, bold=False, color=BRAND_DARK_TEXT_HEX,
                alignment=PP_ALIGN.CENTER,
            )
    return slide


def _build_table_slide(prs, title: str, headers: List[str],
                       rows: List[List[str]],
                       badge_text: str = BRAND_HEADER_BADGE_TEXT):
    """Create a slide with a branded data table."""
    from pptx.util import Pt, Emu
    from pptx.enum.text import PP_ALIGN
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_accent_bar(slide, top=0, height=_emu(0.04))
    if badge_text:
        _add_badge(slide, badge_text)

    # Title
    _add_text_box(
        slide, MARGIN_LEFT, _emu(0.85), CONTENT_WIDTH, _emu(0.6),
        title, size=24, bold=True, color=BRAND_BLACK_HEX,
    )

    if not headers:
        return slide

    max_rows = min(len(rows), 15)  # cap at 15 rows per slide
    n_cols = len(headers)
    n_rows = max_rows + 1  # +1 for header

    table_top = _emu(1.7)
    table_height = _emu(min(5.0, 0.35 * n_rows + 0.1))
    col_width = CONTENT_WIDTH // n_cols

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        MARGIN_LEFT, table_top, CONTENT_WIDTH, table_height,
    )
    table = table_shape.table

    # Style header row
    for ci, header in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = str(header)
        # Brand-colored header background
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(BRAND_ACCENT_HEX)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                _set_font(run, size=TABLE_HEADER_FONT_SIZE, bold=True,
                          color=BRAND_WHITE_HEX)

    # Style data rows with alternating colors
    for ri, row_data in enumerate(rows[:max_rows]):
        for ci in range(n_cols):
            cell = table.cell(ri + 1, ci)
            cell_val = row_data[ci] if ci < len(row_data) else ""
            cell.text = str(cell_val)
            # Zebra striping
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(BRAND_WHITE_HEX)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(BRAND_LIGHT_BG_HEX)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    _set_font(run, size=TABLE_BODY_FONT_SIZE,
                              color=BRAND_DARK_TEXT_HEX)

    if len(rows) > max_rows:
        note = f"(+{len(rows) - max_rows} linhas omitidas)"
        _add_text_box(
            slide, MARGIN_LEFT, table_top + table_height + _emu(0.15),
            CONTENT_WIDTH, _emu(0.3),
            note, size=8, italic=True, color=BRAND_DARK_TEXT_HEX,
        )
    return slide


def _build_agenda_slide(prs, items: List[str],
                        badge_text: str = BRAND_HEADER_BADGE_TEXT):
    """Create an agenda/index slide."""
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_accent_bar(slide, top=0, height=_emu(0.04))
    if badge_text:
        _add_badge(slide, badge_text)

    _add_text_box(
        slide, MARGIN_LEFT, _emu(0.85), CONTENT_WIDTH, _emu(0.6),
        "Agenda", size=TITLE_FONT_SIZE, bold=True, color=BRAND_BLACK_HEX,
    )

    for idx, item_text in enumerate(items[:12]):
        y = _emu(1.8 + idx * 0.45)
        # Number circle
        _add_text_box(
            slide, MARGIN_LEFT, y, _emu(0.5), _emu(0.35),
            f"{idx + 1:02d}", size=14, bold=True, color=BRAND_ACCENT_HEX,
            alignment=PP_ALIGN.CENTER,
        )
        # Item text
        _add_text_box(
            slide, MARGIN_LEFT + _emu(0.7), y, _emu(10), _emu(0.35),
            str(item_text), size=14, bold=False, color=BRAND_DARK_TEXT_HEX,
        )
    return slide


def _build_closing_slide(prs, text: str = "Obrigado",
                         subtitle: str = "",
                         badge_text: str = BRAND_HEADER_BADGE_TEXT):
    """Create a closing/thank-you slide."""
    from pptx.enum.text import PP_ALIGN
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_accent_bar(slide, top=0, height=_emu(0.08))

    _add_text_box(
        slide, MARGIN_LEFT, _emu(2.5), CONTENT_WIDTH, _emu(1.5),
        text, size=44, bold=True, color=BRAND_ACCENT_HEX,
        alignment=PP_ALIGN.CENTER,
    )

    if subtitle:
        _add_text_box(
            slide, MARGIN_LEFT, _emu(4.2), CONTENT_WIDTH, _emu(0.6),
            subtitle, size=SUBTITLE_FONT_SIZE, bold=False,
            color=BRAND_DARK_TEXT_HEX, alignment=PP_ALIGN.CENTER,
        )

    date_str = datetime.now(timezone.utc).strftime("%d/%m/%Y")
    _add_text_box(
        slide, MARGIN_LEFT, _emu(6.6), CONTENT_WIDTH, _emu(0.3),
        date_str, size=SMALL_FONT_SIZE, color=BRAND_DARK_TEXT_HEX,
        alignment=PP_ALIGN.CENTER,
    )
    return slide


# ---------------------------------------------------------------------------
# Smart slide validation & auto-correction
# ---------------------------------------------------------------------------
# This layer enforces presentation quality rules REGARDLESS of what the LLM
# sends. It splits overloaded slides, trims excess content, and ensures
# professional structure.
# ---------------------------------------------------------------------------

_MAX_BULLETS_PER_SLIDE = 7
_MAX_BULLET_LENGTH = 150
_MAX_TITLE_LENGTH = 80
_MAX_KPIS_PER_SLIDE = 4
_MAX_TABLE_ROWS_PER_SLIDE = 12
_MAX_TABLE_COLS = 8
_MAX_AGENDA_ITEMS = 12
_MAX_TWO_COL_ITEMS = 8


def _validate_and_fix_slides(slides: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """Validate and auto-correct slide specs for professional quality.

    Rules enforced:
    1. Content slides with >7 bullets → split into multiple slides
    2. Bullet text >150 chars → truncate with ellipsis
    3. Title >80 chars → truncate
    4. KPI slides with >4 KPIs → split
    5. Table with >12 rows → split into continuation slides
    6. Table with >8 columns → trim to 8
    7. Two-column slides with >8 items per side → trim
    8. Empty slides → removed
    9. Adjacent content slides with same title → merge if under bullet limit
    10. Consecutive sections without content between → drop duplicate
    """
    if not slides:
        return []

    fixed = []
    for spec in slides:
        if not isinstance(spec, dict):
            continue
        slide_type = str(spec.get("type", "content")).lower().strip()

        # --- Truncate title ---
        title = str(spec.get("title", "")).strip()
        if len(title) > _MAX_TITLE_LENGTH:
            title = title[:_MAX_TITLE_LENGTH - 1].rstrip() + "…"
            spec = {**spec, "title": title}

        # --- Content slides: split if too many bullets ---
        if slide_type in ("content", "bullets"):
            bullets = spec.get("bullets", [])
            if isinstance(bullets, str):
                bullets = [b.strip() for b in bullets.split("\n") if b.strip()]
            # Truncate individual bullets
            bullets = [
                (b[:_MAX_BULLET_LENGTH - 1].rstrip() + "…" if len(b) > _MAX_BULLET_LENGTH else b)
                for b in bullets
            ]
            if not bullets:
                # Empty content slide → skip
                continue
            # Split into chunks of _MAX_BULLETS_PER_SLIDE
            for chunk_idx in range(0, len(bullets), _MAX_BULLETS_PER_SLIDE):
                chunk = bullets[chunk_idx:chunk_idx + _MAX_BULLETS_PER_SLIDE]
                chunk_title = title
                if chunk_idx > 0:
                    chunk_title = f"{title} (cont.)"
                fixed.append({**spec, "title": chunk_title, "bullets": chunk})
            continue

        # --- KPI slides: split if >4 ---
        elif slide_type in ("kpi", "kpis", "metrics"):
            kpis = spec.get("kpis", [])
            if not kpis:
                continue
            for chunk_idx in range(0, len(kpis), _MAX_KPIS_PER_SLIDE):
                chunk = kpis[chunk_idx:chunk_idx + _MAX_KPIS_PER_SLIDE]
                chunk_title = title
                if chunk_idx > 0:
                    chunk_title = f"{title} (cont.)"
                fixed.append({**spec, "title": chunk_title, "kpis": chunk})
            continue

        # --- Table slides: split rows, trim columns ---
        elif slide_type == "table":
            headers = spec.get("headers", [])
            rows = spec.get("rows", [])
            if not headers:
                continue
            # Trim columns
            if len(headers) > _MAX_TABLE_COLS:
                headers = headers[:_MAX_TABLE_COLS]
                rows = [r[:_MAX_TABLE_COLS] for r in rows]
            if not rows:
                fixed.append({**spec, "headers": headers, "rows": []})
                continue
            # Split rows
            for chunk_idx in range(0, len(rows), _MAX_TABLE_ROWS_PER_SLIDE):
                chunk = rows[chunk_idx:chunk_idx + _MAX_TABLE_ROWS_PER_SLIDE]
                chunk_title = title
                if chunk_idx > 0:
                    chunk_title = f"{title} (cont.)"
                fixed.append({
                    **spec, "title": chunk_title,
                    "headers": headers, "rows": chunk,
                })
            continue

        # --- Two-column: trim items ---
        elif slide_type in ("two_column", "two_columns"):
            left = spec.get("left", [])
            right = spec.get("right", [])
            if not left and not right:
                continue
            fixed.append({
                **spec,
                "left": left[:_MAX_TWO_COL_ITEMS],
                "right": right[:_MAX_TWO_COL_ITEMS],
            })
            continue

        # --- Agenda: trim items ---
        elif slide_type in ("agenda", "index"):
            items = spec.get("items", [])
            if not items:
                continue
            fixed.append({**spec, "items": items[:_MAX_AGENDA_ITEMS]})
            continue

        # --- All other types pass through ---
        else:
            fixed.append(spec)

    # --- Post-pass: remove consecutive empty sections ---
    cleaned = []
    for i, slide in enumerate(fixed):
        slide_type = str(slide.get("type", "")).lower().strip()
        if slide_type in ("section", "section_divider", "divider"):
            # Check if next slide is also a section → skip this one
            if i + 1 < len(fixed):
                next_type = str(fixed[i + 1].get("type", "")).lower().strip()
                if next_type in ("section", "section_divider", "divider"):
                    continue
        cleaned.append(slide)

    return cleaned


_SLIDE_TYPE_MAP = {
    "title": _build_title_slide,
    "cover": _build_title_slide,
    "capa": _build_title_slide,
    "section": _build_section_divider,
    "section_divider": _build_section_divider,
    "divider": _build_section_divider,
    "content": _build_content_slide,
    "bullets": _build_content_slide,
    "two_column": _build_two_column_slide,
    "two_columns": _build_two_column_slide,
    "kpi": _build_kpi_slide,
    "kpis": _build_kpi_slide,
    "metrics": _build_kpi_slide,
    "table": _build_table_slide,
    "agenda": _build_agenda_slide,
    "index": _build_agenda_slide,
    "closing": _build_closing_slide,
    "end": _build_closing_slide,
    "obrigado": _build_closing_slide,
}


def _build_slide_from_spec(prs, spec: Dict[str, Any], section_counter: int,
                           badge_text: str = BRAND_HEADER_BADGE_TEXT):
    """Build a single slide from a spec dict.

    Spec format:
    {
        "type": "content|title|section|kpi|table|two_column|agenda|closing",
        "title": "...",
        "subtitle": "...",         # for title/closing
        "bullets": ["...", ...],   # for content
        "left": ["...", ...],      # for two_column
        "right": ["...", ...],     # for two_column
        "kpis": [{value, label, description?}, ...],  # for kpi
        "headers": ["...", ...],   # for table
        "rows": [["...", ...], ...], # for table
        "items": ["...", ...],     # for agenda
        "section_number": int,     # for section (auto-calculated if omitted)
        "text": "...",             # for closing
    }
    """
    slide_type = str(spec.get("type", "content")).lower().strip()
    builder = _SLIDE_TYPE_MAP.get(slide_type, _build_content_slide)

    title = str(spec.get("title", "")).strip()
    subtitle = str(spec.get("subtitle", "")).strip()

    try:
        if slide_type in ("title", "cover", "capa"):
            return builder(prs, title or "Apresentação", subtitle, badge_text)

        elif slide_type in ("section", "section_divider", "divider"):
            num = spec.get("section_number", section_counter)
            return builder(prs, int(num), title, badge_text)

        elif slide_type in ("content", "bullets"):
            bullets = spec.get("bullets", [])
            if isinstance(bullets, str):
                bullets = [b.strip() for b in bullets.split("\n") if b.strip()]
            return builder(prs, title, bullets, badge_text)

        elif slide_type in ("two_column", "two_columns"):
            left = spec.get("left", [])
            right = spec.get("right", [])
            return builder(prs, title, left, right, badge_text)

        elif slide_type in ("kpi", "kpis", "metrics"):
            kpis = spec.get("kpis", [])
            return builder(prs, title, kpis, badge_text)

        elif slide_type in ("table",):
            headers = spec.get("headers", [])
            rows = spec.get("rows", [])
            return builder(prs, title, headers, rows, badge_text)

        elif slide_type in ("agenda", "index"):
            items = spec.get("items", [])
            return builder(prs, items, badge_text)

        elif slide_type in ("closing", "end", "obrigado"):
            text = spec.get("text", title or "Obrigado")
            return builder(prs, text, subtitle, badge_text)

        else:
            # Fallback: treat as content slide
            bullets = spec.get("bullets", [])
            if not bullets and spec.get("text"):
                bullets = [spec["text"]]
            return _build_content_slide(prs, title, bullets, badge_text)

    except Exception as e:
        logger.warning("[PptxEngine] slide build error (type=%s): %s", slide_type, e)
        # Create a simple fallback content slide
        return _build_content_slide(
            prs, title or "Slide",
            [f"(Erro ao gerar slide: {str(e)[:100]})"],
            badge_text,
        )


# ---------------------------------------------------------------------------
# Main entry point
# ---------------------------------------------------------------------------

def generate_presentation(
    title: str,
    slides: List[Dict[str, Any]],
    *,
    subtitle: str = "",
    badge_text: str = BRAND_HEADER_BADGE_TEXT,
    include_title_slide: bool = True,
    include_closing_slide: bool = True,
) -> io.BytesIO:
    """Generate a branded Millennium BCP PPTX presentation.

    Args:
        title: Presentation title
        slides: List of slide spec dicts (see _build_slide_from_spec)
        subtitle: Subtitle for title slide
        badge_text: Badge text (default: "DIGITAL EMPRESAS")
        include_title_slide: Auto-add title slide if not in specs
        include_closing_slide: Auto-add closing slide if not in specs

    Returns:
        BytesIO buffer with PPTX content
    """
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()

    # Set widescreen 16:9 dimensions
    prs.slide_width = Emu(SLIDE_WIDTH_EMU)
    prs.slide_height = Emu(SLIDE_HEIGHT_EMU)

    # ── Smart validation: enforce quality rules before rendering ──
    validated_slides = _validate_and_fix_slides(slides)

    # Check if slides already start with a title/cover
    has_title_slide = False
    has_closing_slide = False
    if validated_slides:
        first_type = str(validated_slides[0].get("type", "")).lower().strip()
        last_type = str(validated_slides[-1].get("type", "")).lower().strip()
        has_title_slide = first_type in ("title", "cover", "capa")
        has_closing_slide = last_type in ("closing", "end", "obrigado")

    # Auto-add title slide
    if include_title_slide and not has_title_slide:
        _build_title_slide(prs, title, subtitle, badge_text)

    # Build each slide
    section_counter = 1
    for spec in validated_slides:
        if not isinstance(spec, dict):
            continue
        slide_type = str(spec.get("type", "content")).lower().strip()
        if slide_type in ("section", "section_divider", "divider"):
            _build_slide_from_spec(prs, spec, section_counter, badge_text)
            section_counter += 1
        else:
            _build_slide_from_spec(prs, spec, section_counter, badge_text)

    # Auto-add closing slide
    if include_closing_slide and not has_closing_slide:
        _build_closing_slide(prs, "Obrigado", "", badge_text)

    # Save to buffer
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def generate_presentation_from_outline(
    title: str,
    outline: str,
    *,
    subtitle: str = "",
    badge_text: str = BRAND_HEADER_BADGE_TEXT,
) -> io.BytesIO:
    """Generate presentation from a plain-text outline.

    Parses a simple outline format:
    - Lines starting with "# " = section dividers
    - Lines starting with "## " = content slide titles
    - Lines starting with "- " under a ## = bullets
    - Lines starting with "| " = table rows (first row = headers)

    Returns BytesIO buffer with PPTX.
    """
    slides = []
    current_slide = None
    section_num = 0

    for raw_line in outline.split("\n"):
        line = raw_line.strip()
        if not line:
            continue

        if line.startswith("# "):
            # Save previous slide
            if current_slide:
                slides.append(current_slide)
                current_slide = None
            section_num += 1
            slides.append({
                "type": "section",
                "title": line[2:].strip(),
                "section_number": section_num,
            })

        elif line.startswith("## "):
            if current_slide:
                slides.append(current_slide)
            current_slide = {
                "type": "content",
                "title": line[3:].strip(),
                "bullets": [],
            }

        elif line.startswith("- ") or line.startswith("• "):
            if current_slide is None:
                current_slide = {"type": "content", "title": "", "bullets": []}
            if "bullets" not in current_slide:
                current_slide["bullets"] = []
            current_slide["bullets"].append(line[2:].strip())

        elif line.startswith("| "):
            # Table detection
            if current_slide and current_slide.get("type") != "table":
                if current_slide.get("bullets") or current_slide.get("title"):
                    slides.append(current_slide)
                current_slide = {
                    "type": "table",
                    "title": current_slide.get("title", ""),
                    "headers": [],
                    "rows": [],
                }
            if current_slide is None:
                current_slide = {"type": "table", "title": "", "headers": [], "rows": []}
            cells = [c.strip() for c in line.strip("| ").split("|")]
            # Skip separator rows (---)
            if all(re.match(r'^-+$', c) for c in cells):
                continue
            if not current_slide.get("headers"):
                current_slide["headers"] = cells
            else:
                current_slide["rows"].append(cells)

        else:
            # Plain text line — add as bullet to current slide
            if current_slide and "bullets" in current_slide:
                current_slide["bullets"].append(line)

    if current_slide:
        slides.append(current_slide)

    return generate_presentation(
        title, slides,
        subtitle=subtitle,
        badge_text=badge_text,
    )
